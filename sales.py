
# ================= SAFE SESSION ACCESSORS =================
def get_price_df():
    return st.session_state.get("price_df", pd.DataFrame())



# ================= PROFIT & MARGIN HELPER =================
def calculate_cost_profit(df, price_df):
    df = df.copy()
    price_df = price_df.copy()

    df["_mat_norm"] = df["Material Description"].astype(str).str.strip().str.upper()
    price_df["_mat_norm"] = price_df["Material Description"].astype(str).str.strip().str.upper()

    price_map = price_df.set_index("_mat_norm")

    df["Cost Price"] = df["_mat_norm"].map(price_map["Cost Price"])
    df["Pack Size"] = df["_mat_norm"].map(price_map["Pack Size"])

    def _cost(row):
        if pd.isna(row["Cost Price"]):
            return None
        if row["UOM"] == "KAR":
            return row["Quantity"] * row["Pack Size"] * row["Cost Price"]
        return row["Quantity"] * row["Cost Price"]

    df["Calculated Cost"] = df.apply(_cost, axis=1)
    df["Gross Profit"] = df["Net Value"] - df["Calculated Cost"]
    df["Margin %"] = (df["Gross Profit"] / df["Net Value"]) * 100
    df["⚠ Cost Missing"] = df["Cost Price"].isna()

    return df


import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
from plotly.subplots import make_subplots
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from sklearn.linear_model import LinearRegression
from statsmodels.tsa.holtwinters import ExponentialSmoothing
import os
from datetime import datetime
from prophet import Prophet
import io
import base64
import streamlit_authenticator as stauth
import hashlib
from difflib import SequenceMatcher
from fuzzywuzzy import fuzz
from io import BytesIO
import urllib.parse
from datetime import date
import streamlit.components.v1 as components
import textwrap






# ================= GLOBAL DISPLAY COLUMN RENAME MAP =================
# NOTE: This is ONLY for table headers (display). Do NOT change calculation logic.
COLUMN_RENAME_MAP = {
    # Market -> Retail everywhere
    "Market Target": "Retail Target",
    "Market Sales": "Retail Sales",
    "Market Balance": "Retail Balance",
    "Market % Achieved": "Retail % Achieved",
    "Market": "Retail",

    # Billing Type codes -> Friendly names
    "YKS1": "HHTCancel",
    "YKS2": "WH1 Cancel",
    "ZCAN": "WH2 Cancel",
    "Cancel Total": "Total Cancel",
    "YKRE": "Salesman Return",
    "ZRE": "Presales Return",
}

def rename_col_key(col_name: str) -> str:
    """Convert one column header to display label."""
    try:
        c = str(col_name).strip()
    except Exception:
        c = col_name
    c = COLUMN_RENAME_MAP.get(c, c)
    # Also replace word Market -> Retail inside longer headers (safe)
    try:
        c = c.replace("Market", "Retail")
    except Exception:
        pass
    return c

def apply_header_renames(df: pd.DataFrame) -> pd.DataFrame:
    """Return a copy with display header names applied (safe, no logic change)."""
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return df
    out = df.copy()
    out.columns = [rename_col_key(c) for c in out.columns]
    return out

def rename_format_keys(formats: dict | None) -> dict | None:
    """Rename format dict keys to match renamed display columns."""
    if not formats:
        return formats
    out = {}
    for k, v in formats.items():
        nk = rename_col_key(k)
        # Fix common typo: '{:,0f}' -> '{:,.0f}'
        if isinstance(v, str) and "{:,0f}" in v:
            v = v.replace("{:,0f}", "{:,.0f}")
        out[nk] = v
    return out


# --- Language Selector ---
st.sidebar.header("Language / اللغة")
language = st.sidebar.selectbox("Choose / اختر", ["English", "العربية"])
lang = "en" if language == "English" else "ar"

if lang == "ar":
    st.markdown("""
    <style>
    .stApp {
        direction: rtl;
        text-align: right;
    }
    .stButton > button {
        float: right;
    }
    .dataframe th, .dataframe td {
        text-align: right !important;
        line-height: normal !important;
    }
    </style>
    """, unsafe_allow_html=True)

texts = {
    "en": {
        "page_title": "📊 Khazan Dashboard",
        "layout": "wide",
        "page_icon": "📈",
        "welcome": "Welcome {0} 👋",
        "logout": "Logout",
        "incorrect_login": "❌ Username/password is incorrect",
        "no_login": "⚠️ Please enter your username and password",
        "dark_mode": "🌙 Dark Mode",
        "upload_header": "📂 Upload Excel (one-time)",
        "upload_tooltip": "Upload an Excel file with sheets: sales data, Target, sales channels, and optionally YTD.",
        "clear_data": "🔁 Clear data",
        "file_loaded": "✅ File loaded — now use the menu to go to any page.",
        "menu_title": "🧭 Menu",
        "navigate": "Navigate",
        "home": "Home",
        "sales_tracking": "Sales Tracking",
        "ytd_comparison": "Year to Date Comparison",
        "custom_analysis": "Custom Analysis",
        "target_allocation": "SP/PY Target Allocation",
        "ai_insights": "AI Insights",
        "customer_insights": "Customer Insights",
        "customer_insights_title": "Customer Insights Dashboard",
        "material_forecast": "Material Forecast",
        "material_forecast_title": "📈 Material Forecast",
        "rfm_analysis_sub": "RFM Analysis",
        "rfm_table_sub": "RFM Table",
        "rfm_chart_sub": "RFM Visualization",
        "rfm_no_data": "No data available for RFM analysis.",
        "rfm_download": "Download RFM Report",
        "rfm_cohort_sub": "RFM Cohort Analysis",
        "rfm_cohort_info": "Analyzes how RFM scores evolve over time for customer acquisition cohorts (grouped by first purchase month).",
        "rfm_cohort_table_sub": "Cohort Summary Table",
        "rfm_cohort_insights_sub": "Key Insights",
        "rfm_cohort_download": "Download Cohort Report (Excel)",
        "rfm_cohort_no_data": "Insufficient data for cohort analysis.",
        "product_availability_checker": "Product Availability Checker",
        "home_title": "🏠 Khazan Dashboard",
        "home_welcome": "**Welcome to your Sales Analytical Hub!**\n- 📈 Track Sales&Targets By Salesman, By Customer, By Branch Name\n- 📅 Sales Comparision\n- 📊 Visualize Trends With Interactive Charts (With Forecaste)\n- 🎯 Allocate Customer Monthly Target \n- 💾 Download reports in PPT & Excel",
        "data_loaded_msg": "Data is loaded — choose a page from the menu.",
        "upload_prompt": "Please upload your Excel file in the sidebar to start.",
        "sales_tracking_title": "📊 MTD Tracking",
        "no_data_warning": "⚠️ Please upload the Excel file in the sidebar (one-time).",
        "filters_header": "🔍 Filters (Sales Tracking)",
        "filters_tooltip": "Filter data by salesmen, billing types, PY, SP, and date range.",
        "select_salesmen": "👥 Select Salesmen",
        "select_billing_types": "📋 Select Billing Types",
        "select_py": "🏬 Select PY Name",
        "select_sp": "🏷️ Select SP Name1",
        "date_presets": "📅 Quick Date Presets",
        "date_presets_options": ["Custom Range", "Last 7 Days", "This Month", "YTD"],
        "select_date_range": "📆 Select Date Range",
        "date_error": "❌ Start date must be before end date.",
        "top_n_salesmen": "🏆 Show Top N Salesmen",
        "no_match_warning": "⚠️ No data matches the selected filters.",
        "kpis_tab": "📈 KPIs",
        "tables_tab": "📋 Tables",
        "charts_tab": "📊 Charts",
        "downloads_tab": "💾 Downloads",
        "key_metrics_sub": "🏆 Key Metrics",
        "total_ka_sales": "Total KA Sales",
        "of_ka_target": "{0:.0f}% of KA Target Achieved",
        "ka_other_ecom": "KA & Other E-com",
        "of_ka_target_pct": "{0:.0f}% of KA Target",
        "talabat_sales": "Talabat Sales",
        "of_talabat_target": "{0:.0f}% of Talabat Target Achieved",
        "target_overview_sub": "🎯 Target Overview",
        "ka_target": "KA Target",
        "talabat_target": "Talabat Target",
        "ka_gap": "KA Gap",
        "talabat_gap": "Talabat Gap",
        "channel_sales_sub": "📊 Channel Sales",
        "retail_sales": "Retail Sales",
        "of_total_ka": "{0:.0f}% of Sales",
        "ecom_sales": "E-com Sales",
        "performance_metrics_sub": "📈 Performance Metrics",
        "days_finished": "Days Elapsed",
        "current_sales_per_day": "Current Sales Per Day",
        "forecast_month_end": "Forecasted Month-End KA Sales",
        "sales_targets_summary_sub": "📋 Sales & Targets Summary-Value",
        "download_sales_targets": "⬇️ Download Sales & Targets Summary (Excel)",
        "sales_by_billing_sub": "📊 Billing Type - Sales - Cancellation And Return ( KWD )",
        "download_billing": "⬇️ Download Billing Type Table (Excel)",
        "sales_by_py_sub": "🏬 Sales Summary By Customer-Value",
        "download_py": "⬇️ Download PY Name Table (Excel)",
        "daily_sales_trend_sub": "📊 Daily Sales Trend And Forecast",
        "daily_sales_title": "Daily Sales Trend And Forecast",
        "not_enough_data": "Not enough data to perform a time-series forecast.",
        "market_vs_ecom_sub": "📊 Market vs E-com Sales",
        "market_vs_ecom_title": "Market vs E-com Sales Distribution",
        "daily_ka_target_sub": "📊 Daily KA Target vs Actual Sales",
        "daily_ka_title": "Daily KA Target vs Actual Sales",
        "salesman_ka_sub": "📊 Salesman KA Target vs Actual",
        "salesman_ka_title": "KA Target vs Actual Sales by Salesman",
        "top10_py_sub": "📊 Top 10 Customers",
        "top10_py_title": "Top 10 Customer 1 by Sales",
        "download_reports_sub": "💾 Download Reports",
        "generate_pptx": "📑 Generate PPTX Report",
        "download_pptx": "⬇️ Download PPTX Report",
        "ytd_title": "📊 Year to Date Comparison",
        "ytd_filters_header": "🔍 Filters (YTD)",
        "ytd_filters_tooltip": "Filter data by salesmen, billing types, PY, SP, and date range.",
        "ytd_select_group": "Group By",
        "ytd_select_value": "Value Column",
        "ytd_period1": "Select Period 1",
        "ytd_period2": "Select Period 2",
        "ytd_no_data": "⚠️ No data matches the filters.",
        "ytd_comparison_sub": "📋 YTD Comparison Table",
        "ytd_download": "⬇️ Download YTD Comparison (Excel)",
        "ytd_chart_title": "YTD Comparison Chart",
        "custom_title": "📊 Custom Analysis",
        "custom_upload": "Upload Additional Sheet (optional)",
        "custom_extra_loaded": "✅ Extra sheet loaded.",
        "custom_select_sheet": "📑 Select Sheet for Analysis",
        "custom_sheet_empty": "⚠️ The sheet '{0}' is empty or not available in your file.",
        "custom_explore": "💡 Explore your data by multiple columns & compare two periods.",
        "custom_group_cols": "Group by columns",
        "custom_value_col": "Value to analyze",
        "custom_periods_sub": "📆 Select Two Periods",
        "custom_period1": "Period 1",
        "custom_period2": "Period 2",
        "custom_select_p1": "Select Period 1",
        "custom_select_p2": "Select Period 2",
        "custom_select_prompt": "👉 Please select at least one group column, one value column, and valid date ranges.",
        "custom_comparison_sub": "📋 Comparison of {0} by {1}",
        "custom_download": "⬇️ Download Comparison (Excel)",
        "target_alloc_title": "🎯 SP/PY Target Allocation",
        "target_config_sub": "Configuration",
        "target_alloc_tooltip": "Allocate targets by branch or customer based on historical sales.",
        "target_alloc_type": "Select Target Allocation Type",
        "target_alloc_options": ["By Branch", "By Customer"],
        "target_input_option": "Select Target Input Option",
        "target_input_options": ["Manual", "Auto (from 'Target' sheet)"],
        "target_enter_total": "Enter the Total Target to be Allocated for this Month (KD)",
        "target_auto_info": "Using Total Target from 'Target' sheet: KD {0:,.0f}",
        "target_zero_warning": "Please ensure the total target is greater than 0.",
        "target_hist_sub": "Historical Data Period",
        "target_hist_option": "Select Historical Data Period",
        "target_hist_options": ["Last 6 Months", "Manual Days"],
        "target_select_range": "Select date range",
        "target_date_warning": "Please select both a start and an end date.",
        "target_no_hist": "⚠️ No sales data available in 'YTD' for {0}.",
        "target_analysis_sub": "🎯 Target Analysis",
        "hist_sales_total": "Historical Sales Total",
        "alloc_target_total": "Allocated Target Total",
        "increase_needed": "Increase Needed vs Avg Sales",
        "current_month_sales": "Current Month Sales",
        "target_balance": "Target Balance",
        "alloc_targets_sub": "📊 Auto-Allocated Targets Based on {0}",
        "download_alloc": "💾 Download Target Allocation Table",
        "ai_title": "🤖 AI Insights",
        "ai_scope_sub": "Scope and filters",
        "ai_select_period": "Select period for insights",
        "ai_top_n": "Top-N salesmen spotlight",
        "ai_no_data": "No data in the selected period. Try expanding the date range.",
        "ai_exec_sub": "📜 Executive summary",
        "ai_prescript_sub": "🛠️ Prescriptive Recommendations",
        "ai_visual_sub": "📊 AI-Generated Visuals",
        "ai_section_sub": "🧭 Section insights",
        "ai_download_sub": "📥 Download executive summary",
        "ai_download_button": "💾 Download AI executive summary (TXT)",
        "ai_ask_sub": "💬 Ask a question about your data",
        "ai_ask_prompt": "Type a question (e.g., 'Which salesman is growing fastest?', 'Where are returns highest?', 'Correlation between targets and sales?')",
        "admin_tools": "Admin Tools",
        "view_logs": "View Audit Logs",
        "audit_title": "📋 Audit Logs",
        "download_logs": "⬇️ Download Audit Logs (Excel)",
        "sheet_missing": "❌ Excel file must contain sheets: {0}. Missing: {1}",
        "cols_missing": "❌ Missing required columns: {0}",
        "load_error": "❌ Error loading Excel file: {0}",
        "pptx_title": "Sales & Targets Report",
        "pptx_generated": "Generated on {0}",
        "pptx_kpi_title": "📈 Key Performance Indicators",
        "pptx_summary_title": "📋 Sales & Targets Summary",
        "pptx_billing_title": "📊 Sales by Billing Type per Salesman",
        "pptx_py_title": "🏬 Sales by PY Name 1",
        "pptx_embed_error": "Chart cannot be embedded: {0}. Install kaleido if missing.",
        "generating_pptx": "⏳ Generating PPTX report...",
        "loading_data": "⏳ Loading Excel data...",
    },
    "ar": {
        "page_title": "📊 لوحة تحكم بيانات حنيف",
        "layout": "wide",
        "page_icon": "📈",
        "welcome": "مرحبا {0} 👋",
        "logout": "تسجيل الخروج",
        "incorrect_login": "❌ اسم المستخدم/كلمة المرور غير صحيحة",
        "no_login": "⚠️ يرجى إدخال اسم المستخدم وكلمة المرور",
        "dark_mode": "🌙 الوضع الداكن",
        "upload_header": "📂 تحميل إكسل (مرة واحدة)",
        "upload_tooltip": "تحميل ملف إكسل يحتوي على أوراق: بيانات المبيعات، الهدف، قنوات المبيعات، واختيارياً YTD.",
        "clear_data": "🔁 مسح البيانات",
        "file_loaded": "✅ تم تحميل الملف — استخدم القائمة الآن للذهاب إلى أي صفحة.",
        "menu_title": "🧭 القائمة",
        "navigate": "التنقل",
        "home": "الرئيسية",
        "sales_tracking": "تتبع المبيعات",
        "ytd_comparison": "مقارنة من بداية السنة",
        "custom_analysis": "تحليل مخصص",
        "target_allocation": "تخصيص أهداف SP/PY",
        "ai_insights": "رؤى الذكاء الاصطناعي",
        "customer_insights": "رؤى العملاء",
        "customer_insights_title": "لوحة رؤى العملاء",
        "material_forecast": "توقعات المواد",
        "material_forecast_title": "📈 توقعات المواد",
        "rfm_analysis_sub": "تحليل RFM",
        "rfm_table_sub": "جدول RFM",
        "rfm_chart_sub": "تصور RFM",
        "rfm_no_data": "لا توجد بيانات متاحة لتحليل RFM.",
        "rfm_download": "تنزيل تقرير RFM",
        "rfm_cohort_sub": "تحليل الأتراب RFM",
        "rfm_cohort_info": "يحلل كيفية تطور درجات RFM بمرور الوقت لأتراب اكتساب العملاء (مجمعة حسب شهر الشراء الأول).",
        "rfm_cohort_table_sub": "جدول ملخص الأتراب",
        "rfm_cohort_insights_sub": "الرؤى الرئيسية",
        "rfm_cohort_download": "تنزيل تقرير الأتراب (إكسل)",
        "rfm_cohort_no_data": "بيانات غير كافية لتحليل الأتراب.",
        "product_availability_checker": "مدقق توفر المنتج",
        "home_title": "🏠 لوحة تحكم بيانات حنيف",
        "home_welcome": "**مرحبا بك في مركز تحليلات المبيعات!**\n- 📈 تتبع المبيعات والأهداف حسب البائع، اسم العميل، اسم الفرع\n- 📊 تصور الاتجاهات مع رسوم بيانية تفاعلية (الآن مع تنبؤ متقدم)\n- 💾 تنزيل التقارير في PPTX و Excel\n- 📅 مقارنة المبيعات عبر فترات مخصصة\n- 🎯 تخصيص أهداف SP/PY بناءً على الأداء الأخير\nاستخدم الشريط الجانبي للتنقل وتحميل البيانات مرة واحدة.",
        "data_loaded_msg": "تم تحميل البيانات — اختر صفحة من القائمة.",
        "upload_prompt": "يرجى تحميل ملف الإكسل في الشريط الجانبي للبدء.",
        "sales_tracking_title": "📊 تتبع MTD",
        "no_data_warning": "⚠️ يرجى تحميل ملف الإكسل في الشريط الجانبي (مرة واحدة).",
        "filters_header": "🔍 فلاتر (تتبع المبيعات)",
        "filters_tooltip": "تصفية البيانات حسب البائعين، أنواع الفواتير، PY، SP، ونطاق التاريخ.",
        "select_salesmen": "👥 اختر البائعين",
        "select_billing_types": "📋 اختر أنواع الفواتير",
        "select_py": "🏬 اختر اسم PY",
        "select_sp": "🏷️ اختر اسم SP1",
        "date_presets": "📅 إعدادات تاريخ سريعة",
        "date_presets_options": ["نطاق مخصص", "آخر 7 أيام", "هذا الشهر", "من بداية السنة"],
        "select_date_range": "📆 اختر نطاق التاريخ",
        "date_error": "❌ تاريخ البداية يجب أن يكون قبل تاريخ النهاية.",
        "top_n_salesmen": "🏆 عرض أفضل N بائعين",
        "no_match_warning": "⚠️ لا توجد بيانات تطابق الفلاتر المختارة.",
        "kpis_tab": "📈 المؤشرات الرئيسية",
        "tables_tab": "📋 الجداول",
        "charts_tab": "📊 الرسوم البيانية",
        "downloads_tab": "💾 التنزيلات",
        "key_matrix_sub": "🏆 المقاييس الرئيسية",
        "total_ka_sales": "إجمالي مبيعات KA",
        "of_ka_target": "{0:.0f}% من هدف KA المحقق",
        "ka_other_ecom": "KA و E-com أخرى",
        "of_ka_target_pct": "{0:.0f}% من هدف KA",
        "talabat_sales": "مبيعات طلبات",
        "of_talabat_target": "{0:.0f}% من هدف طلبات المحقق",
        "target_overview_sub": "🎯 نظرة عامة على الأهداف",
        "ka_target": "هدف KA",
        "talabat_target": "هدف طلبات",
        "ka_gap": "فجوة KA",
        "talabat_gap": "فجوة طلبات",
        "channel_sales_sub": "📊 مبيعات القنوات",
        "retail_sales": "مبيعات التجزئة",
        "of_total_ka": "{0:.0f}% من إجمالي مبيعات KA",
        "ecom_sales": "مبيعات E-com",
        "performance_metrics_sub": "📈 مقاييس الأداء",
        "days_finished": "الأيام المكتملة (العملية)",
        "current_sales_per_day": "المبيعات الحالية لكل يوم",
        "forecast_month_end": "مبيعات KA المتوقعة في نهاية الشهر",
        "sales_targets_summary_sub": "📋 ملخص المبيعات والأهداف",
        "download_sales_targets": "⬇️ تنزيل ملخص المبيعات والأهداف (إكسل)",
        "sales_by_billing_sub": "📊 المبيعات حسب نوع الفاتورة لكل بائع",
        "download_billing": "⬇️ تنزيل جدول أنواع الفواتير (إكسل)",
        "sales_by_py_sub": "🏬 المبيعات حسب اسم PY 1",
        "download_py": "⬇️ تنزيل جدول اسم PY (إكسل)",
        "daily_sales_trend_sub": "📊 اتجاه المبيعات اليومي مع تنبؤ Prophet",
        "daily_sales_title": "اتجاه المبيعات اليومي، تنبؤ Prophet والشذوذات",
        "not_enough_data": "لا توجد بيانات كافية لإجراء تنبؤ زمني.",
        "market_vs_ecom_sub": "📊 السوق مقابل مبيعات E-com",
        "market_vs_ecom_title": "توزيع مبيعات السوق مقابل E-com",
        "daily_ka_target_sub": "📊 هدف KA اليومي مقابل المبيعات الفعلية",
        "daily_ka_title": "هدف KA اليومي مقابل المبيعات الفعلية",
        "salesman_ka_sub": "📊 هدف KA للبائع مقابل الفعلي",
        "salesman_ka_title": "هدف KA مقابل المبيعات الفعلية حسب البائع",
        "top10_py_sub": "📊 أفضل 10 أسماء PY 1 حسب المبيعات",
        "top10_py_title": "أفضل 10 أسماء PY 1 حسب المبيعات",
        "download_reports_sub": "💾 تنزيل التقارير",
        "generate_pptx": "📑 توليد تقرير PPTX",
        "download_pptx": "⬇️ تنزيل تقرير PPTX",
        "ytd_title": "📊 مقارنة من بداية السنة",
        "ytd_filters_header": "🔍 فلاتر (YTD)",
        "ytd_filters_tooltip": "تصفية البيانات حسب البائعين، أنواع الفواتير، PY، SP، ونطاق التاريخ.",
        "ytd_select_group": "التجميع حسب",
        "ytd_select_value": "عمود القيمة",
        "ytd_period1": "اختر الفترة 1",
        "ytd_period2": "اختر الفترة 2",
        "ytd_no_data": "⚠️ لا توجد بيانات تطابق الفلاتر.",
        "ytd_comparison_sub": "📋 جدول مقارنة YTD",
        "ytd_download": "⬇️ تنزيل مقارنة YTD (إكسل)",
        "ytd_chart_title": "رسم بياني لمقارنة YTD",
        "custom_title": "📊 تحليل مخصص",
        "custom_upload": "تحميل ورقة إضافية (اختياري)",
        "custom_extra_loaded": "✅ تم تحميل الورقة الإضافية.",
        "custom_select_sheet": "📑 اختر الورقة للتحليل",
        "custom_sheet_empty": "⚠️ الورقة '{0}' فارغة أو غير متوفرة في ملفك.",
        "custom_explore": "💡 استكشف بياناتك حسب أعمدة متعددة وقارن فترتين.",
        "custom_group_cols": "التجميع حسب الأعمدة",
        "custom_value_col": "القيمة للتحليل",
        "custom_periods_sub": "📆 اختر فترتين",
        "custom_period1": "الفترة 1",
        "custom_period2": "الفترة 2",
        "custom_select_p1": "اختر الفترة 1",
        "custom_select_p2": "اختر الفترة 2",
        "custom_select_prompt": "👉 يرجى اختيار عمود تجميع واحد على الأقل، عمود قيمة واحد، ونطاقات تاريخ صالحة.",
        "custom_comparison_sub": "📋 مقارنة {0} حسب {1}",
        "custom_download": "⬇️ تنزيل المقارنة (إكسل)",
        "target_alloc_title": "🎯 تخصيص أهداف SP/PY",
        "target_config_sub": "التكوين",
        "target_alloc_tooltip": "تخصيص الأهداف حسب الفرع أو العميل بناءً على المبيعات التاريخية.",
        "target_alloc_type": "اختر نوع تخصيص الهدف",
        "target_alloc_options": ["حسب الفرع", "حسب العميل"],
        "target_input_option": "اختر خيار إدخال الهدف",
        "target_input_options": ["يدوي", "تلقائي (من ورقة 'Target')"],
        "target_enter_total": "أدخل إجمالي الهدف للتخصيص لهذا الشهر (د.ك)",
        "target_auto_info": "استخدام إجمالي الهدف من ورقة 'Target': د.ك {0:,.0f}",
        "target_zero_warning": "يرجى التأكد من أن إجمالي الهدف أكبر من 0.",
        "target_hist_sub": "فترة البيانات التاريخية",
        "target_hist_option": "اختر فترة البيانات التاريخية",
        "target_hist_options": ["آخر 6 أشهر", "أيام يدوية"],
        "target_select_range": "اختر نطاق التاريخ",
        "target_date_warning": "يرجى اختيار تاريخ بداية ونهاية.",
        "target_no_hist": "⚠️ لا توجد بيانات مبيعات متوفرة في 'YTD' لـ {0}.",
        "target_analysis_sub": "🎯 تحليل الهدف",
        "hist_sales_total": "إجمالي المبيعات التاريخية",
        "alloc_target_total": "إجمالي الهدف المخصص",
        "increase_needed": "الزيادة المطلوبة مقابل متوسط المبيعات",
        "current_month_sales": "مبيعات الشهر الحالي",
        "target_balance": "رصيد الهدف",
        "alloc_targets_sub": "📊 الأهداف المخصصة تلقائياً بناءً على {0}",
        "download_alloc": "💾 تنزيل جدول تخصيص الهدف",
        "ai_title": "🤖 رؤى الذكاء الاصطناعي",
        "ai_scope_sub": "النطاق والفلاتر",
        "ai_select_period": "اختر الفترة للرؤى",
        "ai_top_n": "أبرز أفضل N بائعين",
        "ai_no_data": "لا توجد بيانات في الفترة المختارة. جرب توسيع نطاق التاريخ.",
        "ai_exec_sub": "📜 ملخص تنفيذي",
        "ai_prescript_sub": "🛠️ توصيات إرشادية",
        "ai_visual_sub": "📊 رسوم بيانية مولدة بالذكاء الاصطناعي",
        "ai_section_sub": "🧭 رؤى حسب القسم",
        "ai_download_sub": "📥 تنزيل الملخص التنفيذي",
        "ai_download_button": "💾 تنزيل ملخص تنفيذي AI (TXT)",
        "ai_ask_sub": "💬 اسأل سؤالاً عن بياناتك",
        "ai_ask_prompt": "اكتب سؤالاً (مثل 'أي بائع ينمو أسرع؟'، 'أين المرتجعات أعلى؟'، 'الارتباط بين الأهداف والمبيعات؟')",
        "admin_tools": "أدوات الإدارة",
        "view_logs": "عرض سجلات التدقيق",
        "audit_title": "📋 سجلات التدقيق",
        "download_logs": "⬇️ تنزيل سجلات التدقيق (إكسل)",
        "sheet_missing": "❌ ملف الإكسل يجب أن يحتوي على أوراق: {0}. المفقودة: {1}",
        "cols_missing": "❌ أعمدة مطلوبة مفقودة: {0}",
        "load_error": "❌ خطأ في تحميل ملف الإكسل: {0}",
        "pptx_title": "تقرير المبيعات والأهداف",
        "pptx_generated": "تم توليده في {0}",
        "pptx_kpi_title": "📈 مؤشرات الأداء الرئيسية",
        "pptx_summary_title": "📋 ملخص المبيعات والأهداف",
        "pptx_billing_title": "📊 المبيعات حسب نوع الفاتورة لكل بائع",
        "pptx_py_title": "🏬 المبيعات حسب اسم PY 1",
        "pptx_embed_error": "لا يمكن تضمين الرسم البياني: {0}. قم بتثبيت kaleido إذا كان مفقوداً.",
        "generating_pptx": "⏳ جاري توليد تقرير PPTX...",
        "loading_data": "⏳ جاري تحميل بيانات الإكسل...",
    }
}

# --- Page Config ---
st.set_page_config(page_title=texts[lang]["page_title"], layout=texts[lang]["layout"], page_icon=texts[lang]["page_icon"])

# --- Streamlit Authenticator (v0.4.2) ---
# Hash the plain-text passwords
hashed_passwords = [
    stauth.Hasher.hash("admin123"),
    stauth.Hasher.hash("salesman1"),
    stauth.Hasher.hash("salesman2")
]

# Define credentials
credentials = {
    "usernames": {
        "admin": {
            "email": "admin@example.com",
            "name": "Admin User",
            "password": hashed_passwords[0],
            "role": "admin"
        },
        "salesman1": {
            "email": "sales1@example.com",
            "name": "Salesman One",
            "password": hashed_passwords[1],
            "role": "salesman",
            "salesman_name": "Salesman One"
        },
        "salesman2": {
            "email": "sales2@example.com",
            "name": "Salesman Two",
            "password": hashed_passwords[2],
            "role": "salesman",
            "salesman_name": "Salesman Two"
        }
    }
}

# Initialize authenticator
authenticator = stauth.Authenticate(
    credentials,
    cookie_name="sales_app",
    key="auth_key",
    cookie_expiry_days=30
)

# Initialize variables
user_role = None
salesman_name = None
username = None

# --- Login ---
authenticator.login(location="main")

if st.session_state.get("authentication_status"):
    st.success(texts[lang]["welcome"].format(st.session_state["name"]))
    authenticator.logout(texts[lang]["logout"], location="sidebar")
    username = st.session_state["username"]
    user_role = credentials["usernames"][username]["role"]
    salesman_name = credentials["usernames"][username].get("salesman_name", None)

elif st.session_state.get("authentication_status") is False:
    st.error(texts[lang]["incorrect_login"])
    st.stop()

elif st.session_state.get("authentication_status") is None:
    st.warning(texts[lang]["no_login"])
    st.stop()

# --- Custom CSS for Visual Enhancements ---
st.markdown(
    """
    <style>
    /* General layout and typography */
    .main {
        background-color: #F8FAFC;
        padding: 30px;
        border-radius: 12px;
        box-shadow: 0 10px 25px rgba(0, 0, 0, 0.05);
        transition: background-color 0.3s;
    }
    h1, h2, h3 {
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        color: #0F172A;
        margin-bottom: 15px;
        font-weight: 700;
    }

    /* Buttons */
    .stButton>button {
        background: linear-gradient(135deg, #1565C0, #1E88E5);
        color: white;
        border: 2px solid #0EA5E9;
        border-radius: 12px;
        padding: 12px 24px;
        font-weight: 700;
        letter-spacing: 0.3px;
        transition: all 0.2s ease;
        box-shadow: 0 6px 18px rgba(2,132,199,0.25);
    }
    .stButton>button:hover {
        transform: translateY(-1px);
        box-shadow: 0 10px 20px rgba(2,132,199,0.35);
    }

    /* Dataframe (tables) */
    .dataframe {
        border-radius: 12px !important;
        overflow: hidden !important;
        box-shadow: 0 4px 16px rgba(0, 0, 0, 0.08);
        border: 1px solid #E2E8F0;
    }
    .dataframe th {
        background: #1E3A8A !important; /* Dark blue for headers */
        color: #FFFFFF !important;
        font-weight: 800 !important;
        padding: 12px !important;
        text-transform: uppercase;
        letter-spacing: 0.4px;
        border: 1px solid #E5E7EB !important;
        line-height: normal !important;
    }
    .dataframe td {
        background-color: #FFFFFF;
        border: 1px solid #E5E7EB !important;
        padding: 10px !important;
        font-weight: 600;
        color: #0F172A;
        vertical-align: middle !important;
        line-height: normal !important;
    }

    /* Sidebar */
    .css-1d391kg {
        background-color: #E2E8F0;
        border-radius: 12px;
        padding: 20px;
        box-shadow: 0 4px 10px rgba(0, 0, 0, 0.05);
    }

    /* Metric card styling with pretty border */
    div[data-testid="stMetric"] {
        border: 2px solid #38BDF8;
        border-radius: 14px;
        padding: 16px 14px;
        background: linear-gradient(180deg, #FFFFFF, #F0F9FF);
        box-shadow: 0 10px 20px rgba(56,189,248,0.25);
        white-space: nowrap;
        overflow: visible;
    }
    div[data-testid="stMetric"] > div {
        color: #0F172A !important;
        font-size: 16px;
    }

    /* Dark mode */
    .dark-mode .main { background-color: #1F2937; }
    .dark-mode h1, .dark-mode h2, .dark-mode h3 { color: #F3F4F6; }
    .dark-mode         line-height: normal !important;
.dataframe td { background-color: #111827; color: #F3F4F6; }
    .dark-mode         line-height: normal !important;
.dataframe th { background: #1E3A8A !important; } /* Dark blue headers in dark mode */
    .dark-mode div[data-testid="stMetric"] { background: linear-gradient(180deg,#111827,#0B1220); border-color:#60A5FA; box-shadow: 0 10px 20px rgba(59,130,246,0.25); }

    /* Tooltip */
    .tooltip { position: relative; display: inline-block; cursor: pointer; }
    .tooltip .tooltiptext {
        visibility: hidden; width: 220px; background-color: #0F172A; color: #fff; text-align: center;
        border-radius: 8px; padding: 8px; position: absolute; z-index: 1; bottom: 125%; left: 50%; margin-left: -110px;
        opacity: 0; transition: opacity 0.3s;
    }
    .tooltip:hover .tooltiptext { visibility: visible; opacity: 1; }

    .progress-bar {
        background-color: #e0e0e0;
        border-radius: 10px;
        margin-top: 5px;
    }
    .progress-bar-fill {
        background-color: #4CAF50;
        height: 15px;
        border-radius: 10px;
        text-align: right;
        padding-right: 5px;
        color: white;
        font-weight: bold;
        transition: width 0.5s ease-in-out;
    }

    /* Green caption styling for specific percentage captions */
    .green-caption {
        color: #15803D !important;
        font-weight: 600;
        font-size: 14px;
    }
    .dark-mode .green-caption {
        color: #6EE7B7 !important; /* Lighter green for dark mode */
    }
    </style>
    """,
    unsafe_allow_html=True
)

# --- Dark Mode Toggle ---
if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = False

st.sidebar.checkbox(
    texts[lang]["dark_mode"],
    value=st.session_state.dark_mode,
    key="dark_mode_toggle",
    on_change=lambda: setattr(st.session_state, "dark_mode", not st.session_state.dark_mode)
)

if st.session_state.dark_mode:
    st.markdown('<script>document.body.classList.add("dark-mode");</script>', unsafe_allow_html=True)
else:
    st.markdown('<script>document.body.classList.remove("dark-mode");</script>', unsafe_allow_html=True)

# --- Cache Data Loading ---
@st.cache_data
def load_data(file):
    with st.spinner(texts[lang]["loading_data"]):
        try:
            xls = pd.ExcelFile(file)

            # ✅ DEFINE FIRST (FIX)
            def normalize_series(s):
                try:
                    return s.astype(str).str.strip().str.lower().replace({'nan': ''})
                except Exception:
                    return s

            # ================= REQUIRED SHEETS =================
            required_sheets = ["sales data", "Target", "sales channels"]
            missing = [s for s in required_sheets if s not in xls.sheet_names]
            if missing:
                st.error(texts[lang]["sheet_missing"].format(
                    ', '.join(required_sheets), ', '.join(missing)
                ))
                return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

            # ================= MAIN SHEETS =================
            sales_df    = pd.read_excel(xls, sheet_name="sales data")
            target_df   = pd.read_excel(xls, sheet_name="Target")
            channels_df = pd.read_excel(xls, sheet_name="sales channels")

            # ================= R&R SHEET (FIXED) =================
            rr_df = pd.read_excel(xls, sheet_name="R&R") if "R&R" in xls.sheet_names else pd.DataFrame()

            if not rr_df.empty and "PY Name 1" in rr_df.columns:
                rr_df["_py_name_norm"] = normalize_series(rr_df["PY Name 1"])

                for col in ["Rebate %", "Display Rental value"]:
                    if col in rr_df.columns:
                        rr_df[col] = pd.to_numeric(rr_df[col], errors="coerce").fillna(0)
                    else:
                        rr_df[col] = 0.0

            # ================= OPTIONAL YTD =================
            ytd_df = pd.read_excel(xls, sheet_name="YTD") if "YTD" in xls.sheet_names else pd.DataFrame()

            # ================= NORMALIZATION =================
            sales_df["Billing Date"] = pd.to_datetime(sales_df["Billing Date"], errors="coerce")
            sales_df["_py_name_norm"] = normalize_series(sales_df["PY Name 1"])

            channels_df["_py_name_norm"] = normalize_series(channels_df["PY Name 1"])
            channels_df["_channels_norm"] = normalize_series(channels_df["Channels"])

            if not ytd_df.empty and "PY Name 1" in ytd_df.columns:
                ytd_df["_py_name_norm"] = normalize_series(ytd_df["PY Name 1"])

            # ================= RETURN ALL =================
            return sales_df, target_df, ytd_df, channels_df, rr_df

        except Exception as e:
            st.error(texts[lang]["load_error"].format(e))
            return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame()
        
        
# --- Helpers: Downloads ---
@st.cache_data
def to_excel_bytes(df: pd.DataFrame, sheet_name: str = "Sheet1", index: bool = False) -> bytes:
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name=sheet_name, index=index)
    return output.getvalue()

@st.cache_data
def to_multi_sheet_excel_bytes(dfs, sheet_names) -> bytes:
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        for df, sn in zip(dfs, sheet_names):
            df.to_excel(writer, sheet_name=sn, index=True)
    return output.getvalue()

# --- PPTX Export ---
def create_pptx(report_df, billing_df, py_table, figs_dict, kpi_data, talabat_tables=None):
    with st.spinner(texts[lang]["generating_pptx"]):
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = texts[lang]["pptx_title"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.name = 'Roboto'
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
        try:
            subtitle = slide.placeholders[1]
            subtitle.text = texts[lang]["pptx_generated"].format(datetime.now().strftime('%Y-%m-%d'))
            subtitle.text_frame.paragraphs[0].font.size = Pt(18)
            subtitle.text_frame.paragraphs[0].font.name = 'Roboto'
            subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(55, 65, 81)
        except Exception:
            pass

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = texts[lang]["pptx_kpi_title"]
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
        slide.shapes.title.text_frame.paragraphs[0].font.name = 'Roboto'
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)

        rows = 4
        cols = 3
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(4)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        kpi_list = list(kpi_data.items())
        for i in range(rows):
            for j in range(cols):
                index = i * cols + j
                if index < len(kpi_list):
                    label, value = kpi_list[index]
                    cell = table.cell(i, j)
                    cell.text = f"{label}\n{value}"
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.name = 'Roboto'
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(243, 244, 246)

        def add_table_slide(df, title):
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
            slide.shapes.title.text_frame.paragraphs[0].font.name = 'Roboto'
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
            rows, cols = df.shape
            table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table
            for j, col in enumerate(df.columns):
                cell = table.cell(0, j)
                cell.text = str(col)
                cell.text_frame.paragraphs[0].font.size = Pt(14)
                cell.text_frame.paragraphs[0].font.name = 'Roboto'
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
            for i, row in enumerate(df.itertuples(index=False), start=1):
                for j, val in enumerate(row):
                    cell = table.cell(i, j)
                    if isinstance(val, (int, float, np.integer, np.floating)):
                        cell.text = f"{val:,.0f}"
                    else:
                        cell.text = str(val)
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.name = 'Roboto'
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(243, 244, 246) if i % 2 == 0 else RGBColor(255, 255, 255)

        def add_chart_slide(fig, title):
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
            slide.shapes.title.text_frame.paragraphs[0].font.name = 'Roboto'
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
            img_stream = io.BytesIO()
            try:
                fig.write_image(img_stream, format="png", width=800, height=600)
                img_stream.seek(0)
                slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9))
            except Exception as e:
                slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4)).text_frame.text = (
                    texts[lang]["pptx_embed_error"].format(e)
                )

        add_table_slide(report_df.reset_index(), texts[lang]["pptx_summary_title"])
        add_table_slide(billing_df.reset_index(), texts[lang]["pptx_billing_title"])
        add_table_slide(py_table.reset_index(), texts[lang]["pptx_py_title"])

        # --- Talabat tables (optional) ---
        if isinstance(talabat_tables, dict):
            tb = talabat_tables.get("billing_split")
            if tb is not None and hasattr(tb, "empty") and (not tb.empty):
                add_table_slide(tb, "🛵 Talabat – Billing Split")
            tc = talabat_tables.get("customers")
            if tc is not None and hasattr(tc, "empty") and (not tc.empty):
                add_table_slide(tc, "🛵 Talabat – Customer Summary")
        for key, fig in figs_dict.items():
            add_chart_slide(fig, key)
        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream

# --- Table Rendering Helpers (consistent headers & full-row highlights) ---
def clean_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    # Remove axis names that can create blank header cells
    try:
        df.rename_axis(None, axis=0, inplace=True)
        df.rename_axis(None, axis=1, inplace=True)
    except Exception:
        pass
    # Ensure all column names are strings (avoid missing/blank headers)
    try:
        df.columns = [("" if c is None else str(c)) for c in df.columns]
    except Exception:
        pass
    return df

def render_table(df, *, formats: dict | None = None, total_row_match=None, hide_index: bool = True):
    '''
    df can be a DataFrame or a pandas Styler.
    formats: dict like {"Sales":"{:,.0f}", "%":"{:.0f}%"}
    total_row_match: function(row)->bool, highlights full row (e.g. lambda r: r.get("Salesman Name")=="Total")
    '''
    try:
        from pandas.io.formats.style import Styler as _Styler
        is_styler = isinstance(df, _Styler)
    except Exception:
        is_styler = False

    if (not is_styler) and isinstance(df, pd.DataFrame):
        df = clean_columns(df)
        df = apply_header_renames(df)

        sty = df.style.set_table_styles([
            {'selector': 'th', 'props': [
                ('background', '#1E3A8A'),
                ('color', '#FFFFFF'),
                ('font-weight', '800'),
                ('border', '1px solid #E5E7EB'),
                ('text-align', 'center')
            ]}
        ])

        if total_row_match:
            def _hl(row):
                try:
                    if total_row_match(row):
                        return ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900' for _ in row]
                except Exception:
                    pass
                return ['' for _ in row]
            sty = sty.apply(_hl, axis=1)

        if formats:
            formats = rename_format_keys(formats)
            sty = sty.format(formats)

        st.dataframe(sty, use_container_width=True, hide_index=hide_index)
        return

    # If it's already a Styler, just display (best effort)
    st.dataframe(df, use_container_width=True, hide_index=hide_index)

# --- Positive/Negative Coloring ---
def color_positive_negative(val):
    try:
        v = float(val)
        color = "#15803D" if v > 0 else "#B91C1C" if v < 0 else ""
        return f"color: {color}; font-weight: bold"
    except:
        return ""

def create_progress_bar_html(percentage):
    safe_pct = max(0, min(100, percentage))
    fill_color = "#4CAF50" if safe_pct >= 100 else "#2196F3"
    html = f"""
    <div style="background-color: #f0f0f0; border-radius: 5px; height: 10px; margin-top: 5px;">
        <div style="background-color: {fill_color}; height: 100%; width: {safe_pct}%; border-radius: 5px; text-align: right; font-size: 8px; color: white;">
        </div>
    </div>
    """
    return html

# --- SINGLE SIDEBAR UPLOADER ---
st.sidebar.header(texts[lang]["upload_header"])
st.sidebar.markdown(
    f'<div class="tooltip">ℹ️<span class="tooltiptext">{texts[lang]["upload_tooltip"]}</span></div>',
    unsafe_allow_html=True
)

uploaded = st.sidebar.file_uploader("", type=["xlsx"], key="single_upload")

if st.sidebar.button(texts[lang]["clear_data"]):
    for k in [
        "sales_df", "target_df", "ytd_df", "channels_df",
        "rr_df",               # ← ADD THIS LINE HERE
        "price_df", "data_loaded", "audit_log"
    ]:
        if k in st.session_state:
            del st.session_state[k]
    st.experimental_rerun()


# ================= PRICE LIST (FOR PROFIT & MARGIN) =================
# IMPORTANT: keep in session_state (prevents NameError on reruns)
if "price_df" not in st.session_state:
    st.session_state["price_df"] = pd.DataFrame()

if uploaded is not None:
    try:
        xls = pd.ExcelFile(uploaded)
        if "price list" in xls.sheet_names:
            st.session_state["price_df"] = pd.read_excel(
                xls, sheet_name="price list"
            )
    except Exception:
        st.session_state["price_df"] = pd.DataFrame()


# ================= LOAD MAIN DATA (ONCE) =================
if uploaded is not None and "data_loaded" not in st.session_state:
    sales_df, target_df, ytd_df, channels_df, rr_df = load_data(uploaded)

    st.session_state["sales_df"] = sales_df
    st.session_state["target_df"] = target_df
    st.session_state["ytd_df"] = ytd_df
    st.session_state["channels_df"] = channels_df
    st.session_state["rr_df"] = rr_df   # ← ADD THIS
    st.session_state["data_loaded"] = True
    st.session_state["audit_log"] = []  # Initialize audit log

    st.success(texts[lang]["file_loaded"])

    # Log upload action
    st.session_state["audit_log"].append({
        "user": username,
        "action": "upload_file",
        "details": uploaded.name,
        "timestamp": datetime.now().isoformat()
    })

# --- Sidebar Menu with multilingual support ---
st.sidebar.title(texts[lang]["menu_title"])

menu = [
    texts[lang]["home"],
    texts[lang]["sales_tracking"],
    texts[lang]["ytd_comparison"],
    texts[lang]["custom_analysis"],
    texts[lang]["target_allocation"],
    texts[lang]["ai_insights"],
    texts[lang]["customer_insights"],
    texts[lang]["material_forecast"],
    "💰 Profit & Margin",          # 👈 ADD THIS
    "🧭 Management Command Center"
]


choice = st.sidebar.selectbox(texts[lang]["navigate"], menu)
# Role-based filtering for data
if "data_loaded" in st.session_state:
    sales_df = st.session_state["sales_df"].copy()
    target_df = st.session_state["target_df"].copy()
    ytd_df = st.session_state["ytd_df"].copy()
    channels_df = st.session_state["channels_df"].copy()

    if user_role == "salesman" and salesman_name:
        sales_df = sales_df[sales_df["Driver Name EN"] == salesman_name]
        ytd_df = ytd_df[ytd_df.get("Driver Name EN", pd.Series()) == salesman_name]
        target_df = target_df[target_df.get("Driver Name EN", pd.Series()) == salesman_name]

# --- Home Page ---
if choice == texts[lang]["home"]:
    st.title(texts[lang]["home_title"])
    with st.container():
        st.markdown(
            texts[lang]["home_welcome"],
            unsafe_allow_html=True
        )
    if "data_loaded" in st.session_state: st.success(texts[lang]["data_loaded_msg"])
    else: st.info(texts[lang]["upload_prompt"])


# --- Sales Tracking Page ---
elif choice == texts[lang]["sales_tracking"]:
    st.title(texts[lang]["sales_tracking_title"])
    if "data_loaded" not in st.session_state:
        st.warning(texts[lang]["no_data_warning"])
    else:
        # Filters
        st.sidebar.subheader(texts[lang]["filters_header"])
        st.sidebar.markdown(f'<div class="tooltip">ℹ️<span class="tooltiptext">{texts[lang]["filters_tooltip"]}</span></div>', unsafe_allow_html=True)
        salesmen = st.sidebar.multiselect(
            texts[lang]["select_salesmen"],
            options=sorted(sales_df["Driver Name EN"].dropna().unique()),
            default=sorted(sales_df["Driver Name EN"].dropna().unique()),
            key="st_salesmen"
        )
        billing_types = st.sidebar.multiselect(
            texts[lang]["select_billing_types"],
            options=sorted(sales_df["Billing Type"].dropna().unique()),
            default=sorted(sales_df["Billing Type"].dropna().unique()),
            key="st_billing_types"
        )
        py_filter = st.sidebar.multiselect(
            texts[lang]["select_py"],
            options=sorted(sales_df["PY Name 1"].dropna().unique()),
            default=sorted(sales_df["PY Name 1"].dropna().unique()),
            key="st_py_filter"
        )
        sp_filter = st.sidebar.multiselect(
            texts[lang]["select_sp"],
            options=sorted(sales_df["SP Name1"].dropna().unique()),
            default=sorted(sales_df["SP Name1"].dropna().unique()),
            key="st_sp_filter"
        )

        preset = st.sidebar.radio(texts[lang]["date_presets"], texts[lang]["date_presets_options"], key="st_preset")
        today = pd.Timestamp.today().normalize()
        if preset == texts[lang]["date_presets_options"][1]:  # Last 7 Days
            date_range = [today - pd.Timedelta(days=7), today]
        elif preset == texts[lang]["date_presets_options"][2]:  # This Month
            month_start = today.replace(day=1)
            month_end = month_start + pd.offsets.MonthEnd(0)
            date_range = [month_start, month_end]
        elif preset == texts[lang]["date_presets_options"][3]:  # YTD
            date_range = [today.replace(month=1, day=1), today]
        else:
            date_range = st.sidebar.date_input(
                texts[lang]["select_date_range"],
                [sales_df["Billing Date"].min(), sales_df["Billing Date"].max()],
                key="st_date_range"
            )
            if isinstance(date_range, tuple) and len(date_range) == 2:
                date_range = [pd.to_datetime(date_range[0]), pd.to_datetime(date_range[1])]
            else:
                date_range = [sales_df["Billing Date"].min(), sales_df["Billing Date"].max()]

        if date_range[0] > date_range[1]:
            st.error(texts[lang]["date_error"])
        else:
            top_n = st.sidebar.slider(
                texts[lang]["top_n_salesmen"],
                min_value=1,
                max_value=max(1, len(sales_df["Driver Name EN"].dropna().unique())),
                value=min(5, max(1, len(sales_df["Driver Name EN"].dropna().unique()))),
                key="st_topn"
            )

            df_filtered = sales_df[
                (sales_df["Driver Name EN"].isin(salesmen))
                & (sales_df["Billing Type"].isin(billing_types))
                & (sales_df["Billing Date"] >= date_range[0])
                & (sales_df["Billing Date"] <= date_range[1])
                & (sales_df["PY Name 1"].isin(py_filter))
                & (sales_df["SP Name1"].isin(sp_filter))
            ].copy()

            if df_filtered.empty:
                st.warning(texts[lang]["no_match_warning"])
            else:
                billing_start = df_filtered["Billing Date"].min().normalize()
                billing_end = df_filtered["Billing Date"].max().normalize()
                all_days = pd.date_range(billing_start, billing_end, freq="D")
                days_finish = int(sum(1 for d in all_days if d.weekday() != 4))

                current_month_start = today.replace(day=1)
                current_month_end = current_month_start + pd.offsets.MonthEnd(0)
                month_days = pd.date_range(current_month_start, current_month_end, freq="D")
                working_days_current_month = int(sum(1 for d in month_days if d.weekday() != 4))

                # --- Base aggregates ---
                total_sales = df_filtered.groupby("Driver Name EN")["Net Value"].sum()
                talabat_df = df_filtered[df_filtered["PY Name 1"] == "STORES SERVICES KUWAIT CO."]
                talabat_sales = talabat_df.groupby("Driver Name EN")["Net Value"].sum()

                ka_targets = target_df.set_index("Driver Name EN")["KA Target"] if "KA Target" in target_df.columns else pd.Series(dtype=float)
                talabat_targets = target_df.set_index("Driver Name EN")["Talabat Target"] if "Talabat Target" in target_df.columns else pd.Series(dtype=float)

                all_salesmen_idx = total_sales.index.union(talabat_sales.index).union(ka_targets.index).union(talabat_targets.index)
                total_sales = total_sales.reindex(all_salesmen_idx, fill_value=0).astype(float)
                talabat_sales = talabat_sales.reindex(all_salesmen_idx, fill_value=0).astype(float)
                ka_targets = ka_targets.reindex(all_salesmen_idx, fill_value=0).astype(float)
                talabat_targets = talabat_targets.reindex(all_salesmen_idx, fill_value=0).astype(float)

                ka_gap = (ka_targets - total_sales).clip(lower=0)
                talabat_gap = (talabat_targets - talabat_sales).clip(lower=0)

                top_salesmen = total_sales.sort_values(ascending=False).head(top_n).index
                total_sales_top = total_sales.loc[top_salesmen]
                talabat_sales_top = talabat_sales.loc[top_salesmen]
                ka_gap_top = ka_gap.loc[top_salesmen]
                talabat_gap_top = talabat_gap.loc[top_salesmen]

                # --- Talabat detailed tables (for Tables / Downloads / PPTX) ---
                TALABAT_PY = "STORES SERVICES KUWAIT CO."
                talabat_df_detail = df_filtered[df_filtered["PY Name 1"] == TALABAT_PY].copy()

                def _talabat_group_row(row):
                    bt = str(row.get("Billing Type", "")).strip().upper()
                    try:
                        nv = float(row.get("Net Value", 0) or 0)
                    except Exception:
                        nv = 0.0
                    if bt == "ZFR":
                        return "ZFR"
                    if bt == "HHT":
                        return "HHT"
                    # Returns: negative value OR common return codes / patterns
                    if (nv < 0) or ("RE" in bt) or bt in {"YKF2", "YKRE", "ZRE", "ZCR", "CR"}:
                        return "Returns"
                    return "Other"

                if not talabat_df_detail.empty:
                    talabat_df_detail["Talabat Billing Group"] = talabat_df_detail.apply(_talabat_group_row, axis=1)
                else:
                    talabat_df_detail["Talabat Billing Group"] = pd.Series(dtype=str)

                talabat_billing_split = (
                    talabat_df_detail
                    .pivot_table(
                        index="Driver Name EN",
                        columns="Talabat Billing Group",
                        values="Net Value",
                        aggfunc="sum",
                        fill_value=0
                    )
                    .reindex(all_salesmen_idx, fill_value=0)
                )
                for _c in ["ZFR", "HHT", "Returns", "Other"]:
                    if _c not in talabat_billing_split.columns:
                        talabat_billing_split[_c] = 0.0
                talabat_billing_split = talabat_billing_split[["ZFR", "HHT", "Returns", "Other"]]
                talabat_billing_split["Total"] = talabat_billing_split.sum(axis=1)
                talabat_billing_split = talabat_billing_split.reset_index().rename(columns={"Driver Name EN": "Salesman"})

                def _pick_customer_col(df):
                    for _col in ["Customer Name", "Branch Name", "PY Name 1"]:
                        if _col in df.columns and df[_col].notna().any():
                            return _col
                    # Fallback (if you don't have outlet/branch columns)
                    return "SP Name1" if "SP Name1" in df.columns else "PY Name 1"

                _cust_col = _pick_customer_col(talabat_df_detail)
                if (not talabat_df_detail.empty) and (_cust_col in talabat_df_detail.columns):
                    talabat_customer_table = (
                        talabat_df_detail
                        .groupby(_cust_col, dropna=False)
                        .agg(
                            **{
                                "Talabat Sales": ("Net Value", "sum"),
                                "Orders": ("Net Value", "size"),
                            }
                        )
                        .sort_values("Talabat Sales", ascending=False)
                        .reset_index()
                        .rename(columns={_cust_col: "Customer"})
                    )
                else:
                    talabat_customer_table = pd.DataFrame(columns=["Customer", "Talabat Sales", "Orders"])

                if (not talabat_df_detail.empty) and ("Billing Date" in talabat_df_detail.columns):
                    talabat_df_detail["Talabat Date"] = pd.to_datetime(talabat_df_detail["Billing Date"], errors="coerce").dt.date
                    talabat_daily_trend = (
                        talabat_df_detail
                        .groupby(["Talabat Date", "Talabat Billing Group"])["Net Value"].sum()
                        .reset_index()
                        .pivot_table(
                            index="Talabat Date",
                            columns="Talabat Billing Group",
                            values="Net Value",
                            aggfunc="sum",
                            fill_value=0
                        )
                        .sort_index()
                    )
                    for _c in ["ZFR", "HHT", "Returns", "Other"]:
                        if _c not in talabat_daily_trend.columns:
                            talabat_daily_trend[_c] = 0.0
                    talabat_daily_trend = talabat_daily_trend[["ZFR", "HHT", "Returns", "Other"]]
                    talabat_daily_trend["Total"] = talabat_daily_trend.sum(axis=1)
                    talabat_daily_trend = talabat_daily_trend.reset_index().rename(columns={"Talabat Date": "Date"})
                else:
                    talabat_daily_trend = pd.DataFrame(columns=["Date", "ZFR", "HHT", "Returns", "Other", "Total"])

                total_ka_target_all = float(ka_targets.sum())
                total_tal_target_all = float(talabat_targets.sum())
                per_day_ka_target = (total_ka_target_all / working_days_current_month) if working_days_current_month > 0 else 0
                current_sales_per_day = (total_sales.sum() / days_finish) if days_finish > 0 else 0
                forecast_month_end_ka = current_sales_per_day * working_days_current_month

                # --- Channels mapping: Market (Retail) vs E-com ---
                df_py_sales = df_filtered.groupby("_py_name_norm")["Net Value"].sum().reset_index()
                df_channels_merged = df_py_sales.merge(
                    channels_df[["_py_name_norm", "Channels"]],
                    on="_py_name_norm",
                    how="left"
                )
                df_channels_merged["Channels"] = df_channels_merged["Channels"].str.strip().str.lower().fillna("uncategorized")
                channel_sales = df_channels_merged.groupby("Channels")["Net Value"].sum()
                total_retail_sales = float(channel_sales.get("market", 0.0) + channel_sales.get("uncategorized", 0.0))
                total_ecom_sales = float(channel_sales.get("e-com", 0.0))
                total_channel_sales = total_retail_sales + total_ecom_sales
                retail_sales_pct = (total_retail_sales / total_channel_sales * 100) if total_channel_sales > 0 else 0
                ecom_sales_pct = (total_ecom_sales / total_channel_sales * 100) if total_channel_sales > 0 else 0

                # --- KA & Other E-com Calculation ---
                ka_other_ecom_sales = total_sales.sum() - talabat_sales.sum()
                ka_other_ecom_pct = (ka_other_ecom_sales / total_ka_target_all * 100) if total_ka_target_all > 0 else 0

                # --- KPI Data for PPTX ---
                kpi_data = {
                    texts[lang]["ka_target"]: f"KD {total_ka_target_all:,.0f}",
                    texts[lang]["talabat_target"]: f"KD {total_tal_target_all:,.0f}",
                    texts[lang]["ka_gap"]: f"KD {(total_ka_target_all - total_sales.sum()):,.0f}",
                    "Total Talabat Gap": f"KD {talabat_gap.sum():,.0f}",
                    texts[lang]["total_ka_sales"]: f"KD {total_sales.sum():,.0f} ({((total_sales.sum() / total_ka_target_all) * 100):.0f}%)" if total_ka_target_all else f"KD {total_sales.sum():,.0f} (0%)",
                    "Total Talabat Sales": f"KD {talabat_sales.sum():,.0f} ({((talabat_sales.sum() / total_tal_target_all) * 100):.0f}%)" if total_tal_target_all else f"KD {talabat_sales.sum():,.0f} (0%)",
                    texts[lang]["ka_other_ecom"]: f"KD {ka_other_ecom_sales:,.0f} ({ka_other_ecom_pct:.0f}%)",
                    texts[lang]["retail_sales"]: f"KD {total_retail_sales:,.0f} ({retail_sales_pct:.0f}%)",
                    texts[lang]["ecom_sales"]: f"KD {total_ecom_sales:,.0f} ({ecom_sales_pct:.0f}%)",
                    texts[lang]["days_finished"]: f"{days_finish}",
                    "Per Day KA Target": f"KD {per_day_ka_target:,.0f}",
                    texts[lang]["current_sales_per_day"]: f"KD {current_sales_per_day:,.0f}",
                    texts[lang]["forecast_month_end"]: f"KD {forecast_month_end_ka:,.0f}"
                }

                tabs = st.tabs([texts[lang]["kpis_tab"], texts[lang]["tables_tab"], texts[lang]["charts_tab"], texts[lang]["downloads_tab"]])

 # --- KPIs with progress bars ---
                with tabs[0]:
                    st.subheader(texts[lang]["key_metrics_sub"])

                    # ROW 1: Big KPI (Full width)
                    r1c1 = st.columns(1)[0]

                    with r1c1:
                        DATE_COL = "Billing Date"
                        _d = df_filtered.copy()
                        _d[DATE_COL] = pd.to_datetime(_d[DATE_COL], errors="coerce")

                        from_dt = _d[DATE_COL].min()
                        to_dt   = _d[DATE_COL].max()

                        from_txt = from_dt.strftime("%d %b %Y") if pd.notna(from_dt) else "-"
                        to_txt   = to_dt.strftime("%d %b %Y") if pd.notna(to_dt) else "-"

                        ka_sales_value = float(total_sales.sum()) if total_sales is not None else 0.0

                      
                        # --- Total KA Sales card (FULL width + date inside) ---
                        st.markdown(
                            """
                            <style>
                            /* Force full width */
                            .metric-card {width:100% !important; max-width:100% !important; display:block !important; box-sizing:border-box;}
                            </style>
                            """,
                            unsafe_allow_html=True
                        )

                        card_html = (
                            f'<div style="'
                            f'border:2px solid #38BDF8;'
                            f'border-radius:18px;'
                            f'padding:12px 16px;'
                            f'background: linear-gradient(180deg, #FFFFFF, #F0F9FF);'
                            f'box-shadow:0 10px 22px rgba(56,189,248,0.25);'
                            f'width:100%;'
                            f'box-sizing:border-box;'
                            f'">'
                            f'<div style="font-size:18px;font-weight:700;color:#0F172A;margin-bottom:8px;">{texts[lang]["total_ka_sales"]}</div>'
                            f'<div style="font-size:44px;font-weight:900;color:#111827;line-height:1.1;margin-bottom:14px;">KD {ka_sales_value:,.0f}</div>'
                            f'<div style="display:flex;justify-content:space-between;font-size:13px;color:#334155;opacity:0.9;">'
                            f'<span><b>From:</b> {from_txt}</span>'
                            f'<span><b>To:</b> {to_txt}</span>'
                            f'</div>'
                            f'</div>'
                        )

                        st.markdown(card_html, unsafe_allow_html=True)

                        # Progress bar (will now visually match because card is full width)
                        progress_pct_ka = (ka_sales_value / total_ka_target_all * 100) if total_ka_target_all > 0 else 0
                        st.markdown(create_progress_bar_html(progress_pct_ka), unsafe_allow_html=True)
                        st.markdown(
                            f'<div class="green-caption">{texts[lang]["of_ka_target"].format(progress_pct_ka)}</div>',
                            unsafe_allow_html=True
                        )

                    # ROW 2: Two KPIs (Side by side)  ✅ MUST be outside r1c1
                    r2c1, r2c2 = st.columns(2)

                    with r2c1:
                        st.metric(texts[lang]["ka_other_ecom"], f"KD {ka_other_ecom_sales:,.0f}")
                        st.markdown(create_progress_bar_html(ka_other_ecom_pct), unsafe_allow_html=True)
                        st.markdown(
                            f'<div class="green-caption">{texts[lang]["of_ka_target_pct"].format(ka_other_ecom_pct)}</div>',
                            unsafe_allow_html=True
                        )

                    with r2c2:
                        st.metric(texts[lang]["talabat_sales"], f"KD {talabat_sales.sum():,.0f}")
                        progress_pct_talabat = (talabat_sales.sum() / total_tal_target_all * 100) if total_tal_target_all > 0 else 0
                        st.markdown(create_progress_bar_html(progress_pct_talabat), unsafe_allow_html=True)
                        st.markdown(
                            f'<div class="green-caption">{texts[lang]["of_talabat_target"].format(progress_pct_talabat)}</div>',
                            unsafe_allow_html=True
                        )

                    # ---- Rest stays same (Target Overview etc.) ----
                    st.subheader(texts[lang]["target_overview_sub"])
                    r3c1, r3c2, r3c3, r3c4 = st.columns(4)
                    r3c1.metric(texts[lang]["ka_target"], f"KD {total_ka_target_all:,.0f}")
                    r3c2.metric(texts[lang]["talabat_target"], f"KD {total_tal_target_all:,.0f}")
                    r3c3.metric(texts[lang]["ka_gap"], f"KD {(total_ka_target_all - total_sales.sum()):,.0f}")
                    r3c4.metric(texts[lang]["talabat_gap"], f"KD {talabat_gap.sum():,.0f}")

                    st.subheader(texts[lang]["channel_sales_sub"])
                    r4c1, r4c2 = st.columns(2)
                    with r4c1:
                        st.metric(texts[lang]["retail_sales"], f"KD {total_retail_sales:,.0f}")
                        retail_contribution_pct = (total_retail_sales / total_sales.sum() * 100) if total_sales.sum() > 0 else 0
                        st.caption(texts[lang]["of_total_ka"].format(retail_contribution_pct))
                    with r4c2:
                        st.metric(texts[lang]["ecom_sales"], f"KD {total_ecom_sales:,.0f}")
                        ecom_contribution_pct = (total_ecom_sales / total_sales.sum() * 100) if total_sales.sum() > 0 else 0
                        st.caption(texts[lang]["of_total_ka"].format(ecom_contribution_pct))

                    st.subheader(texts[lang]["performance_metrics_sub"])
                    r5c1, r5c2, r5c3 = st.columns(3)
                    r5c1.metric(texts[lang]["days_finished"], days_finish)
                    r5c2.metric(texts[lang]["current_sales_per_day"], f"KD {current_sales_per_day:,.0f}")
                    r5c3.metric(texts[lang]["forecast_month_end"], f"KD {forecast_month_end_ka:,.0f}")
                    
                    
                    # --- TABLES ---
                    with tabs[1]:

                        st.subheader(texts[lang]["sales_targets_summary_sub"])

                        # ================= BASE INDEX =================
                        idx = ka_targets.index

                        # ================= KA =================
                        ka_target = ka_targets.reindex(idx, fill_value=0).astype(float)
                        ka_sales = total_sales.reindex(idx, fill_value=0).astype(float)

                        ka_balance = (ka_target - ka_sales).clip(lower=0)
                        ka_percent = np.where(ka_target > 0, (ka_sales / ka_target * 100).round(0), 0)

                        # ================= CHANNEL SPLIT (SAFE DERIVATION) =================
                        df_sales = df_filtered[["Driver Name EN", "PY Name 1", "Net Value"]].copy()

                        # normalize PY names
                        df_sales["_py_norm"] = (
                            df_sales["PY Name 1"].astype(str).str.strip().str.lower()
                        )

                        ch_tmp = channels_df.copy()
                        ch_tmp["_py_norm"] = (
                            ch_tmp["PY Name 1"].astype(str).str.strip().str.lower()
                        )

                        df_sales = df_sales.merge(
                            ch_tmp[["_py_norm", "Channels"]],
                            on="_py_norm",
                            how="left"
                        )

                        df_sales["Channels"] = (
                            df_sales["Channels"]
                            .astype(str)
                            .str.lower()
                            .str.strip()
                        )

                        # default channel
                        df_sales.loc[
                            df_sales["Channels"].isin(["", "nan", "none"]),
                            "Channels"
                        ] = "market"

                        # ================= E-COM / MARKET =================
                        ecom_mask = df_sales["Channels"].str.contains(
                            "e-com|ecom|ecommerce|online|talabat",
                            regex=True,
                            na=False
                        )

                        ecom_sales = (
                            df_sales[ecom_mask]
                            .groupby("Driver Name EN")["Net Value"]
                            .sum()
                            .reindex(idx, fill_value=0)
                        )

                        market_sales = (
                            df_sales[~ecom_mask]
                            .groupby("Driver Name EN")["Net Value"]
                            .sum()
                            .reindex(idx, fill_value=0)
                        )

                        # ================= E-COM TARGET =================
                        ecom_target = pd.Series(0.0, index=idx)

                        if not target_df.empty and "Driver Name EN" in target_df.columns:
                            col_map = {c.lower().strip(): c for c in target_df.columns}

                            for k in [
                                "e-com target", "ecom target", "e-commerce target",
                                "ecom target kd", "e-com target kd"
                            ]:
                                if k in col_map:
                                    ecom_target = (
                                        target_df
                                        .set_index("Driver Name EN")[col_map[k]]
                                        .apply(pd.to_numeric, errors="coerce")
                                        .fillna(0)
                                        .reindex(idx, fill_value=0)
                                    )
                                    break

                        # ================= MARKET =================
                        market_target = (ka_target - ecom_target).clip(lower=0)

                        market_balance = (market_target - market_sales).clip(lower=0)
                        market_percent = np.where(
                            market_target > 0,
                            (market_sales / market_target * 100).round(0),
                            0
                        )

                        ecom_balance = (ecom_target - ecom_sales).clip(lower=0)
                        ecom_percent = np.where(
                            ecom_target > 0,
                            (ecom_sales / ecom_target * 100).round(0),
                            0
                        )

                        # ================= FINAL TABLE =================
                        report_df = pd.DataFrame({
                            "Salesman Name": idx,

                            "Total Target": ka_target.values,
                            "Total Sales": ka_sales.values,
                            "Total Balance": ka_balance.values,
                            "Total % Achieved": ka_percent,

                            "Market Target": market_target.values,
                            "Market Sales": market_sales.values,
                            "Market Balance": market_balance.values,
                            "Market % Achieved": market_percent,

                            "E-Com Target": ecom_target.values,
                            "E-Com Sales": ecom_sales.values,
                            "E-Com Balance": ecom_balance.values,
                            "E-Com % Achieved": ecom_percent,
                        })

                        # ================= TOTAL ROW =================
                        total_row = report_df.sum(numeric_only=True).to_frame().T
                        total_row["Salesman Name"] = "Total"

                        def pct(a, b):
                            return round(a / b * 100, 0) if b > 0 else 0

                        total_row["Total % Achieved"] = pct(
                            total_row["Total Sales"].iloc[0],
                            total_row["Total Target"].iloc[0]
                        )
                        total_row["Market % Achieved"] = pct(
                            total_row["Market Sales"].iloc[0],
                            total_row["Market Target"].iloc[0]
                        )
                        total_row["E-Com % Achieved"] = pct(
                            total_row["E-Com Sales"].iloc[0],
                            total_row["E-Com Target"].iloc[0]
                        )

                        report_df = pd.concat([report_df, total_row], ignore_index=True)

                        # ================= SORT =================
                        data_part = (
                            report_df[report_df["Salesman Name"] != "Total"]
                            .sort_values("Total % Achieved", ascending=False)
                        )

                        total_part = report_df[report_df["Salesman Name"] == "Total"]

                        report_df = pd.concat([data_part, total_part], ignore_index=True)

                        # ================= STYLING =================
                        def row_style(row):
                            if row["Salesman Name"] == "Total":
                                return ["background-color:#BFDBFE; color:#1E3A8A; font-weight:900"] * len(row)
                            return ["background-color:#F9FAFB" if row.name % 2 == 0 else ""] * len(row)

                        def pct_color(v):
                            if v >= 100:
                                return "color:#166534; font-weight:700"
                            elif v >= 80:
                                return "color:#92400E; font-weight:600"
                            return "color:#991B1B; font-weight:700"

                        # Apply display header renames ONLY for showing the table (no logic change)
                        report_df_disp = apply_header_renames(report_df)

                        num_cols = [
                            "Total Target","Total Sales","Total Balance",
                            "Market Target","Market Sales","Market Balance",
                            "E-Com Target","E-Com Sales","E-Com Balance"
                        ]
                        pct_cols = ["Total % Achieved","Market % Achieved","E-Com % Achieved"]

                        styled = (
                            report_df_disp.style
                            .apply(row_style, axis=1)
                            .format("{:,.0f}", subset=[rename_col_key(c) for c in num_cols])
                            .format("{:.0f}%", subset=[rename_col_key(c) for c in pct_cols])
                            .applymap(pct_color, subset=[rename_col_key(c) for c in pct_cols])
                        )

                        st.dataframe(styled, use_container_width=True, hide_index=True)

                        # --- Sales by Billing Type per Salesman ---
                        st.subheader(texts[lang]["sales_by_billing_sub"])
                        billing_wide = df_filtered.pivot_table(
                            index="Driver Name EN",
                            columns="Billing Type",
                            values="Net Value",
                            aggfunc="sum",
                            fill_value=0
                        )

                        required_cols_raw = ["ZFR", "YKF2", "YKRE", "YKS1", "YKS2", "ZCAN", "ZRE"]
                        billing_wide = billing_wide.reindex(columns=required_cols_raw, fill_value=0)
                        display_df = billing_wide.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})
                        display_df["Sales Total"] = billing_wide.sum(axis=1)
                        display_df["Return"] = billing_wide["YKRE"] + billing_wide["ZRE"]
                        display_df["Return %"] = np.where(display_df["Sales Total"] != 0,
                                                        (display_df["Return"] / display_df["Sales Total"] * 100).round(0), 0)
                        display_df["Cancel Total"] = billing_wide[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)

                        ordered_cols = ["Presales", "HHT", "Sales Total", "YKS1", "YKS2", "ZCAN",
                                        "Cancel Total", "YKRE", "ZRE", "Return", "Return %"]
                        display_df = display_df.reindex(columns=ordered_cols, fill_value=0)

                        total_row = pd.DataFrame(display_df.sum(numeric_only=True)).T
                        total_row.index = ["Total"]
                        total_row["Return %"] = round((total_row["Return"] / total_row["Sales Total"] * 100), 0) if total_row["Sales Total"].iloc[0] != 0 else 0
                        billing_df = pd.concat([display_df, total_row])
                        billing_df.index.name = "Salesman"
                        # --- Styling + Display (fixed Total row color & header) ---
                        billing_df_show = billing_df.reset_index()  # brings 'Salesman' as a real column
                        render_table(
                            billing_df_show,
                            hide_index=True,
                            total_row_match=lambda r: str(r.get("Salesman", "")).strip() == "Total",
                            formats={
                                "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                                "YKS1": "{:,.0f}", "YKS2": "{:,.0f}", "ZCAN": "{:,.0f}", "Cancel Total": "{:,.0f}",
                                "YKRE": "{:,.0f}", "ZRE": "{:,.0f}", "Return": "{:,.0f}", "Return %": "{:.0f}%"
                            }
                        )

                        if st.download_button(
                            texts[lang]["download_billing"],
                            data=to_excel_bytes(billing_df.reset_index(), sheet_name="Billing_Types", index=False),
                            file_name=f"Billing_Types_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        ):
                            st.session_state["audit_log"].append({
                                "user": username,
                                "action": "download",
                                "details": "Billing Type Table Excel",
                                "timestamp": datetime.now()
                            })


                        # --- Sales Summary By Customer – Value ---
                        st.subheader("📌 Sales Summary By Customer – Value")

                        # Sales grouped by Customer
                        py_table = df_filtered.groupby("PY Name 1")["Net Value"].sum().sort_values(ascending=False).to_frame(name="Sales")

                        # Returns grouped by Customer (Billing Type = YKRE or ZRE)
                        returns_df = df_filtered[df_filtered["Billing Type"].isin(["YKRE", "ZRE"])] \
                            .groupby("PY Name 1")["Net Value"].sum()

                        # Add returns into table
                        py_table["Returns"] = py_table.index.map(returns_df).fillna(0.0)

                        # Calculate return %
                        py_table["Return %"] = np.where(
                            py_table["Sales"] > 0,
                            (py_table["Returns"] / py_table["Sales"] * 100).round(1),
                            0
                        )

                        # Contribution %
                        py_table["Contribution %"] = np.where(
                            py_table["Sales"] > 0,
                            (py_table["Sales"] / py_table["Sales"].sum() * 100).round(1),
                            0
                        )

                        # Total Row
                        total_row = pd.DataFrame({
                            "Sales": [py_table["Sales"].sum()],
                            "Returns": [py_table["Returns"].sum()],
                            "Return %": [(py_table["Returns"].sum() / py_table["Sales"].sum() * 100).round(1) if py_table["Sales"].sum() > 0 else 0],
                            "Contribution %": [100.0]
                        }, index=["Total"])

                        py_table_with_total = pd.concat([py_table, total_row])

                        py_table_with_total.index.name = "Customer Name"

                        # Styling function
                        # --- Styling + Display (fixed Total row color & header) ---
                        py_show = py_table_with_total.reset_index()  # brings 'Customer Name' as a real column
                        # Ensure the first column is clearly named
                        if py_show.columns[0] == "":
                            py_show = py_show.rename(columns={py_show.columns[0]: "Customer Name"})
                        render_table(
                            py_show,
                            hide_index=True,
                            total_row_match=lambda r: str(r.get("Customer Name", "")).strip() == "Total",
                            formats={
                                "Sales": "{:,.0f}",
                                "Returns": "{:,.0f}",
                                "Return %": "{:.1f}%",
                                "Contribution %": "{:.1f}%"
                            }
                        )

                        # Download Button
                        if st.download_button(
                            "⬇️ Download Customer Summary (Excel)",
                            data=to_excel_bytes(py_table_with_total.reset_index(), sheet_name="Sales_by_Customer", index=False),
                            file_name=f"Sales_by_Customer_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        ):
                            st.session_state["audit_log"].append({
                                "user": username,
                                "action": "download",
                                "details": "Customer Summary Excel",
                                "timestamp": datetime.now()
                            })


                        # --- Return by SP Name1 ---
                        st.subheader("🔄 Sales Vs Return's Summary By Branch-Value")
                        sp_billing = df_filtered.pivot_table(
                            index="SP Name1",
                            columns="Billing Type",
                            values="Net Value",
                            aggfunc="sum",
                            fill_value=0
                        )
                        sp_billing = sp_billing.reindex(columns=required_cols_raw, fill_value=0)
                        sp_billing["Sales Total"] = sp_billing.sum(axis=1)
                        sp_billing["Return"] = sp_billing["YKRE"] + sp_billing["ZRE"]
                        sp_billing["Cancel Total"] = sp_billing[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)
                        sp_billing = sp_billing.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})
                        sp_billing["Return %"] = np.where(sp_billing["Sales Total"] != 0,
                                                        (sp_billing["Return"] / sp_billing["Sales Total"] * 100).round(0), 0)
                        sp_billing = sp_billing.reindex(columns=ordered_cols, fill_value=0).astype(int)

                        total_row = pd.DataFrame(sp_billing.sum(numeric_only=True)).T
                        total_row.index = ["Total"]
                        total_row["Return %"] = round((total_row["Return"]/total_row["Sales Total"]*100), 0) if total_row["Sales Total"].iloc[0]!=0 else 0
                        sp_billing = pd.concat([sp_billing, total_row])
                        # --- Styling + Display (fixed missing header & Total row color) ---
                        sp_billing_show = sp_billing.reset_index().rename(columns={"SP Name1": "Branch Name"})
                        render_table(
                            sp_billing_show,
                            hide_index=True,
                            total_row_match=lambda r: str(r.get("Branch Name", r.get("SP Name1", r.get("Branch", "")))).strip() == "Total",
                            formats={
                                "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                                "YKS1": "{:,.0f}", "YKS2": "{:,.0f}", "ZCAN": "{:,.0f}", "Cancel Total": "{:,.0f}",
                                "YKRE": "{:,.0f}", "ZRE": "{:,.0f}", "Return": "{:,.0f}", "Return %": "{:.0f}%"
                            }
                        )                        # --- Return by Material Description ---
                        st.subheader("🔄 Return Summary By SKU")
                        if "Material Description" in df_filtered.columns:
                            material_billing = df_filtered.pivot_table(
                                index="Material Description",
                                columns="Billing Type",
                                values="Net Value",
                                aggfunc="sum",
                                fill_value=0
                            )

                            material_cols_raw = ["ZFR", "YKF2", "YKRE", "YKS1", "YKS2", "ZCAN", "ZRE"]
                            material_billing = material_billing.reindex(columns=material_cols_raw, fill_value=0)
                            material_billing["Sales Total"] = material_billing.sum(axis=1)
                            material_billing["Return"] = material_billing["YKRE"] + material_billing["ZRE"]
                            material_billing["Cancel Total"] = material_billing[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)
                            material_billing = material_billing.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})
                            material_billing["Return %"] = np.where(
                                material_billing["Sales Total"] != 0,
                                (material_billing["Return"] / material_billing["Sales Total"] * 100).round(0),
                                0
                            )

                            ordered_cols_material = [
                                "Presales", "HHT", "Sales Total",
                                "YKS1", "YKS2", "ZCAN",
                                "Cancel Total", "YKRE", "ZRE",
                                "Return", "Return %"
                            ]
                            material_billing = material_billing.reindex(columns=ordered_cols_material, fill_value=0)

                            total_row = pd.DataFrame(material_billing.sum(numeric_only=True)).T
                            total_row.index = ["Total"]
                            total_row["Return %"] = round(
                                (total_row["Return"] / total_row["Sales Total"] * 100), 0
                            ) if total_row["Sales Total"].iloc[0] != 0 else 0
                            material_billing = pd.concat([material_billing, total_row])

                            # ✅ Display-only header renames (no logic change)
                            material_show = apply_header_renames(material_billing)

                            def highlight_total_row_material(row):
                                if str(row.name).strip() == "Total":
                                    return ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900'] * len(row)
                                return ['' for _ in row]

                            styled_material = (
                                material_show.style
                                .set_table_styles([
                                    {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                                 ('font-weight', '800'), ('height', '40px'),
                                                                 ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                                ])
                                .apply(highlight_total_row_material, axis=1)
                                .format({
                                    "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                                    "HHTCancel": "{:,.0f}", "WH1 Cancel": "{:,.0f}", "WH2 Cancel": "{:,.0f}", "Total Cancel": "{:,.0f}",
                                    "Salesman Return": "{:,.0f}", "Presales Return": "{:,.0f}",
                                    "Return": "{:,.0f}", "Return %": "{:.0f}%"
                                })
                            )
                            st.dataframe(styled_material, use_container_width=True, hide_index=False)

                            if st.download_button(
                                texts[lang].get("download_material", "Download Return by Material Description"),
                                data=to_excel_bytes(material_billing.reset_index(), sheet_name="Return_by_Material", index=False),
                                file_name=f"Return_by_Material_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                            ):
                                st.session_state["audit_log"].append({
                                    "user": username,
                                    "action": "download",
                                    "details": "Return by Material Description Excel",
                                    "timestamp": datetime.now()
                                })
                        else:
                            st.info("No 'Material Description' column found in data — skipping Material Description table.")

                        # --- Return by SP Name1 + Material Description ---
                        st.subheader("🔄 Return Summary By Branch & Product ")
                        required_cols = {"SP Name1", "Material Description", "Billing Type", "Net Value"}
                        if required_cols.issubset(df_filtered.columns):
                            # Pivot table
                            sp_mat_table = pd.pivot_table(
                                df_filtered,
                                index=["SP Name1", "Material Description"],
                                columns="Billing Type",
                                values="Net Value",
                                aggfunc="sum",
                                fill_value=0
                            )

                            # Ensure all billing columns exist
                            billing_cols = ["ZFR", "YKF2", "YKRE", "YKS1", "YKS2", "ZCAN", "ZRE"]
                            for col in billing_cols:
                                if col not in sp_mat_table.columns:
                                    sp_mat_table[col] = 0

                            # Rename sales columns only (keep raw codes for logic)
                            sp_mat_table = sp_mat_table.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})

                            # Calculate totals (logic unchanged)
                            sp_mat_table["Sales Total"] = sp_mat_table.sum(axis=1, numeric_only=True)
                            sp_mat_table["Return"] = sp_mat_table["YKRE"] + sp_mat_table["ZRE"]
                            sp_mat_table["Cancel Total"] = sp_mat_table[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)
                            sp_mat_table["Return %"] = np.where(
                                sp_mat_table["Sales Total"] != 0,
                                (sp_mat_table["Return"] / sp_mat_table["Sales Total"] * 100).round(0),
                                0
                            )

                            # Reorder columns (raw names)
                            ordered_cols_spm = [
                                "Presales", "HHT", "Sales Total",
                                "YKS1", "YKS2", "ZCAN",
                                "Cancel Total", "YKRE", "ZRE",
                                "Return", "Return %"
                            ]
                            sp_mat_table = sp_mat_table.reindex(columns=ordered_cols_spm, fill_value=0)

                            # Add total row
                            total_row = pd.DataFrame(sp_mat_table.sum(numeric_only=True)).T
                            total_row.index = [("Total", "")]
                            total_row["Return %"] = round(
                                (total_row["Return"] / total_row["Sales Total"] * 100), 0
                            ) if total_row["Sales Total"].iloc[0] != 0 else 0
                            sp_mat_table = pd.concat([sp_mat_table, total_row])

                            # ✅ Display-only header renames (no logic change)
                            sp_mat_show = apply_header_renames(sp_mat_table)

                            def highlight_sp_mat(row):
                                styles = []
                                for col in row.index:
                                    if row.name == ("Total", ""):
                                        styles.append('background-color: #BFDBFE; color: #1E3A8A; font-weight: 900')
                                    elif col == "Return" and row.get(col, 0) != 0:
                                        styles.append('background-color: #FECACA; color: #991B1B; font-weight: 700')
                                    elif col == "Total Cancel" and row.get(col, 0) != 0:
                                        styles.append('background-color: #FDE68A; color: #92400E; font-weight: 700')
                                    elif col == "Sales Total" and row.get(col, 0) != 0:
                                        styles.append('background-color: #D1FAE5; color: #065F46; font-weight: 700')
                                    else:
                                        styles.append('')
                                return styles

                            styled_sp_mat = (
                                sp_mat_show.style
                                .set_table_styles([
                                    {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                                 ('font-weight', '800'), ('height', '40px'),
                                                                 ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                                ])
                                .apply(highlight_sp_mat, axis=1)
                                .format({
                                    "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                                    "HHTCancel": "{:,.0f}", "WH1 Cancel": "{:,.0f}", "WH2 Cancel": "{:,.0f}", "Total Cancel": "{:,.0f}",
                                    "Salesman Return": "{:,.0f}", "Presales Return": "{:,.0f}",
                                    "Return": "{:,.0f}", "Return %": "{:.0f}%"
                                })
                            )

                            st.dataframe(styled_sp_mat, use_container_width=True, hide_index=True)

                            if st.download_button(
                                texts[lang].get("download_sp_material", "Download Return by SP+Material"),
                                data=to_excel_bytes(sp_mat_table.reset_index(), sheet_name="Return_by_SP_Material", index=False),
                                file_name=f"Return_by_SP_Material_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                            ):
                                st.session_state["audit_log"].append({
                                    "user": username,
                                    "action": "download",
                                    "details": "Return by SP+Material Excel",
                                    "timestamp": datetime.now()
                                })
                        else:
                            st.info("Required columns are missing in your data for SP+Material table.")

                            
                            


                        # ------------------------------------------------
                        # 🛵 Talabat – MTD details (Billing split / Customers / Daily trend + Excel)
                        # Place this inside: with tabs[1]:
                        # ------------------------------------------------

                        TALABAT_PY = "STORES SERVICES KUWAIT CO."

                        talabat_only_df = df_filtered[df_filtered["PY Name 1"] == TALABAT_PY].copy()

                        # Defaults (so display won't break)
                        talabat_billing_split = pd.DataFrame()
                        talabat_customer_table = pd.DataFrame()
                        talabat_daily_trend = pd.DataFrame()

                        if not talabat_only_df.empty:

                            # ---- Billing Type Group Mapping (FIXED) ----
                            SALES_ZFR = {"ZFR"}
                            SALES_HHT = {"YKF2"}              # HHT in your app = YKF2
                            RETURNS_CODES = {"YKRE", "ZRE"}   # ONLY these are returns

                            def _bt_group(bt):
                                bt = str(bt).strip().upper()
                                if bt in SALES_ZFR:
                                    return "ZFR"
                                if bt in SALES_HHT:
                                    return "HHT"
                                if bt in RETURNS_CODES:
                                    return "Returns"
                                return "Other"

                            talabat_only_df["_bt_group"] = talabat_only_df["Billing Type"].apply(_bt_group)

                            # =========================
                            # 1) Billing split by Salesman
                            # =========================
                            talabat_billing_split = (
                                talabat_only_df
                                .groupby(["Driver Name EN", "_bt_group"])["Net Value"]
                                .sum()
                                .unstack(fill_value=0)
                            )

                            for c in ["ZFR", "HHT", "Returns", "Other"]:
                                if c not in talabat_billing_split.columns:
                                    talabat_billing_split[c] = 0

                            talabat_billing_split = talabat_billing_split[["ZFR", "HHT", "Returns", "Other"]]
                            talabat_billing_split["Total"] = talabat_billing_split.sum(axis=1)
                            talabat_billing_split = talabat_billing_split.reset_index().rename(columns={"Driver Name EN": "Salesman"})

                            # Total row
                            talabat_billing_split = pd.concat([
                                talabat_billing_split,
                                pd.DataFrame([{
                                    "Salesman": "Total",
                                    "ZFR": talabat_billing_split["ZFR"].sum(),
                                    "HHT": talabat_billing_split["HHT"].sum(),
                                    "Returns": talabat_billing_split["Returns"].sum(),
                                    "Other": talabat_billing_split["Other"].sum(),
                                    "Total": talabat_billing_split["Total"].sum(),
                                }])
                            ], ignore_index=True)

                            # =========================
                            # 2) Customer / Outlet summary
                            # =========================
                            candidate_customer_cols = [
                                "Customer", "Customer Name", "Outlet", "Outlet Name", "Branch Name",
                                "Ship-to Name", "Sold-to Name", "SP Name1", "PY Name 1"
                            ]
                            customer_col = next((c for c in candidate_customer_cols if c in talabat_only_df.columns), None)

                            candidate_order_cols = ["Billing Document", "Invoice", "Invoice No", "Sales Document", "Document No"]
                            order_col = next((c for c in candidate_order_cols if c in talabat_only_df.columns), None)

                            if customer_col:
                                if order_col:
                                    orders_series = talabat_only_df.groupby(customer_col)[order_col].nunique()
                                else:
                                    orders_series = talabat_only_df.groupby(customer_col).size()

                                talabat_customer_table = (
                                    talabat_only_df.groupby(customer_col)["Net Value"].sum()
                                    .to_frame("Talabat Sales")
                                    .join(orders_series.to_frame("Orders"))
                                    .reset_index()
                                    .rename(columns={customer_col: "Customer"})
                                    .sort_values("Talabat Sales", ascending=False)
                                )

                                talabat_customer_table = pd.concat([
                                    talabat_customer_table,
                                    pd.DataFrame([{
                                        "Customer": "Total",
                                        "Talabat Sales": talabat_customer_table["Talabat Sales"].sum(),
                                        "Orders": talabat_customer_table["Orders"].sum()
                                    }])
                                ], ignore_index=True)

                            # =========================
                            # 3) Daily trend (MTD)
                            # =========================
                            if "Billing Date" in talabat_only_df.columns:
                                talabat_only_df["Billing Date"] = pd.to_datetime(talabat_only_df["Billing Date"], errors="coerce")

                                daily = (
                                    talabat_only_df
                                    .dropna(subset=["Billing Date"])
                                    .groupby([talabat_only_df["Billing Date"].dt.date, "_bt_group"])["Net Value"]
                                    .sum()
                                    .unstack(fill_value=0)
                                )

                                for c in ["ZFR", "HHT", "Returns", "Other"]:
                                    if c not in daily.columns:
                                        daily[c] = 0

                                daily = daily[["ZFR", "HHT", "Returns", "Other"]]
                                daily["Total"] = daily.sum(axis=1)
                                talabat_daily_trend = daily.reset_index().rename(columns={"Billing Date": "Date"})
                                talabat_daily_trend.columns = ["Date"] + [c for c in talabat_daily_trend.columns if c != "Date"]


                        # ------------------------------------------------
                        # 🛵 Talabat – DISPLAY (your same style)
                        # ------------------------------------------------
                        st.markdown("---")
                        st.subheader("🛵 Talabat -( Sales  / Returns / Cancellation  )- ( By Billing Type )")

                        if isinstance(talabat_billing_split, pd.DataFrame) and (not talabat_billing_split.empty):
                            st.dataframe(
                                talabat_billing_split.style.format({
                                    "ZFR": "{:,.0f}",
                                    "HHT": "{:,.0f}",
                                    "Returns": "{:,.0f}",
                                    "Other": "{:,.0f}",
                                    "Total": "{:,.0f}",
                                }),
                                use_container_width=True,
                                hide_index=True
                            )
                        else:
                            st.info("No Talabat data found in the selected filters.")

                        st.subheader("🛵 Talabat – Customer / Outlet Summary")
                        if isinstance(talabat_customer_table, pd.DataFrame) and (not talabat_customer_table.empty):
                            st.dataframe(
                                talabat_customer_table.style.format({
                                    "Talabat Sales": "{:,.0f}",
                                    "Orders": "{:,.0f}",
                                }),
                                use_container_width=True,
                                hide_index=True
                            )
                        else:
                            st.info("No Talabat customer-level data available (missing customer/outlet columns or no rows).")

                        st.subheader("🛵 Talabat – Daily Trend (MTD)")
                        if isinstance(talabat_daily_trend, pd.DataFrame) and (not talabat_daily_trend.empty):
                            st.dataframe(
                                talabat_daily_trend.style.format({
                                    "ZFR": "{:,.0f}",
                                    "HHT": "{:,.0f}",
                                    "Returns": "{:,.0f}",
                                    "Other": "{:,.0f}",
                                    "Total": "{:,.0f}",
                                }),
                                use_container_width=True,
                                hide_index=True
                            )

                            # Simple line chart (Total)
                            try:
                                fig_tal_daily = px.line(
                                    talabat_daily_trend,
                                    x="Date",
                                    y="Total",
                                    markers=True,
                                    title="Talabat Daily Sales (Total)"
                                )
                                st.plotly_chart(fig_tal_daily, use_container_width=True)
                            except Exception:
                                pass
                        else:
                            st.info("No Talabat daily trend data available.")


                        # ------------------------------------------------
                        # ⬇️ Talabat Excel Download (3 sheets)
                        # Uses your helper: to_multi_sheet_excel_bytes(dfs, sheet_names)
                        # ------------------------------------------------
                        if (isinstance(talabat_billing_split, pd.DataFrame) and not talabat_billing_split.empty) or \
                        (isinstance(talabat_customer_table, pd.DataFrame) and not talabat_customer_table.empty) or \
                        (isinstance(talabat_daily_trend, pd.DataFrame) and not talabat_daily_trend.empty):

                            dfs = [
                                talabat_billing_split.set_index("Salesman") if (not talabat_billing_split.empty and "Salesman" in talabat_billing_split.columns) else pd.DataFrame(),
                                talabat_customer_table.set_index("Customer") if (not talabat_customer_table.empty and "Customer" in talabat_customer_table.columns) else pd.DataFrame(),
                                talabat_daily_trend.set_index("Date") if (not talabat_daily_trend.empty and "Date" in talabat_daily_trend.columns) else pd.DataFrame(),
                            ]
                            sheet_names = ["Talabat_Billing_Split", "Talabat_Customers", "Talabat_Daily_Trend"]

                            tal_excel = to_multi_sheet_excel_bytes(dfs, sheet_names)

                            st.download_button(
                                "⬇️ Download Talabat Details (Excel)",
                                data=tal_excel,
                                file_name=f"Talabat_MTD_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                            )


                        # Download Talabat Excel (3 sheets)
                        talabat_xlsx_bytes = to_multi_sheet_excel_bytes(
                            dfs=[
                                talabat_billing_split,
                                talabat_customer_table,
                                talabat_daily_trend,
                            ],
                            sheet_names=[
                                "Talabat_Billing_Split",
                                "Talabat_Customers",
                                "Talabat_Daily_Trend",
                            ]
                        )

                        if st.download_button(
                            "⬇️ Download Talabat Details (Excel)",
                            data=talabat_xlsx_bytes,
                            file_name=f"Talabat_MTD_Details_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        ):
                            st.session_state["audit_log"].append({
                                "user": username,
                                "action": "download",
                                "details": "Talabat MTD Details (Excel)",
                                "timestamp": datetime.now()
                            })

                    # ================================
                    # 📊 CHARTS TAB (GM Premium Visuals)
                    # ================================
                    # ================================
                    # 📊 CHARTS TAB (D: Data-heavy Analytical Command Center) - UPDATED (Retail rename + MTD/Forecast fix)
                    # ================================
                    with tabs[2]:

                        from calendar import monthrange

                        st.subheader("📈 Sales Trends & Channel Performance – Management Command Center")
                        st.caption("Analytical view: Toggle period • Forecast • Run-rate • Channel • Salesman • Drivers")

                        if df_filtered is None or df_filtered.empty:
                            st.info("No data available for charts.")
                            st.stop()

                        if "Billing Date" not in df_filtered.columns:
                            st.error("⚠️ 'Billing Date' column not found!")
                            st.stop()

                        df = df_filtered.copy()

                        # ✅ Normalize dates (prevents month filter issues)
                        df["Billing Date"] = pd.to_datetime(df["Billing Date"], errors="coerce").dt.normalize()
                        df["Net Value"] = pd.to_numeric(df.get("Net Value", 0), errors="coerce").fillna(0.0)
                        df = df.dropna(subset=["Billing Date"])

                        if df.empty:
                            st.info("No valid Billing Date rows to plot.")
                            st.stop()

                        # ---------------------------------------------------
                        # 0️⃣ Period Toggle (Daily/Weekly/Monthly)
                        # ---------------------------------------------------
                        gran = st.radio("View", ["Daily", "Weekly", "Monthly"], horizontal=True, index=0)
                        freq_map = {"Daily": "D", "Weekly": "W-MON", "Monthly": "MS"}
                        freq = freq_map[gran]

                        # ---------------------------------------------------
                        # 1️⃣ Month-end forecast + run-rate (✅ USE LATEST DATE IN DATA, NOT SYSTEM TODAY)
                        # ---------------------------------------------------
                        as_of = df["Billing Date"].max()          # ✅ latest date from filtered data
                        cm, cy = int(as_of.month), int(as_of.year)
                        total_days_month = monthrange(cy, cm)[1]

                        month_start = as_of.replace(day=1)        # ✅ month start for the as_of month
                        df_cm = df[(df["Billing Date"] >= month_start) & (df["Billing Date"] <= as_of)]

                        mtd_sales = float(df_cm["Net Value"].sum())
                        days_passed = int(df_cm["Billing Date"].dt.date.nunique())
                        remaining_days = max(total_days_month - days_passed, 0)

                        avg_per_day = (mtd_sales / days_passed) if days_passed > 0 else 0.0
                        month_end_forecast = mtd_sales + (avg_per_day * remaining_days)

                        per_day_target = float(per_day_ka_target) if "per_day_ka_target" in globals() else 0.0
                        monthly_target = per_day_target * total_days_month
                        remaining_target = monthly_target - mtd_sales
                        required_run_rate = (remaining_target / remaining_days) if remaining_days > 0 else 0.0

                        achievement_mtd = (mtd_sales / monthly_target * 100) if monthly_target > 0 else 0.0

                        # ---------------------------------------------------
                        # 2️⃣ Channel totals (Retail vs E-com)  ✅ rename Market -> Retail
                        # ---------------------------------------------------
                        df_channel_temp = df.groupby("PY Name 1")["Net Value"].sum().reset_index()

                        df_ch_merge = df_channel_temp.merge(
                            channels_df[["PY Name 1", "Channels"]],
                            on="PY Name 1",
                            how="left"
                        )

                        df_ch_merge["Channels"] = (
                            df_ch_merge["Channels"]
                            .astype(str).str.strip().str.lower()
                            .replace({"": "retail", "nan": "retail", "none": "retail", "market": "retail"})
                        )

                        total_ecom = float(df_ch_merge.loc[df_ch_merge["Channels"] == "e-com", "Net Value"].sum())
                        total_retail = float(df_ch_merge.loc[df_ch_merge["Channels"] != "e-com", "Net Value"].sum())
                        total_all = total_retail + total_ecom

                        ecom_share = (total_ecom / total_all * 100) if total_all > 0 else 0.0
                        retail_share = (total_retail / total_all * 100) if total_all > 0 else 0.0

                        # ---------------------------------------------------
                        # 3️⃣ KPI Strip (data-heavy)
                        # ---------------------------------------------------
                        active_days = int(df["Billing Date"].dt.date.nunique())

                        k1, k2, k3, k4, k5, k6 = st.columns(6)
                        k1.metric("Total Sales", f"KD {total_all:,.0f}")
                        k2.metric("Retail", f"KD {total_retail:,.0f}", f"{retail_share:.1f}%")
                        k3.metric("E-com", f"KD {total_ecom:,.0f}", f"{ecom_share:.1f}%")
                        k4.metric("MTD Sales", f"KD {mtd_sales:,.0f}", f"{achievement_mtd:.1f}% of M target")
                        k5.metric("Month-End Forecast", f"KD {month_end_forecast:,.0f}")
                        k6.metric("Required / Day", f"KD {required_run_rate:,.0f}")

                        st.caption(f"📅 As-of Date (from data): {as_of.date()} | Active Days in Filter: {active_days:,}")
                        st.markdown("---")

                        # ---------------------------------------------------
                        # 4️⃣ Period Trend (Actual vs Target) + Forecast line
                        # ---------------------------------------------------
                        c1, c2 = st.columns([0.62, 0.38], gap="large")

                        with c1:
                            st.markdown(f"### 📌 {gran} Trend (Actual vs Target) + Month-End Forecast")

                            df_period = (df.groupby(pd.Grouper(key="Billing Date", freq=freq))["Net Value"]
                                        .sum()
                                        .reset_index()
                                        .rename(columns={"Billing Date": "Period", "Net Value": "Sales"}))

                            days_in_period = (
                                df.assign(_date_only=df["Billing Date"].dt.date)
                                .groupby(pd.Grouper(key="Billing Date", freq=freq))["_date_only"]
                                .nunique()
                                .reset_index(name="Days")
                                .rename(columns={"Billing Date": "Period"})
                            )

                            df_period = df_period.merge(days_in_period, on="Period", how="left")
                            df_period["Days"] = pd.to_numeric(df_period["Days"], errors="coerce").fillna(1)

                            df_period["Target"] = df_period["Days"] * per_day_target

                            # ✅ make it look clean (avoid weird x-axis behavior)
                            df_period = df_period.sort_values("Period")

                            fig_trend = go.Figure()
                            fig_trend.add_trace(go.Scatter(
                                x=df_period["Period"], y=df_period["Sales"],
                                mode="lines+markers", name="Actual",
                                line=dict(width=3),
                                marker=dict(size=6)
                            ))
                            fig_trend.add_trace(go.Scatter(
                                x=df_period["Period"], y=df_period["Target"],
                                mode="lines", name="Target",
                                line=dict(width=2, dash="dot")
                            ))

                            # Forecast line (reference)
                            fig_trend.add_hline(
                                y=month_end_forecast,
                                line_dash="dash",
                                annotation_text=f"Month-End Forecast: KD {month_end_forecast:,.0f}",
                                annotation_position="top right"
                            )

                            fig_trend.update_layout(
                                height=420,
                                template="plotly_white",
                                margin=dict(l=10, r=10, t=35, b=10),
                                xaxis_title="Period",
                                yaxis_title="Net Value (KD)",
                                hovermode="x unified"
                            )
                            fig_trend.update_xaxes(type="date", tickformat="%d-%b")
                            st.plotly_chart(fig_trend, use_container_width=True)

                            with st.expander("📋 Trend Summary Table", expanded=False):
                                view_tbl = df_period.copy()
                                view_tbl["Ach%"] = np.where(view_tbl["Target"] > 0, (view_tbl["Sales"] / view_tbl["Target"]) * 100, 0)
                                view_tbl["Gap"] = view_tbl["Sales"] - view_tbl["Target"]
                                st.dataframe(view_tbl[["Period", "Sales", "Target", "Ach%", "Gap"]], use_container_width=True)

                        with c2:
                            st.markdown("### 🍕 Channel Split (Retail vs E-com)")

                            # ✅ Donut pie
                            fig_pizza = go.Figure(go.Pie(
                                labels=["Retail", "E-com"],
                                values=[total_retail, total_ecom],
                                hole=0.62,
                                sort=False,
                                direction="clockwise",
                                rotation=210,
                                pull=[0.05, 0.10],
                                textinfo="percent+label",
                                textposition="outside",
                                marker=dict(line=dict(color="white", width=3))
                            ))

                            fig_pizza.update_layout(
                                height=300,
                                template="plotly_white",
                                margin=dict(l=10, r=10, t=35, b=10),
                                showlegend=False,
                                annotations=[dict(
                                    text=f"Total<br><b>KD {total_all:,.0f}</b>",
                                    x=0.5, y=0.5, showarrow=False, font=dict(size=14)
                                )]
                            )
                            st.plotly_chart(fig_pizza, use_container_width=True)

                            # ✅ Channel values table (Retail, E-com, Total)
                            ch_tbl = pd.DataFrame({
                                "Channel": ["Retail", "E-com", "TOTAL"],
                                "Sales (KD)": [total_retail, total_ecom, total_all],
                                "Share %": [retail_share, ecom_share, 100.0 if total_all > 0 else 0.0]
                            })
                            st.dataframe(ch_tbl, use_container_width=True, hide_index=True)

                            # ✅ Optional mini bar for clarity (clean)
                            fig_bar = go.Figure(go.Bar(
                                x=["Retail", "E-com", "TOTAL"],
                                y=[total_retail, total_ecom, total_all],
                                text=[f"KD {total_retail:,.0f}", f"KD {total_ecom:,.0f}", f"KD {total_all:,.0f}"],
                                textposition="outside",
                                name="Sales"
                            ))
                            fig_bar.update_layout(
                                height=220,
                                template="plotly_white",
                                margin=dict(l=10, r=10, t=10, b=10),
                                xaxis_title="",
                                yaxis_title="KD"
                            )
                            st.plotly_chart(fig_bar, use_container_width=True)

                        st.markdown("---")

                        # ---------------------------------------------------
                        # 5️⃣ Channel Trend over time (stacked area) ✅ Retail rename
                        # ---------------------------------------------------
                        st.markdown(f"### 📊 {gran} Channel Trend (Retail vs E-com)")

                        tx = df.merge(channels_df[["PY Name 1", "Channels"]], on="PY Name 1", how="left")
                        tx["Channels"] = (
                            tx["Channels"].astype(str).str.strip().str.lower()
                            .replace({"": "retail", "nan": "retail", "none": "retail", "market": "retail"})
                        )
                        tx["Ch2"] = np.where(tx["Channels"] == "e-com", "E-com", "Retail")

                        ch_period = (tx.groupby([pd.Grouper(key="Billing Date", freq=freq), "Ch2"])["Net Value"]
                                    .sum()
                                    .reset_index()
                                    .rename(columns={"Billing Date": "Period"}))

                        pivot = ch_period.pivot(index="Period", columns="Ch2", values="Net Value").fillna(0.0)
                        for col in ["Retail", "E-com"]:
                            if col not in pivot.columns:
                                pivot[col] = 0.0
                        pivot = pivot[["Retail", "E-com"]].reset_index().sort_values("Period")

                        fig_area = go.Figure()
                        fig_area.add_trace(go.Scatter(
                            x=pivot["Period"], y=pivot["Retail"], mode="lines", name="Retail", stackgroup="one"
                        ))
                        fig_area.add_trace(go.Scatter(
                            x=pivot["Period"], y=pivot["E-com"], mode="lines", name="E-com", stackgroup="one"
                        ))
                        fig_area.update_layout(
                            height=420,
                            template="plotly_white",
                            margin=dict(l=10, r=10, t=35, b=10),
                            xaxis_title="Period",
                            yaxis_title="Net Value (KD)",
                            hovermode="x unified"
                        )
                        fig_area.update_xaxes(type="date", tickformat="%d-%b")
                        st.plotly_chart(fig_area, use_container_width=True)

                        st.markdown("---")

                        # ---------------------------------------------------
                        # 6️⃣ Salesman Performance (Achievement %, Gap, Contribution)
                        # ---------------------------------------------------
                        st.markdown("### 💪 Salesman Performance (Analytical)")

                        salesman_col = None
                        for c in ["Driver Name EN", "Salesman", "Sales Rep", "Salesperson"]:
                            if c in df.columns:
                                salesman_col = c
                                break

                        if salesman_col is None:
                            st.info("Salesman column not found, skipping salesman charts.")
                        else:
                            sales_by_sm = df.groupby(salesman_col)["Net Value"].sum().sort_values(ascending=False)

                            if (target_df is not None) and (salesman_col in target_df.columns) and ("KA Target" in target_df.columns):
                                tgt = target_df[[salesman_col, "KA Target"]].copy()
                                tgt["KA Target"] = pd.to_numeric(tgt["KA Target"], errors="coerce").fillna(0.0)

                                sm = (sales_by_sm.reset_index()
                                    .rename(columns={salesman_col: "Salesman", "Net Value": "Sales"})
                                    .merge(tgt.rename(columns={salesman_col: "Salesman"}), on="Salesman", how="left")
                                    .fillna({"KA Target": 0.0}))
                            else:
                                sm = sales_by_sm.reset_index()
                                sm.columns = ["Salesman", "Sales"]
                                sm["KA Target"] = 0.0

                            sm["Achievement %"] = np.where(sm["KA Target"] > 0, (sm["Sales"] / sm["KA Target"]) * 100, 0.0)
                            sm["Gap"] = sm["Sales"] - sm["KA Target"]
                            total_sm_sales = float(sm["Sales"].sum()) if float(sm["Sales"].sum()) != 0 else 1.0
                            sm["Contribution %"] = (sm["Sales"] / total_sm_sales) * 100

                            sm_view = sm.sort_values("Sales", ascending=False).head(20)

                            cc1, cc2 = st.columns([0.55, 0.45], gap="large")

                            with cc1:
                                st.markdown("#### 🎯 Achievement % (Top 20)")
                                fig_ach = go.Figure(go.Bar(
                                    x=sm_view["Salesman"],
                                    y=sm_view["Achievement %"],
                                    text=[f"{v:.1f}%" for v in sm_view["Achievement %"]],
                                    textposition="outside",
                                    name="Achievement %"
                                ))
                                fig_ach.update_layout(
                                    height=420,
                                    template="plotly_white",
                                    margin=dict(l=10, r=10, t=25, b=10),
                                    yaxis_title="Achievement %",
                                    xaxis_title="Salesman"
                                )
                                st.plotly_chart(fig_ach, use_container_width=True)

                            with cc2:
                                st.markdown("#### 📌 Top 3 / Bottom 3 Quick View")
                                top3 = sm.sort_values("Sales", ascending=False).head(3)[["Salesman", "Sales", "Achievement %"]]
                                bot3 = sm.sort_values("Sales", ascending=True).head(3)[["Salesman", "Sales", "Achievement %"]]

                                st.write("✅ Top 3")
                                st.dataframe(top3, use_container_width=True, hide_index=True)
                                st.write("⚠️ Bottom 3")
                                st.dataframe(bot3, use_container_width=True, hide_index=True)

                            with st.expander("📋 Salesman Detail Table (All)", expanded=False):
                                sm_tbl = sm.sort_values("Sales", ascending=False)[
                                    ["Salesman", "Sales", "KA Target", "Achievement %", "Gap", "Contribution %"]
                                ]
                                st.dataframe(sm_tbl, use_container_width=True, hide_index=True)

                        st.markdown("---")

                        # ---------------------------------------------------
                        # 7️⃣ Key Drivers
                        # ---------------------------------------------------
                        st.markdown("### 🏆 Key Drivers")

                        d1, d2 = st.columns(2, gap="large")

                        with d1:
                            st.markdown("#### 👤 Top 10 Customers by Sales")
                            if "PY Name 1" not in df.columns:
                                st.error("⚠️ 'PY Name 1' column not found!")
                            else:
                                top10_c = (df.groupby("PY Name 1")["Net Value"]
                                        .sum().sort_values(ascending=False).head(10)
                                        .reset_index()
                                        .rename(columns={"PY Name 1": "Customer", "Net Value": "Sales"}))
                                top10_c = top10_c.sort_values("Sales", ascending=True)

                                fig_top10c = go.Figure(go.Bar(
                                    x=top10_c["Sales"],
                                    y=top10_c["Customer"],
                                    orientation="h",
                                    text=[f"KD {v:,.0f}" for v in top10_c["Sales"]],
                                    textposition="outside"
                                ))
                                fig_top10c.update_layout(
                                    height=450,
                                    template="plotly_white",
                                    margin=dict(l=10, r=10, t=25, b=10),
                                    xaxis_title="Net Value (KD)",
                                    yaxis_title=""
                                )
                                st.plotly_chart(fig_top10c, use_container_width=True)

                        with d2:
                            st.markdown("#### 🧾 Top 10 SKU by Sales (optional)")
                            sku_col = None
                            for c in ["Material Description", "Material", "SKU", "Material Code"]:
                                if c in df.columns:
                                    sku_col = c
                                    break

                            if sku_col is None:
                                st.info("SKU column not found (Material Description / Material / SKU).")
                            else:
                                top10_sku = (df.groupby(sku_col)["Net Value"]
                                            .sum().sort_values(ascending=False).head(10)
                                            .reset_index()
                                            .rename(columns={sku_col: "SKU", "Net Value": "Sales"}))
                                top10_sku = top10_sku.sort_values("Sales", ascending=True)

                                fig_top10s = go.Figure(go.Bar(
                                    x=top10_sku["Sales"],
                                    y=top10_sku["SKU"],
                                    orientation="h",
                                    text=[f"KD {v:,.0f}" for v in top10_sku["Sales"]],
                                    textposition="outside"
                                ))
                                fig_top10s.update_layout(
                                    height=450,
                                    template="plotly_white",
                                    margin=dict(l=10, r=10, t=25, b=10),
                                    xaxis_title="Net Value (KD)",
                                    yaxis_title=""
                                )
                                st.plotly_chart(fig_top10s, use_container_width=True)

                    # --- DOWNLOADS ---
                    with tabs[3]:
                        st.subheader(texts[lang]["download_reports_sub"])
                        col1, col2 = st.columns(2)

                        # --------------------------
                        # COL 1: PPTX GENERATION
                        # --------------------------
                        with col1:
                            if st.button(texts[lang]["generate_pptx"]):

                                # ✅ FIX: Build figs_dict safely (avoid NameError if any fig not created)
                                figs_dict = {}

                                if "fig_trend_new" in globals() and fig_trend_new is not None:
                                    figs_dict["Daily Sales Trend"] = fig_trend_new

                                if "fig_channel_new" in globals() and fig_channel_new is not None:
                                    figs_dict["Market vs E-com Sales"] = fig_channel_new

                                if "fig_target_new" in globals() and fig_target_new is not None:
                                    figs_dict["Daily KA Target vs Actual"] = fig_target_new

                                if "fig_salesman_new" in globals() and fig_salesman_new is not None:
                                    figs_dict["Salesman KA Target vs Actual"] = fig_salesman_new

                                if "fig_top10_new" in globals() and fig_top10_new is not None:
                                    figs_dict["Top 10 Customers by Sales"] = fig_top10_new

                                talabat_ppt_tables = {
                                    "billing_split": talabat_billing_split if "talabat_billing_split" in locals() else pd.DataFrame(),
                                    "customers": talabat_customer_table.head(25) if ("talabat_customer_table" in locals() and isinstance(talabat_customer_table, pd.DataFrame)) else pd.DataFrame(),
                                }

                                pptx_stream = create_pptx(
                                    report_df_with_total,
                                    billing_df,
                                    py_table_with_total,
                                    figs_dict,
                                    kpi_data,
                                    talabat_tables=talabat_ppt_tables
                                )

                                st.download_button(
                                    texts[lang]["download_pptx"],
                                    pptx_stream,
                                    file_name=f"sales_report_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )

                                # Audit safe
                                if "audit_log" in st.session_state:
                                    st.session_state["audit_log"].append({
                                        "user": username,
                                        "action": "download",
                                        "details": "PPTX Report",
                                        "timestamp": datetime.now()
                                    })

                        # --------------------------
                        # COL 2: EXCEL DOWNLOADS
                        # --------------------------
                        with col2:
                            # ==============================
                            # 🎯 KA Target Daily Sheet (Full Month) - DOWNLOAD
                            # ==============================
                            st.markdown("### 🎯 KA Target – Daily Sheet (Full Month)")

                            # Safety checks (avoid crash if variables missing)
                            if "df_filtered" not in locals() or not isinstance(df_filtered, pd.DataFrame) or df_filtered.empty:
                                st.info("No data available to generate KA Target Daily file.")
                            else:
                                # Build month days
                                month_start_ds = current_month_start
                                month_end_ds = current_month_end
                                days_ds = pd.date_range(month_start_ds, month_end_ds, freq="D")

                                # Completed days: up to today (future days stay blank)
                                cutoff = min(today, month_end_ds)

                                # Prepare working df for daily totals
                                df_ds = df_filtered.copy()
                                df_ds["__date"] = pd.to_datetime(df_ds["Billing Date"], errors="coerce").dt.normalize()

                                # --- Sales type masks (adjust if your billing codes differ) ---
                                hht_mask = df_ds["Billing Type"].astype(str).str.upper().isin(["YKF2", "HHT"])
                                presales_mask = df_ds["Billing Type"].astype(str).str.upper().isin(["ZFR", "PRESALES"])
                                tal_mask = df_ds["PY Name 1"].astype(str).str.strip().str.upper() == str(TALABAT_PY).strip().upper()

                                # Group daily sums
                                g_hht = df_ds.loc[hht_mask].groupby(["Driver Name EN", "__date"])["Net Value"].sum()
                                g_pre = df_ds.loc[presales_mask].groupby(["Driver Name EN", "__date"])["Net Value"].sum()
                                g_tal = df_ds.loc[tal_mask].groupby(["Driver Name EN", "__date"])["Net Value"].sum()

                                def _get(g, sm, d):
                                    try:
                                        return float(g.get((sm, d), 0.0))
                                    except Exception:
                                        return 0.0

                                # Day columns: Date 1..Date N (full month)
                                day_cols = [f"Date {i}" for i in range(1, len(days_ds) + 1)]

                                # Salesmen list
                                all_salesmen_month = sorted(df_ds["Driver Name EN"].dropna().unique().tolist())

                                rows = []
                                for sm in all_salesmen_month:
                                    ka_t = float(ka_targets.get(sm, 0.0)) if "ka_targets" in locals() else 0.0

                                    # Achieved KA (exclude Talabat)
                                    sm_total = float(df_ds.loc[df_ds["Driver Name EN"] == sm, "Net Value"].sum())
                                    sm_tal = float(df_ds.loc[(df_ds["Driver Name EN"] == sm) & tal_mask, "Net Value"].sum())
                                    achieved_ka = sm_total - sm_tal

                                    per_day_target = (ka_t / working_days_current_month) if ("working_days_current_month" in locals() and working_days_current_month) else 0.0
                                    current_sales_per_day = (achieved_ka / days_finish) if ("days_finish" in locals() and days_finish) else 0.0
                                    balance_vs_target = ka_t - achieved_ka

                                    # Monthly totals per type
                                    sum_hht = float(df_ds.loc[(df_ds["Driver Name EN"] == sm) & hht_mask, "Net Value"].sum())
                                    sum_pre = float(df_ds.loc[(df_ds["Driver Name EN"] == sm) & presales_mask, "Net Value"].sum())
                                    sum_tot = float(df_ds.loc[(df_ds["Driver Name EN"] == sm) & (hht_mask | presales_mask), "Net Value"].sum())
                                    sum_tal2 = float(df_ds.loc[(df_ds["Driver Name EN"] == sm) & tal_mask, "Net Value"].sum())

                                    def _add_row(label, daily_fn, summary_val):
                                        row = {
                                            "KA Target": "KA Target",
                                            "KA Target Value": round(ka_t, 0),
                                            "Per day Target (Total target / Working days)": round(per_day_target, 0),
                                            "Achieved value": round(achieved_ka, 0),
                                            "Current Sales Per day": round(current_sales_per_day, 0),
                                            "Salesman Name": sm,
                                            "Sales Type": label,
                                            "Sales Summary": round(summary_val, 0),
                                            "Balance": round(balance_vs_target, 0),
                                        }
                                        for idx_d, d in enumerate(days_ds, start=1):
                                            col = f"Date {idx_d}"
                                            row[col] = round(daily_fn(sm, d), 0) if d <= cutoff else ""
                                        rows.append(row)

                                    _add_row("HHT", lambda s, d: _get(g_hht, s, d), sum_hht)
                                    _add_row("PRESALES", lambda s, d: _get(g_pre, s, d), sum_pre)
                                    _add_row("Sales Total", lambda s, d: (_get(g_hht, s, d) + _get(g_pre, s, d)), sum_tot)
                                    _add_row("Talabat Sales", lambda s, d: _get(g_tal, s, d), sum_tal2)

                                ka_daily_df = pd.DataFrame(rows)

                                # Order columns nicely
                                base_cols = [
                                    "KA Target", "KA Target Value",
                                    "Per day Target (Total target / Working days)",
                                    "Achieved value", "Current Sales Per day",
                                    "Salesman Name", "Sales Type", "Sales Summary", "Balance"
                                ]
                                ka_daily_df = ka_daily_df.reindex(columns=base_cols + day_cols)

                                ka_daily_bytes = to_excel_bytes(ka_daily_df, sheet_name="KA_Target_Daily", index=False)

                                if st.download_button(
                                    "⬇️ Download KA Target Daily (Excel)",
                                    data=ka_daily_bytes,
                                    file_name=f"KA_Target_Daily_{datetime.now().strftime('%Y-%m')}.xlsx",
                                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                                ):
                                    if "audit_log" in st.session_state:
                                        st.session_state["audit_log"].append({
                                            "user": username,
                                            "action": "download",
                                            "details": "KA Target Daily (Excel)",
                                            "timestamp": datetime.now()
                                        })



elif choice == texts[lang]["sales_tracking"]:
    st.title(texts[lang]["sales_tracking_title"])
    if "data_loaded" not in st.session_state:
        st.warning(texts[lang]["no_data_warning"])
    else:
        # Filters
        st.sidebar.subheader(texts[lang]["filters_header"])
        st.sidebar.markdown(f'<div class="tooltip">ℹ️<span class="tooltiptext">{texts[lang]["filters_tooltip"]}</span></div>', unsafe_allow_html=True)
        salesmen = st.sidebar.multiselect(
            texts[lang]["select_salesmen"],
            options=sorted(sales_df["Driver Name EN"].dropna().unique()),
            default=sorted(sales_df["Driver Name EN"].dropna().unique()),
            key="st_salesmen"
        )
        billing_types = st.sidebar.multiselect(
            texts[lang]["select_billing_types"],
            options=sorted(sales_df["Billing Type"].dropna().unique()),
            default=sorted(sales_df["Billing Type"].dropna().unique()),
            key="st_billing_types"
        )
        py_filter = st.sidebar.multiselect(
            texts[lang]["select_py"],
            options=sorted(sales_df["PY Name 1"].dropna().unique()),
            default=sorted(sales_df["PY Name 1"].dropna().unique()),
            key="st_py_filter"
        )
        sp_filter = st.sidebar.multiselect(
            texts[lang]["select_sp"],
            options=sorted(sales_df["SP Name1"].dropna().unique()),
            default=sorted(sales_df["SP Name1"].dropna().unique()),
            key="st_sp_filter"
        )

        preset = st.sidebar.radio(texts[lang]["date_presets"], texts[lang]["date_presets_options"], key="st_preset")
        today = pd.Timestamp.today().normalize()
        if preset == texts[lang]["date_presets_options"][1]:  # Last 7 Days
            date_range = [today - pd.Timedelta(days=7), today]
        elif preset == texts[lang]["date_presets_options"][2]:  # This Month
            month_start = today.replace(day=1)
            month_end = month_start + pd.offsets.MonthEnd(0)
            date_range = [month_start, month_end]
        elif preset == texts[lang]["date_presets_options"][3]:  # YTD
            date_range = [today.replace(month=1, day=1), today]
        else:
            date_range = st.sidebar.date_input(
                texts[lang]["select_date_range"],
                [sales_df["Billing Date"].min(), sales_df["Billing Date"].max()],
                key="st_date_range"
            )
            if isinstance(date_range, tuple) and len(date_range) == 2:
                date_range = [pd.to_datetime(date_range[0]), pd.to_datetime(date_range[1])]
            else:
                date_range = [sales_df["Billing Date"].min(), sales_df["Billing Date"].max()]

        if date_range[0] > date_range[1]:
            st.error(texts[lang]["date_error"])
        else:
            top_n = st.sidebar.slider(
                texts[lang]["top_n_salesmen"],
                min_value=1,
                max_value=max(1, len(sales_df["Driver Name EN"].dropna().unique())),
                value=min(5, max(1, len(sales_df["Driver Name EN"].dropna().unique()))),
                key="st_topn"
            )

            df_filtered = sales_df[
                (sales_df["Driver Name EN"].isin(salesmen))
                & (sales_df["Billing Type"].isin(billing_types))
                & (sales_df["Billing Date"] >= date_range[0])
                & (sales_df["Billing Date"] <= date_range[1])
                & (sales_df["PY Name 1"].isin(py_filter))
                & (sales_df["SP Name1"].isin(sp_filter))
            ].copy()

            if df_filtered.empty:
                st.warning(texts[lang]["no_match_warning"])
            else:
                billing_start = df_filtered["Billing Date"].min().normalize()
                billing_end = df_filtered["Billing Date"].max().normalize()
                all_days = pd.date_range(billing_start, billing_end, freq="D")
                days_finish = int(sum(1 for d in all_days if d.weekday() != 4))

                current_month_start = today.replace(day=1)
                current_month_end = current_month_start + pd.offsets.MonthEnd(0)
                month_days = pd.date_range(current_month_start, current_month_end, freq="D")
                working_days_current_month = int(sum(1 for d in month_days if d.weekday() != 4))

                # --- Base aggregates ---
                total_sales = df_filtered.groupby("Driver Name EN")["Net Value"].sum()
                talabat_df = df_filtered[df_filtered["PY Name 1"] == "STORES SERVICES KUWAIT CO."]
                talabat_sales = talabat_df.groupby("Driver Name EN")["Net Value"].sum()

                ka_targets = target_df.set_index("Driver Name EN")["KA Target"] if "KA Target" in target_df.columns else pd.Series(dtype=float)
                talabat_targets = target_df.set_index("Driver Name EN")["Talabat Target"] if "Talabat Target" in target_df.columns else pd.Series(dtype=float)

                all_salesmen_idx = total_sales.index.union(talabat_sales.index).union(ka_targets.index).union(talabat_targets.index)
                total_sales = total_sales.reindex(all_salesmen_idx, fill_value=0).astype(float)
                talabat_sales = talabat_sales.reindex(all_salesmen_idx, fill_value=0).astype(float)
                ka_targets = ka_targets.reindex(all_salesmen_idx, fill_value=0).astype(float)
                talabat_targets = talabat_targets.reindex(all_salesmen_idx, fill_value=0).astype(float)

                ka_gap = (ka_targets - total_sales).clip(lower=0)
                talabat_gap = (talabat_targets - talabat_sales).clip(lower=0)

                top_salesmen = total_sales.sort_values(ascending=False).head(top_n).index
                total_sales_top = total_sales.loc[top_salesmen]
                talabat_sales_top = talabat_sales.loc[top_salesmen]
                ka_gap_top = ka_gap.loc[top_salesmen]
                talabat_gap_top = talabat_gap.loc[top_salesmen]

                total_ka_target_all = float(ka_targets.sum())
                total_tal_target_all = float(talabat_targets.sum())
                per_day_ka_target = (total_ka_target_all / working_days_current_month) if working_days_current_month > 0 else 0
                current_sales_per_day = (total_sales.sum() / days_finish) if days_finish > 0 else 0
                forecast_month_end_ka = current_sales_per_day * working_days_current_month

                # --- Channels mapping: Market (Retail) vs E-com ---
                df_py_sales = df_filtered.groupby("_py_name_norm")["Net Value"].sum().reset_index()
                df_channels_merged = df_py_sales.merge(
                    channels_df[["_py_name_norm", "Channels"]],
                    on="_py_name_norm",
                    how="left"
                )
                df_channels_merged["Channels"] = df_channels_merged["Channels"].str.strip().str.lower().fillna("uncategorized")
                channel_sales = df_channels_merged.groupby("Channels")["Net Value"].sum()
                total_retail_sales = float(channel_sales.get("market", 0.0) + channel_sales.get("uncategorized", 0.0))
                total_ecom_sales = float(channel_sales.get("e-com", 0.0))
                total_channel_sales = total_retail_sales + total_ecom_sales
                retail_sales_pct = (total_retail_sales / total_channel_sales * 100) if total_channel_sales > 0 else 0
                ecom_sales_pct = (total_ecom_sales / total_channel_sales * 100) if total_channel_sales > 0 else 0

                # --- KA & Other E-com Calculation ---
                ka_other_ecom_sales = total_sales.sum() - talabat_sales.sum()
                ka_other_ecom_pct = (ka_other_ecom_sales / total_ka_target_all * 100) if total_ka_target_all > 0 else 0

                # --- KPI Data for PPTX ---
                kpi_data = {
                    texts[lang]["ka_target"]: f"KD {total_ka_target_all:,.0f}",
                    texts[lang]["talabat_target"]: f"KD {total_tal_target_all:,.0f}",
                    texts[lang]["ka_gap"]: f"KD {(total_ka_target_all - total_sales.sum()):,.0f}",
                    "Total Talabat Gap": f"KD {talabat_gap.sum():,.0f}",
                    texts[lang]["total_ka_sales"]: f"KD {total_sales.sum():,.0f} ({((total_sales.sum() / total_ka_target_all) * 100):.0f}%)" if total_ka_target_all else f"KD {total_sales.sum():,.0f} (0%)",
                    "Total Talabat Sales": f"KD {talabat_sales.sum():,.0f} ({((talabat_sales.sum() / total_tal_target_all) * 100):.0f}%)" if total_tal_target_all else f"KD {talabat_sales.sum():,.0f} (0%)",
                    texts[lang]["ka_other_ecom"]: f"KD {ka_other_ecom_sales:,.0f} ({ka_other_ecom_pct:.0f}%)",
                    texts[lang]["retail_sales"]: f"KD {total_retail_sales:,.0f} ({retail_sales_pct:.0f}%)",
                    texts[lang]["ecom_sales"]: f"KD {total_ecom_sales:,.0f} ({ecom_sales_pct:.0f}%)",
                    texts[lang]["days_finished"]: f"{days_finish}",
                    "Per Day KA Target": f"KD {per_day_ka_target:,.0f}",
                    texts[lang]["current_sales_per_day"]: f"KD {current_sales_per_day:,.0f}",
                    texts[lang]["forecast_month_end"]: f"KD {forecast_month_end_ka:,.0f}"
                }

                tabs = st.tabs([texts[lang]["kpis_tab"], texts[lang]["tables_tab"], texts[lang]["charts_tab"], texts[lang]["downloads_tab"]])

                # --- KPIs with progress bars ---
                with tabs[0]:
                    st.subheader(texts[lang]["key_metrics_sub"])
                    r1c1 = st.columns(1)[0]
                    with r1c1:
                        st.metric(texts[lang]["total_ka_sales"], f"KD {total_sales.sum():,.0f}")
                        progress_pct_ka = (total_sales.sum() / total_ka_target_all * 100) if total_ka_target_all > 0 else 0
                        st.markdown(create_progress_bar_html(progress_pct_ka), unsafe_allow_html=True)
                        st.markdown(f'<div class="green-caption">{texts[lang]["of_ka_target"].format(progress_pct_ka)}</div>', unsafe_allow_html=True)

                    r2c1, r2c2 = st.columns(2)
                    with r2c1:
                        st.metric(texts[lang]["ka_other_ecom"], f"KD {ka_other_ecom_sales:,.0f}")
                        st.markdown(create_progress_bar_html(ka_other_ecom_pct), unsafe_allow_html=True)
                        st.markdown(f'<div class="green-caption">{texts[lang]["of_ka_target_pct"].format(ka_other_ecom_pct)}</div>', unsafe_allow_html=True)
                    with r2c2:
                        st.metric(texts[lang]["talabat_sales"], f"KD {talabat_sales.sum():,.0f}")
                        progress_pct_talabat = (talabat_sales.sum() / total_tal_target_all * 100) if total_tal_target_all > 0 else 0
                        st.markdown(create_progress_bar_html(progress_pct_talabat), unsafe_allow_html=True)
                        st.markdown(f'<div class="green-caption">{texts[lang]["of_talabat_target"].format(progress_pct_talabat)}</div>', unsafe_allow_html=True)

                    st.subheader(texts[lang]["target_overview_sub"])
                    r3c1, r3c2, r3c3, r3c4 = st.columns(4)
                    r3c1.metric(texts[lang]["ka_target"], f"KD {total_ka_target_all:,.0f}")
                    r3c2.metric(texts[lang]["talabat_target"], f"KD {total_tal_target_all:,.0f}")
                    r3c3.metric(texts[lang]["ka_gap"], f"KD {(total_ka_target_all - total_sales.sum()):,.0f}")
                    r3c4.metric(texts[lang]["talabat_gap"], f"KD {talabat_gap.sum():,.0f}")

                    st.subheader(texts[lang]["channel_sales_sub"])
                    r4c1, r4c2 = st.columns(2)
                    with r4c1:
                        st.metric(texts[lang]["retail_sales"], f"KD {total_retail_sales:,.0f}")
                        retail_contribution_pct = (total_retail_sales / total_sales.sum() * 100) if total_sales.sum() > 0 else 0
                        st.caption(texts[lang]["of_total_ka"].format(retail_contribution_pct))
                    with r4c2:
                        st.metric(texts[lang]["ecom_sales"], f"KD {total_ecom_sales:,.0f}")
                        ecom_contribution_pct = (total_ecom_sales / total_sales.sum() * 100) if total_sales.sum() > 0 else 0
                        st.caption(texts[lang]["of_total_ka"].format(ecom_contribution_pct))

                    st.subheader(texts[lang]["performance_metrics_sub"])
                    r5c1, r5c2, r5c3 = st.columns(3)
                    r5c1.metric(texts[lang]["days_finished"], days_finish)
                    r5c2.metric(texts[lang]["current_sales_per_day"], f"KD {current_sales_per_day:,.0f}")
                    r5c3.metric(texts[lang]["forecast_month_end"], f"KD {forecast_month_end_ka:,.0f}")

                    # --- TABLES ---
                    with tabs[1]:
                        # --- Sales & Targets Summary ---
                        st.subheader(texts[lang]["sales_targets_summary_sub"])
                        report_df = pd.DataFrame({
                            "Salesman": ka_targets.index,
                            "KA Target": ka_targets.values,
                            "KA Sales": total_sales.values,
                            "KA Remaining": ka_gap.values,
                            "KA % Achieved": np.where(ka_targets.values != 0, (total_sales.values / ka_targets.values * 100).round(0), 0),
                            "Talabat Target": talabat_targets.values,
                            "Talabat Sales": talabat_sales.values,
                            "Talabat Remaining": talabat_gap.values,
                            "Talabat % Achieved": np.where(talabat_targets.values != 0, (talabat_sales.values / talabat_targets.values * 100).round(0), 0)
                        })

                        ttotal_row = report_df.sum(numeric_only=True).to_frame().T
                        total_row.index = ["Total"]

                        total_row["Target % Achieved"] = round(
                            total_row["Sales"]/total_row["Target"]*100,0
                        ) if total_row["Target"].iloc[0] != 0 else 0

                        total_row["E-Com % Achieved"] = round(
                            total_row["E-Com Sales"]/total_row["E-Com Target"]*100,0
                        ) if total_row["E-Com Target"].iloc[0] != 0 else 0

                        total_row = total_row.reset_index(drop=True)
                        total_row["Salesman Name"] = "Total"

                        total_row = total_row[report_df.columns]
                        report_df_with_total = pd.concat([report_df, total_row], ignore_index=True)

                        report_df_with_total = pd.concat([report_df, total_row], ignore_index=True)

                        def highlight_total_row(row):
                            return ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900' if row['Salesman'] == "Total" else '' for _ in row]

                        styled_report = (
                            report_df_with_total.style
                            .set_table_styles([
                                {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                            ('font-weight', '800'), ('height', '40px'),
                                                            ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                            ])
                            .apply(highlight_total_row, axis=1)
                            .format("{:,.0f}", subset=["KA Target","KA Sales","KA Remaining",
                                                    "Talabat Target","Talabat Sales","Talabat Remaining"])
                            .format("{:.0f}%", subset=["KA % Achieved","Talabat % Achieved"])
                        )
                        st.dataframe(styled_report, use_container_width=True, hide_index=True)
                        if st.download_button(
                            texts[lang]["download_sales_targets"],
                            data=to_excel_bytes(report_df_with_total, sheet_name="Sales_Targets_Summary", index=False),
                            file_name=f"Sales_Targets_Summary_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        ):
                            st.session_state["audit_log"].append({
                                "user": username,
                                "action": "download",
                                "details": "Sales & Targets Summary Excel",
                                "timestamp": datetime.now()
                            })

                        # --- Sales by Billing Type per Salesman ---
                        st.subheader(texts[lang]["sales_by_billing_sub"])
                        billing_wide = df_filtered.pivot_table(
                            index="Driver Name EN",
                            columns="Billing Type",
                            values="Net Value",
                            aggfunc="sum",
                            fill_value=0
                        )

                        required_cols_raw = ["ZFR", "YKF2", "YKRE", "YKS1", "YKS2", "ZCAN", "ZRE"]
                        billing_wide = billing_wide.reindex(columns=required_cols_raw, fill_value=0)
                        display_df = billing_wide.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})
                        display_df["Sales Total"] = billing_wide.sum(axis=1)
                        display_df["Return"] = billing_wide["YKRE"] + billing_wide["ZRE"]
                        display_df["Return %"] = np.where(display_df["Sales Total"] != 0,
                                                        (display_df["Return"] / display_df["Sales Total"] * 100).round(0), 0)
                        display_df["Cancel Total"] = billing_wide[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)
                        ordered_cols = ["Presales", "HHT", "Sales Total", "YKS1", "YKS2", "ZCAN",
                                        "Cancel Total", "YKRE", "ZRE", "Return", "Return %"]
                        display_df = display_df.reindex(columns=ordered_cols, fill_value=0)

                        total_row = pd.DataFrame(display_df.sum(numeric_only=True)).T
                        total_row.index = ["Total"]
                        total_row["Return %"] = round((total_row["Return"] / total_row["Sales Total"] * 100), 0) if total_row["Sales Total"].iloc[0] != 0 else 0
                        billing_df = pd.concat([display_df, total_row])
                        billing_df.index.name = "Salesman"

                        def highlight_total_row_billing(row):
                            return ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900' if row.name == "Total" else '' for _ in row]

                        styled_billing = (
                            billing_df.style
                            .set_table_styles([
                                {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                            ('font-weight', '800'), ('height', '40px'),
                                                            ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                            ])
                            .apply(highlight_total_row_billing, axis=1)
                            .format({
                                "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                                "YKS1": "{:,.0f}", "YKS2": "{:,.0f}", "ZCAN": "{:,.0f}", "Cancel Total": "{:,.0f}",
                                "YKRE": "{:,.0f}", "ZRE": "{:,.0f}", "Return": "{:,.0f}", "Return %": "{:.0f}%"
                            })
                        )
                        st.dataframe(styled_billing, use_container_width=True, hide_index=False)
                        if st.download_button(
                            texts[lang]["download_billing"],
                            data=to_excel_bytes(billing_df.reset_index(), sheet_name="Billing_Types", index=False),
                            file_name=f"Billing_Types_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        ):
                            st.session_state["audit_log"].append({
                                "user": username,
                                "action": "download",
                                "details": "Billing Type Table Excel",
                                "timestamp": datetime.now()
                            })

                        # --- Sales by PY Name 1 ---
                        st.subheader(texts[lang]["sales_by_py_sub"])
                        py_table = df_filtered.groupby("PY Name 1")["Net Value"].sum().sort_values(ascending=False).to_frame(name="Sales")
                        py_table["Contribution %"] = np.where(py_table["Sales"] != 0,
                                                            (py_table["Sales"]/py_table["Sales"].sum()*100).round(0), 0)

                        total_row = py_table.sum(numeric_only=True).to_frame().T
                        total_row.index = ["Total"]
                        py_table_with_total = pd.concat([py_table, total_row])
                        py_table_with_total.index.name = "PY Name 1"

                        def highlight_total_row_py(row):
                            return ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900' if row.name == "Total" else '' for _ in row]

                        styled_py = (
                            py_table_with_total.style
                            .set_table_styles([
                                {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                            ('font-weight', '800'), ('height', '40px'),
                                                            ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                            ])
                            .apply(highlight_total_row_py, axis=1)
                            .format("{:,.0f}", subset=["Sales"])
                            .format("{:.0f}%", subset=["Contribution %"])
                        )
                        st.dataframe(styled_py, use_container_width=True, hide_index=False)
                        if st.download_button(
                            texts[lang]["download_py"],
                            data=to_excel_bytes(py_table_with_total.reset_index(), sheet_name="Sales_by_PY_Name", index=False),
                            file_name=f"Sales_by_PY_Name_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        ):
                            st.session_state["audit_log"].append({
                                "user": username,
                                "action": "download",
                                "details": "PY Name Table Excel",
                                "timestamp": datetime.now()
                            })

                        # --- Return by PY Name 1 ---
                        st.subheader("🔄 Return by PY Name 1")
                        py_billing = df_filtered.pivot_table(
                            index="PY Name 1",
                            columns="Billing Type",
                            values="Net Value",
                            aggfunc="sum",
                            fill_value=0
                        )

                        py_billing = py_billing.reindex(columns=required_cols_raw, fill_value=0)
                        py_billing["Sales Total"] = py_billing.sum(axis=1)
                        py_billing["Return"] = py_billing["YKRE"] + py_billing["ZRE"]
                        py_billing["Cancel Total"] = py_billing[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)

                        # Reorder columns like Sales by Billing table
                        ordered_cols = ["Presales", "HHT", "Sales Total", "YKS1", "YKS2", "ZCAN",
                                        "Cancel Total", "YKRE", "ZRE", "Return", "Return %"]
                        # Map the original column names to the ordered display names
                        py_billing = py_billing.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})
                        py_billing["Return %"] = np.where(py_billing["Sales Total"] != 0,
                                                        (py_billing["Return"] / py_billing["Sales Total"] * 100).round(0), 0)
                        py_billing = py_billing.reindex(columns=ordered_cols, fill_value=0).astype(int)

                        # Add total row
                        total_row = pd.DataFrame(py_billing.sum(numeric_only=True)).T
                        total_row.index = ["Total"]
                        total_row["Return %"] = round((total_row["Return"]/total_row["Sales Total"]*100), 0) if total_row["Sales Total"].iloc[0]!=0 else 0
                        py_billing = pd.concat([py_billing, total_row])

                        def highlight_total_row_py_return(row):
                            return ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900' if row.name == "Total" else '' for _ in row]

                        styled_py_return = (
                            py_billing.style
                            .set_table_styles([
                                {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                            ('font-weight', '800'), ('height', '40px'),
                                                            ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                            ])
                            .apply(highlight_total_row_py_return, axis=1)
                            .format({
                                "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                                "YKS1": "{:,.0f}", "YKS2": "{:,.0f}", "ZCAN": "{:,.0f}", "Cancel Total": "{:,.0f}",
                                "YKRE": "{:,.0f}", "ZRE": "{:,.0f}", "Return": "{:,.0f}", "Return %": "{:.0f}%"
                            })
                        )
                        st.dataframe(styled_py_return, use_container_width=True, hide_index=False)

                        # --- Return by SP Name1 ---
                        st.subheader("🔄 Return By Branch")
                        sp_billing = df_filtered.pivot_table(
                            index="SP Name1",
                            columns="Billing Type",
                            values="Net Value",
                            aggfunc="sum",
                            fill_value=0
                        )

                        sp_billing = sp_billing.reindex(columns=required_cols_raw, fill_value=0)
                        sp_billing["Sales Total"] = sp_billing.sum(axis=1)
                        sp_billing["Return"] = sp_billing["YKRE"] + sp_billing["ZRE"]
                        sp_billing["Cancel Total"] = sp_billing[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)

                        # Reorder and rename like Sales by Billing table
                        sp_billing = sp_billing.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})
                        sp_billing["Return %"] = np.where(sp_billing["Sales Total"] != 0,
                                                        (sp_billing["Return"] / sp_billing["Sales Total"] * 100).round(0), 0)
                        sp_billing = sp_billing.reindex(columns=ordered_cols, fill_value=0).astype(int)

                        # Add total row
                        total_row = pd.DataFrame(sp_billing.sum(numeric_only=True)).T
                        total_row.index = ["Total"]
                        total_row["Return %"] = round((total_row["Return"]/total_row["Sales Total"]*100), 0) if total_row["Sales Total"].iloc[0]!=0 else 0
                        sp_billing = pd.concat([sp_billing, total_row])

                        def highlight_total_row_sp_return(row):
                            return ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900' if row.name == "Total" else '' for _ in row]

                        styled_sp_return = (
                            sp_billing.style
                            .set_table_styles([
                                {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                            ('font-weight', '800'), ('height', '40px'),
                                                            ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                            ])
                            .apply(highlight_total_row_sp_return, axis=1)
                            .format({
                                "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                                "YKS1": "{:,.0f}", "YKS2": "{:,.0f}", "ZCAN": "{:,.0f}", "Cancel Total": "{:,.0f}",
                                "YKRE": "{:,.0f}", "ZRE": "{:,.0f}", "Return": "{:,.0f}", "Return %": "{:.0f}%"
                            })
                        )
                        st.dataframe(styled_sp_return, use_container_width=True, hide_index=False)

                # --- CHARTS (GM Premium Visuals Pack) ---
                with tabs[2]:
                    st.subheader("📈 Daily Sales Trend & Forecast (GM View)")
                    df_time = df_filtered.groupby(pd.Grouper(key="Billing Date", freq="D"))["Net Value"].sum().reset_index()
                    df_time.rename(columns={"Billing Date": "ds", "Net Value": "y"}, inplace=True)

                    if len(df_time) > 2:
                        # Prophet Forecast
                        m = Prophet()
                        m.fit(df_time)
                        future = m.make_future_dataframe(periods=30)
                        forecast = m.predict(future)

                        # Identify Anomalies (Based on Rolling Mean & Std)
                        df_time['y_mean'] = df_time['y'].rolling(window=7).mean()
                        df_time['y_std'] = df_time['y'].rolling(window=7).std()
                        df_time['upper'] = df_time['y_mean'] + 2 * df_time['y_std']
                        df_time['lower'] = df_time['y_mean'] - 2 * df_time['y_std']
                        df_time['anomaly'] = np.where(
                            (df_time['y'] > df_time['upper']) | (df_time['y'] < df_time['lower']),
                            df_time['y'], np.nan
                        )

                        fig_trend = go.Figure()
                        fig_trend.add_trace(go.Scatter(
                            x=df_time["ds"], y=df_time["y"],
                            mode="lines+markers",
                            name="Actual Sales",
                            line=dict(color="#0052CC", width=3)
                        ))
                        fig_trend.add_trace(go.Scatter(
                            x=forecast["ds"], y=forecast["yhat"],
                            mode="lines",
                            name="Sales Forecast",
                            line=dict(color="#22C55E", width=2, dash="dash")
                        ))
                        fig_trend.add_trace(go.Scatter(
                            x=df_time["ds"], y=df_time["anomaly"],
                            mode="markers",
                            name="Anomaly",
                            marker=dict(color="red", size=12, symbol="x")
                        ))

                        fig_trend.update_layout(
                            xaxis_title="Date",
                            yaxis_title="Net Value (KD)",
                            hovermode="x unified",
                            template="plotly_white"
                        )
                        st.plotly_chart(fig_trend, use_container_width=True)
                    else:
                        st.info("Not enough trend data")

                    st.markdown("---")
                    st.subheader("🛒 Market vs E-com Performance (Value & Share)")

                    # Market vs E-com
                    market_ecom_df = pd.DataFrame({
                        "Channel": ["Market", "E-com"],
                        "Sales": [total_retail_sales, total_ecom_sales]
                    })

                    fig_channel = go.Figure()

                    fig_channel.add_trace(go.Bar(
                        x=market_ecom_df["Channel"],
                        y=market_ecom_df["Sales"],
                        name="Sales",
                        text=market_ecom_df["Sales"].apply(lambda x: f"KD {x:,.0f}"),
                        textposition="outside",
                        marker=dict(
                            color=["#0EA5E9", "#8B5CF6"],
                            line=dict(color="black", width=1)
                        )
                    ))

                    fig_channel.add_trace(go.Pie(
                        labels=market_ecom_df["Channel"],
                        values=market_ecom_df["Sales"],
                        hole=0.55,
                        name="Share %",
                        textinfo="percent",
                        domain=dict(x=[0.55, 1.0])
                    ))

                    fig_channel.update_layout(
                        barmode="group",
                        showlegend=False,
                        template="plotly_white"
                    )
                    st.plotly_chart(fig_channel, use_container_width=True)

                    st.markdown("---")
                    st.subheader("🎯 Daily KA Target vs Actual Sales")

                    df_time_target = df_filtered.groupby(pd.Grouper(key="Billing Date", freq="D"))["Net Value"].sum().reset_index()
                    df_time_target.rename(columns={"Billing Date": "Date", "Net Value": "Sales"}, inplace=True)
                    df_time_target["Daily KA Target"] = per_day_ka_target

                    fig_target = go.Figure()
                    fig_target.add_trace(go.Scatter(
                        x=df_time_target["Date"], y=df_time_target["Sales"],
                        name="Actual Sales",
                        mode="lines+markers",
                        line=dict(color="#22C55E", width=3)
                    ))
                    fig_target.add_trace(go.Scatter(
                        x=df_time_target["Date"], y=df_time_target["Daily KA Target"],
                        name="Daily KA Target",
                        mode="lines",
                        line=dict(color="#FACC15", width=2, dash="dot")
                    ))

                    fig_target.update_layout(
                        xaxis_title="Date",
                        yaxis_title="Net Value (KD)",
                        hovermode="x unified",
                        template="plotly_white"
                    )
                    st.plotly_chart(fig_target, use_container_width=True)

                    st.markdown("---")
                    st.subheader("💪 Salesman KA Target vs Actual")

                    target_actual_df = pd.DataFrame({
                        "Salesman": ka_targets.index,
                        "KA Target": ka_targets.values,
                        "KA Sales": sales_by_sm.values
                    }).sort_values("KA Sales", ascending=False)

                    fig_salesman = go.Figure()
                    fig_salesman.add_trace(go.Bar(
                        x=target_actual_df["Salesman"],
                        y=target_actual_df["KA Sales"],
                        name="Sales",
                        text=target_actual_df["KA Sales"].apply(lambda x: f"{x:,.0f}"),
                        textposition="inside",
                        marker=dict(color=[
                            "#22C55E" if s >= t else "#EF4444"
                            for s, t in zip(target_actual_df["KA Sales"], target_actual_df["KA Target"])
                        ])
                    ))

                    fig_salesman.add_trace(go.Scatter(
                        x=target_actual_df["Salesman"],
                        y=target_actual_df["KA Target"],
                        name="Target",
                        mode="lines+markers",
                        line=dict(color="#1E3A8A", width=2)
                    ))

                    fig_salesman.update_layout(
                        xaxis_title="Salesman",
                        yaxis_title="Net Value (KD)",
                        template="plotly_white"
                    )
                    st.plotly_chart(fig_salesman, use_container_width=True)

                    st.markdown("---")
                    st.subheader("🏆 Top 10 Customers by Sales (KD)")

                    top10 = df_filtered.groupby("PY Name 1")["Net Value"].sum().sort_values(ascending=False).head(10)

                    fig_top10 = go.Figure(go.Bar(
                        x=top10.index,
                        y=top10.values,
                        text=[f"KD {v:,.0f}" for v in top10.values],
                        textposition="outside",
                        marker=dict(color="#2563EB")
                    ))

                    fig_top10.update_layout(
                        xaxis_title="Customer",
                        yaxis_title="Net Value (KD)",
                        template="plotly_white"
                    )
                    st.plotly_chart(fig_top10, use_container_width=True)
                    
                    

# --- YTD Comparison Page ---
elif choice == "Year to Date Comparison":
    if "ytd_df" in st.session_state and not st.session_state["ytd_df"].empty:
        ytd_df = st.session_state["ytd_df"].copy()

        if "Billing Date" not in ytd_df.columns:
            st.error("❌ 'Billing Date' column not found in YTD sheet.")
            st.stop()

        ytd_df["Billing Date"] = pd.to_datetime(ytd_df["Billing Date"], errors="coerce")
        ytd_df = ytd_df.dropna(subset=["Billing Date"])

        if "Net Value" not in ytd_df.columns:
            st.error("❌ 'Net Value' column not found in YTD sheet.")
            st.stop()

        st.title("📅 Year to Date Comparison")
        st.markdown(
            '<div class="tooltip">ℹ️<span class="tooltiptext">Compare sales across two periods by a selected dimension.</span></div>',
            unsafe_allow_html=True
        )

        # --- Select Dimension ---
        st.subheader("📊 Choose Dimension")

        # Map friendly labels to actual column names (added By Material)
        dim_options = {
            "By Customer": "PY Name 1",          # Customer Name column
            "By Salesman": "Driver Name EN",     # Salesman Name
            "By Branch": "SP Name1",             # Branch Name column
            "By Material": "Material Description"  # 👈 NEW: Material Description
        }

        # Filter only dimensions where the column actually exists
        dim_options_available = {
            label: col for label, col in dim_options.items() if col in ytd_df.columns
        }

        if not dim_options_available:
            st.error("❌ None of the expected dimension columns (Customer, Salesman, Branch, Material) are present.")
            st.stop()

        # User sees friendly names; program uses real column names
        dim_label = st.selectbox(
            "Choose dimension",
            list(dim_options_available.keys()),
            index=0
        )
        dimension = dim_options_available[dim_label]

        # --- Select Two Periods ---
        st.subheader("📆 Select Two Periods")
        col1, col2 = st.columns(2)

        min_date = ytd_df["Billing Date"].min().date()
        max_date = ytd_df["Billing Date"].max().date()

        with col1:
            st.write("Period 1")
            period1_range = st.date_input(
                "Select Date Range ( Click Twice Start & End Date)",
                value=(min_date, max_date),
                key="ytd_p1_range"
            )
        with col2:
            st.write("Period 2")
            period2_range = st.date_input(
                "Select Date Range ( Click Twice Start & End Date)",
                value=(min_date, max_date),
                key="ytd_p2_range"
            )

        if period1_range and period2_range and len(period1_range) == 2 and len(period2_range) == 2:
            period1_start, period1_end = period1_range
            period2_start, period2_end = period2_range

            df_p1 = ytd_df[
                (ytd_df["Billing Date"] >= pd.to_datetime(period1_start)) &
                (ytd_df["Billing Date"] <= pd.to_datetime(period1_end))
            ]
            df_p2 = ytd_df[
                (ytd_df["Billing Date"] >= pd.to_datetime(period2_start)) &
                (ytd_df["Billing Date"] <= pd.to_datetime(period2_end))
            ]

            if df_p1.empty or df_p2.empty:
                st.warning("⚠️ One of the selected periods has no data.")
            else:
                # --- YTD Comparison Table ---
                col_p1_name = f"{period1_start.strftime('%Y-%m-%d')} to {period1_end.strftime('%Y-%m-%d')} Sales"
                col_p2_name = f"{period2_start.strftime('%Y-%m-%d')} to {period2_end.strftime('%Y-%m-%d')} Sales"

                summary_p1 = (
                    df_p1.groupby(dimension)["Net Value"]
                    .sum()
                    .reset_index()
                    .rename(columns={"Net Value": col_p1_name})
                )
                summary_p2 = (
                    df_p2.groupby(dimension)["Net Value"]
                    .sum()
                    .reset_index()
                    .rename(columns={"Net Value": col_p2_name})
                )

                ytd_comparison = pd.merge(summary_p1, summary_p2, on=dimension, how="outer")

                # 🔒 Safety: ensure both period columns exist (avoid KeyError)
                if col_p1_name not in ytd_comparison.columns:
                    ytd_comparison[col_p1_name] = 0
                if col_p2_name not in ytd_comparison.columns:
                    ytd_comparison[col_p2_name] = 0

                ytd_comparison = ytd_comparison.fillna(0)

                # Difference = Period2 - Period1
                ytd_comparison["Difference"] = ytd_comparison[col_p2_name] - ytd_comparison[col_p1_name]

                # Rename dimension column to generic "Name"
                ytd_comparison.rename(columns={dimension: "Name"}, inplace=True)

                # Total row
                total_row = {
                    "Name": "Total",
                    col_p1_name: ytd_comparison[col_p1_name].sum(),
                    col_p2_name: ytd_comparison[col_p2_name].sum(),
                    "Difference": ytd_comparison["Difference"].sum()
                }
                ytd_comparison = pd.concat(
                    [ytd_comparison, pd.DataFrame([total_row])],
                    ignore_index=True
                )

                # Use friendly label in the heading
                st.subheader(f"📋 Comparison by {dim_label}")

                styled_ytd = (
                    ytd_comparison.style
                    .set_table_styles([
                        {
                            'selector': 'th',
                            'props': [
                                ('background', '#1E3A8A'),
                                ('color', 'white'),
                                ('font-weight', '800'),
                                ('height', '40px'),
                                ('line-height', '40px'),
                                ('border', '1px solid #E5E7EB')
                            ]
                        }
                    ])
                    .apply(
                        lambda x: [
                            'background-color: #BFDBFE; color: #1E3A8A; font-weight: 900'
                            if x.name == len(x) - 1 else ''  # last row = Total
                            for _ in x
                        ],
                        axis=1
                    )
                    .format(
                        {
                            col_p1_name: "{:,.0f}",
                            col_p2_name: "{:,.0f}",
                            'Difference': "{:,.0f}"
                        }
                    )
                )
                st.dataframe(styled_ytd, use_container_width=True, hide_index=False)

                st.download_button(
                    "⬇️ Download YTD Comparison (Excel)",
                    data=to_excel_bytes(ytd_comparison, sheet_name="YTD_Comparison", index=False),
                    file_name=f"YTD_Comparison_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

                # --- NEW: Sales by Month (YTD) ---
                st.subheader("📈 Sales by Month (YTD)")
                monthly = ytd_df.copy()
                monthly["YearMonth"] = monthly["Billing Date"].dt.to_period("M").astype(str)

                monthly_sales = (
                    monthly.groupby("YearMonth")["Net Value"]
                    .sum()
                    .reset_index()
                    .sort_values("YearMonth")
                )

                if not monthly_sales.empty:
                    fig_month = px.bar(
                        monthly_sales,
                        x="YearMonth",
                        y="Net Value",
                        title="Sales by Month (YTD)",
                        labels={"YearMonth": "Year-Month", "Net Value": "Net Value (KD)"}
                    )
                    fig_month.update_layout(height=400)
                    st.plotly_chart(fig_month, use_container_width=True)
                else:
                    st.info("No monthly data available in YTD sheet.")

                # --- Top 10 Customers: Last Year vs Current Year ---
                st.subheader("🏆 Top 10 Customers – Last Year vs Current Year")

                # Make sure Billing Date is datetime
                if not np.issubdtype(ytd_df["Billing Date"].dtype, np.datetime64):
                    ytd_df["Billing Date"] = pd.to_datetime(ytd_df["Billing Date"], errors="coerce")

                ytd_df["Year"] = ytd_df["Billing Date"].dt.year
                available_years = sorted(ytd_df["Year"].dropna().unique().tolist())
                if not available_years:
                    st.info("⚠️ No valid years found in YTD data.")
                    st.stop()
                default_year = max(available_years)
                selected_current_year = st.selectbox("Select Current Year:", available_years, index=available_years.index(default_year))
                current_year = int(selected_current_year)
                last_year = current_year - 1

                # Aggregate sales by Customer + Year
                cust_sales = (
                    ytd_df[ytd_df["Year"].isin([last_year, current_year])]
                    .groupby(["PY Name 1", "Year"])["Net Value"]
                    .sum()
                    .reset_index()
                )

                if cust_sales.empty:
                    st.info("⚠️ No customer sales found for last year or current year.")
                else:
                    cust_pivot = cust_sales.pivot(
                        index="PY Name 1",
                        columns="Year",
                        values="Net Value"
                    ).fillna(0)

                    # Use only year columns that actually exist (avoid KeyError)
                    year_cols = [y for y in [last_year, current_year] if y in cust_pivot.columns]

                    if not year_cols:
                        st.info("⚠️ No valid year columns found to build Top 10 customers chart.")
                    else:
                        cust_pivot["Total"] = cust_pivot[year_cols].sum(axis=1)
                        top10_cust = cust_pivot.sort_values("Total", ascending=False).head(10).reset_index()

                        # Melt for plotting
                        top10_melt = top10_cust.melt(
                            id_vars="PY Name 1",
                            value_vars=year_cols,
                            var_name="Year",
                            value_name="Sales",
                        )

                        # Merge back the wide data for status calculation
                        merge_cols = ["PY Name 1"] + year_cols
                        top10_melt = top10_melt.merge(
                            top10_cust[merge_cols],
                            on="PY Name 1",
                            how="left",
                        )

                        def classify_status(row):
                            # Both last_year & current_year exist → Achieved vs Not Achieved
                            if (last_year in year_cols) and (current_year in year_cols) and row["Year"] == current_year:
                                return "Achieved" if row.get(current_year, 0) >= row.get(last_year, 0) else "Not Achieved"
                            # Only current year data exists
                            if row["Year"] == current_year:
                                return "Current"
                            # Only previous year / others
                            return "Previous"

                        top10_melt["Status"] = top10_melt.apply(classify_status, axis=1)

                        color_map = {
                            "Achieved": "green",
                            "Not Achieved": "red",
                            "Previous": "gray",
                            "Current": "blue",
                        }

                        fig_top10 = px.bar(
                            top10_melt,
                            x="PY Name 1",
                            y="Sales",
                            color="Status",
                            color_discrete_map=color_map,
                            barmode="group",
                            text=top10_melt["Sales"].apply(lambda x: f"{x:,.0f}"),
                        )
                        fig_top10.update_traces(
                            textposition="inside",
                            insidetextanchor="middle",
                            textfont=dict(color="white", size=12),
                        )
                        fig_top10.update_layout(
                            title=f"Top 10 Customers: {last_year} vs {current_year}",
                            xaxis_title="Customer",
                            yaxis_title="Sales (KD)",
                            font=dict(family="Roboto", size=12),
                            plot_bgcolor="#F3F4F6",
                            paper_bgcolor="#F3F4F6",
                        )
                        st.plotly_chart(fig_top10, use_container_width=True)

                # --- Return by SP Name1 + Material Description (YTD) ---
                st.subheader("🔄 Return By Branch + Material Description (YTD)")
                required_cols = {"SP Name1", "Material Description", "Billing Type", "Net Value"}
                if required_cols.issubset(ytd_df.columns):
                    sp_mat_ytd = pd.pivot_table(
                        ytd_df,
                        index=["SP Name1", "Material Description"],
                        columns="Billing Type",
                        values="Net Value",
                        aggfunc="sum",
                        fill_value=0
                    )
                    billing_cols = ["ZFR", "YKF2", "YKRE", "YKS1", "YKS2", "ZCAN", "ZRE"]
                    for col in billing_cols:
                        if col not in sp_mat_ytd.columns:
                            sp_mat_ytd[col] = 0
                    sp_mat_ytd = sp_mat_ytd.rename(columns={"ZFR": "Presales", "YKF2": "HHT"})
                    sp_mat_ytd["Sales Total"] = sp_mat_ytd.sum(axis=1, numeric_only=True)
                    sp_mat_ytd["Return"] = sp_mat_ytd["YKRE"] + sp_mat_ytd["ZRE"]
                    sp_mat_ytd["Cancel Total"] = sp_mat_ytd[["YKS1", "YKS2", "ZCAN"]].sum(axis=1)
                    sp_mat_ytd["Return %"] = np.where(
                        sp_mat_ytd["Sales Total"] != 0,
                        (sp_mat_ytd["Return"] / sp_mat_ytd["Sales Total"] * 100).round(0),
                        0
                    )
                    ordered_cols = [
                        "Presales", "HHT", "Sales Total", "YKS1", "YKS2", "ZCAN",
                        "Cancel Total", "YKRE", "ZRE", "Return", "Return %"
                    ]
                    sp_mat_ytd = sp_mat_ytd.reindex(columns=ordered_cols, fill_value=0)
                    total_row = pd.DataFrame(sp_mat_ytd.sum(numeric_only=True)).T
                    total_row.index = [("Total", "")]
                    total_row["Return %"] = (
                        round((total_row["Return"] / total_row["Sales Total"] * 100), 0)
                        if total_row["Sales Total"].iloc[0] != 0 else 0
                    )
                    sp_mat_ytd = pd.concat([sp_mat_ytd, total_row])

                    # Conditional highlights for easy read
                    def highlight_sp_mat(row):
                        styles = []
                        for col in row.index:
                            if row.name == ("Total", ""):
                                styles.append('background-color: #BFDBFE; color: #1E3A8A; font-weight: 900')
                            elif col == "Return" and row[col] > 0:
                                styles.append('background-color: #FECACA; color: #991B1B; font-weight: 700')
                            elif col == "Cancel Total" and row[col] > 0:
                                styles.append('background-color: #FDE68A; color: #92400E; font-weight: 700')
                            elif col == "Sales Total" and row[col] > 0:
                                styles.append('background-color: #D1FAE5; color: #065F46; font-weight: 700')
                            else:
                                styles.append('')
                        return styles

                    styled_sp_mat = (
                        sp_mat_ytd.style
                        .set_table_styles([
                            {
                                'selector': 'th',
                                'props': [
                                    ('background', '#1E3A8A'),
                                    ('color', 'white'),
                                    ('font-weight', '800'),
                                    ('height', '40px'),
                                    ('line-height', '40px'),
                                    ('border', '1px solid #E5E7EB')
                                ]
                            }
                        ])
                        .apply(highlight_sp_mat, axis=1)
                        .format({
                            "Presales": "{:,.0f}", "HHT": "{:,.0f}", "Sales Total": "{:,.0f}",
                            "YKS1": "{:,.0f}", "YKS2": "{:,.0f}", "ZCAN": "{:,.0f}", "Cancel Total": "{:,.0f}",
                            "YKRE": "{:,.0f}", "ZRE": "{:,.0f}", "Return": "{:,.0f}", "Return %": "{:.0f}%"
                        })
                    )
                    st.dataframe(styled_sp_mat, use_container_width=True, hide_index=True)
                    st.download_button(
                        "⬇️ Download Return by SP+Material (YTD)",
                        data=to_excel_bytes(sp_mat_ytd.reset_index(), sheet_name="Return_by_SP_Material_YTD", index=False),
                        file_name=f"Return_by_SP_Material_YTD_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    )
                else:
                    st.info("Required columns for SP+Material YTD table are missing.")
    else:
        st.warning("⚠️ Please ensure the 'YTD' sheet is present in your uploaded file.")


# --- Custom Analysis Page ---
elif choice == texts[lang]["custom_analysis"]:
    st.title(texts[lang]["custom_title"])
    if "data_loaded" not in st.session_state:
        st.warning(texts[lang]["no_data_warning"])
    else:
        # ✅ Ensure Extra sheet is loaded into session state
        if "Extra_sheet_df" not in st.session_state:
            try:
                extra_df = pd.read_excel(uploaded, sheet_name="Extra sheet")
            except Exception:
                extra_df = pd.DataFrame()
            st.session_state["Extra_sheet_df"] = extra_df

        # Available sheet options
        sheet_options = {
            "Sales Data": st.session_state.get("sales_df", pd.DataFrame()),
            "YTD": st.session_state.get("ytd_df", pd.DataFrame()),
            "Target": st.session_state.get("target_df", pd.DataFrame()),
            "Sales Channels": st.session_state.get("channels_df", pd.DataFrame()),
            "Extra sheet": st.session_state.get("Extra_sheet_df", pd.DataFrame())
        }

        selected_sheet_name = st.selectbox(texts[lang]["custom_select_sheet"], list(sheet_options.keys()))
        df = sheet_options[selected_sheet_name]

        if df.empty:
            st.warning(texts[lang]["custom_sheet_empty"].format(selected_sheet_name))
        else:
            st.subheader(texts[lang]["custom_explore"])

            available_cols = list(df.columns)
            group_cols = st.multiselect(texts[lang]["custom_group_cols"], available_cols)
            value_col = st.selectbox(texts[lang]["custom_value_col"], available_cols)

            if "Billing Date" in df.columns:
                st.subheader(texts[lang]["custom_periods_sub"])
                col1, col2 = st.columns(2)
                with col1:
                    st.write(texts[lang]["custom_period1"])
                    period1_range = st.date_input(
                        texts[lang]["custom_select_p1"],
                        [df["Billing Date"].min(), df["Billing Date"].max()],
                        key="ca_p1_range"
                    )
                with col2:
                    st.write(texts[lang]["custom_period2"])
                    period2_range = st.date_input(
                        texts[lang]["custom_select_p2"],
                        [df["Billing Date"].min(), df["Billing Date"].max()],
                        key="ca_p2_range"
                    )
            else:
                period1_range = period2_range = None
                st.info("⚠️ No 'Billing Date' column found. Period comparison disabled.")

            if group_cols and value_col and period1_range and period2_range and len(period1_range) == 2 and len(period2_range) == 2:
                # --- Period 1 ---
                p1_start, p1_end = pd.to_datetime(period1_range[0]), pd.to_datetime(period1_range[1])
                df_p1 = df[(df["Billing Date"] >= p1_start) & (df["Billing Date"] <= p1_end)]
                summary_p1 = df_p1.groupby(group_cols)[value_col].sum().reset_index()
                summary_p1.rename(columns={value_col: "Period 1"}, inplace=True)

                # --- Period 2 ---
                p2_start, p2_end = pd.to_datetime(period2_range[0]), pd.to_datetime(period2_range[1])
                df_p2 = df[(df["Billing Date"] >= p2_start) & (df["Billing Date"] <= p2_end)]
                summary_p2 = df_p2.groupby(group_cols)[value_col].sum().reset_index()
                summary_p2.rename(columns={value_col: "Period 2"}, inplace=True)

                # --- Merge & Compare ---
                comparison_df = pd.merge(summary_p1, summary_p2, on=group_cols, how="outer").fillna(0)
                comparison_df["Difference"] = comparison_df["Period 2"] - comparison_df["Period 1"]

                st.subheader(texts[lang]["custom_comparison_sub"].format(value_col, ", ".join(group_cols)))
                styled_custom = (
                    comparison_df.style
                    .set_table_styles([
                        {'selector': 'th', 'props': [('background', '#1E3A8A'), ('color', 'white'),
                                                    ('font-weight', '800'), ('height', '40px'),
                                                    ('line-height', '40px'), ('border', '1px solid #E5E7EB')]}
                    ])
                    .format({
                        "Period 1": "{:,.0f}",
                        "Period 2": "{:,.0f}",
                        "Difference": "{:,.0f}"
                    })
                )
                st.dataframe(styled_custom, use_container_width=True, hide_index=True)

                # --- Plotly Chart Fix ---
                df_plot = comparison_df.sort_values(by="Period 2", ascending=False).copy()

                if len(group_cols) == 1:
                    df_plot["Group"] = df_plot[group_cols[0]].astype(str)
                elif len(group_cols) > 1:
                    df_plot["Group"] = df_plot[group_cols].astype(str).agg(" | ".join, axis=1)
                else:
                    df_plot["Group"] = "All Data"

                df_plot_melted = df_plot.melt(
                    id_vars=["Group"],
                    value_vars=["Period 1", "Period 2"],
                    var_name="Period",
                    value_name="Value"
                )

                fig = px.bar(
                    df_plot_melted,
                    x="Group",
                    y="Value",
                    color="Period",
                    barmode="group",
                    title=f"Comparison of {value_col} by {', '.join(group_cols) if group_cols else 'All'}",
                    color_discrete_sequence=px.colors.qualitative.Set2
                )
                st.plotly_chart(fig, use_container_width=True)

                # --- Download ---
                if st.download_button(
                    texts[lang]["custom_download"],
                    data=to_excel_bytes(comparison_df, sheet_name="Custom_Comparison", index=False),
                    file_name=f"Custom_Comparison_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                ):
                    st.session_state["audit_log"].append({
                        "user": username,
                        "action": "download",
                        "details": "Custom Comparison Excel",
                        "timestamp": datetime.now()
                    })
            else:
                st.info(texts[lang]["custom_select_prompt"])

# --- SP/PY Target Allocation Page ---
elif choice == "SP/PY Target Allocation":
    st.title("🎯 By Customer & By Branch Target Allocation")
    if "data_loaded" not in st.session_state:
        st.warning("⚠️ Please upload the Excel file from the sidebar first.")
        st.stop()

    sales_df = st.session_state["sales_df"]
    ytd_df = st.session_state["ytd_df"]
    target_df = st.session_state.get("target_df", pd.DataFrame())

    st.subheader("Configuration")
    st.markdown('<div class="tooltip">ℹ️<span class="tooltiptext">Allocate targets by branch or customer based on historical sales.</span></div>', unsafe_allow_html=True)
        # Target allocation grouping choice
    allocation_type = st.radio(
        "Select Target Allocation Type",
        ["By Branch", "By Customer"]
    )

    # Map choice to actual column
    group_col = "SP Name1" if allocation_type == "By Branch" else "PY Name 1"

    target_option = st.radio("Select Target Input Option", ["Manual", "Auto (from 'Target' sheet)"])

    total_target = 0
    if target_option == "Manual":
        total_target = st.number_input("Enter the Total Target to be Allocated for this Month (KD)", min_value=0, value=1000000, step=1000)
    else:
        if target_df.empty or "KA Target" not in target_df.columns:
            st.error("❌ 'Target' sheet or 'KA Target' column not found. Please upload a file with this sheet for 'Auto' mode.")
            st.stop()
        total_target = target_df["KA Target"].sum()
        st.info(f"Using Total Target from 'Target' sheet: KD {total_target:,.0f}")

    if total_target <= 0:
        st.warning("Please ensure the total target is greater than 0.")
        st.stop()

    st.subheader("Historical Data Period")
    today = pd.Timestamp.today().normalize()
    data_period_option = st.radio("Select Historical Data Period", ["Last 6 Months", "Manual Days"], index=1)

    if data_period_option == "Last 6 Months":
        lookback_period = pd.DateOffset(months=6)
        days_label = "6 Months"; months_count = 6
        end_date_selected = today; start_date_selected = today - lookback_period
    else:
        start_date_manual = today - pd.DateOffset(days=180)
        selected_dates = st.date_input("Select date range", value=(start_date_manual, today))
        if len(selected_dates) == 2:
            start_date_selected, end_date_selected = selected_dates
            lookback_period = end_date_selected - start_date_selected
            days_label = f"From {start_date_selected.strftime('%Y-%m-%d')} to {end_date_selected.strftime('%Y-%m-%d')}"
            months_count = lookback_period.days / 30
        else:
            st.warning("Please select both a start and an end date.")
            st.stop()

    historical_df = ytd_df[(ytd_df["Billing Date"] >= pd.Timestamp(start_date_selected)) & (ytd_df["Billing Date"] <= pd.Timestamp(end_date_selected))].copy()
    if historical_df.empty:
        st.warning(f"⚠️ No sales data available in 'YTD' for {days_label}.")
        st.stop()

    historical_sales = historical_df.groupby(group_col)["Net Value"].sum()
    total_historical_sales_value = historical_sales.sum()
    current_month_sales_df = sales_df[(sales_df["Billing Date"].dt.month == today.month) & (sales_df["Billing Date"].dt.year == today.year)].copy()
    current_month_sales = current_month_sales_df.groupby(group_col)["Net Value"].sum()
    total_current_month_sales = current_month_sales.sum()

    target_balance = total_target - total_current_month_sales

    if total_target > 0:
        average_historical_sales = total_historical_sales_value / months_count
        st.subheader("🎯 Target Analysis")
        col1, col2, col3 = st.columns(3)
        col4, col5 = st.columns(2)
        with col1: st.metric("Historical Sales Total", f"KD {total_historical_sales_value:,.0f}")
        with col2: st.metric("Allocated Target Total", f"KD {total_target:,.0f}")
        with col3:
            if average_historical_sales > 0:
                percentage_increase_needed = ((total_target - average_historical_sales) / average_historical_sales) * 100
                delta_value = total_target - average_historical_sales
                st.metric("Increase Needed vs Avg Sales", f"{percentage_increase_needed:.0f}%", delta=f"KD {delta_value:,.0f}")
            else:
                st.metric("Increase Needed vs Avg Sales", "N/A", delta="Historical = 0")
        st.markdown("---")
        with col4: st.metric("Current Month Sales", f"KD {total_current_month_sales:,.0f}")
        with col5: st.metric("Target Balance", f"KD {target_balance:,.0f}")

    allocation_table = pd.DataFrame(index=historical_sales.index.union(current_month_sales.index).unique())
    allocation_table.index.name = "Name"
    allocation_table[f"Last {days_label} Total Sales"] = historical_sales.reindex(allocation_table.index, fill_value=0)
    allocation_table[f"Last {days_label} Average Sales"] = allocation_table[f"Last {days_label} Total Sales"] / months_count
    if total_historical_sales_value > 0:
        allocation_table["This Month Auto-Allocated Target"] = allocation_table[f"Last {days_label} Total Sales"] / total_historical_sales_value * total_target
    else:
        allocation_table["This Month Auto-Allocated Target"] = 0
    allocation_table["Current Month Sales"] = current_month_sales.reindex(allocation_table.index, fill_value=0)
    allocation_table["Target Balance"] = allocation_table["This Month Auto-Allocated Target"] - allocation_table["Current Month Sales"]

    total_row = allocation_table.sum().to_frame().T
    total_row.index = ["Total"]
    total_row["Target Balance"] = total_target - total_current_month_sales
    allocation_table_with_total = pd.concat([allocation_table, total_row])

    def color_target_balance(val):
        if isinstance(val, (int, float)):
            color = 'red' if val > 0 else 'green'
            return f'color: {color}'
        return ''

    st.subheader(f"📊 Auto-Allocated Targets Based on {days_label}")

    # Show the first column (Branch/Customer) with a proper header
    name_col = "Branch" if allocation_type == "By Branch" else "Customer"

    allocation_display = allocation_table_with_total.reset_index()
    first_col = allocation_display.columns[0]
    allocation_display = allocation_display.rename(columns={first_col: name_col})

    numeric_cols = [c for c in allocation_display.columns if c != name_col]
    for c in numeric_cols:
        allocation_display[c] = (
            pd.to_numeric(allocation_display[c], errors="coerce")
            .fillna(0)
            .round(0)
            .astype(int)
        )

    styled_allocation = (
        allocation_display.style
        .set_table_styles([
            {'selector': 'th', 'props': [
                ('background', '#1E3A8A'),
                ('color', 'white'),
                ('font-weight', '800'),
                ('height', '40px'),
                ('line-height', '40px'),
                ('border', '1px solid #E5E7EB')
            ]}
        ])
        .apply(
            lambda r: ['background-color: #BFDBFE; color: #1E3A8A; font-weight: 900' if str(r.get(name_col)) == 'Total' else '' for _ in r],
            axis=1
        )
        .format({c: '{:,.0f}' for c in numeric_cols})
    )
    st.dataframe(styled_allocation, use_container_width=True, hide_index=True)

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    excel_data = to_excel_bytes(allocation_table, sheet_name="Allocated_Targets")
    if st.download_button(
        "💾 Download Target Allocation Table",
        data=excel_data,
        file_name=f"target_allocation_{allocation_type.replace(' ', '_')}_{timestamp}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        st.session_state["audit_log"].append({
            "user": username,
            "action": "download",
            "details": f"target_allocation_{allocation_type.replace(' ', '_')}_{timestamp}.xlsx",
            "timestamp": datetime.now()
        })

# --- AI Insights Page (GM Executive View - Reset) ---
    def gm_tag(return_pct, cancel_pct):
        if return_pct >= 4:
            return "🔴 High Return"
        if cancel_pct >= 4:
            return "🟠 High Cancel"
        if return_pct >= 2 or cancel_pct >= 2:
            return "🟡 Watch"
        return "🟢 Normal"
    
# --- AI Insights Page (GM Executive View - No Duplicate Key Metrics) ---
elif choice == "AI Insights":
    st.title("🧠 GM Insights – Executive View")

    if "data_loaded" not in st.session_state:
        st.warning("⚠️ Please upload the Excel file first.")
    else:
        import calendar

        # ------------------------------------------------
        # 1) Base Data
        # ------------------------------------------------
        sales_df = st.session_state["sales_df"].copy()
        target_df = st.session_state.get("target_df", pd.DataFrame()).copy()
        channels_df = st.session_state.get("channels_df", pd.DataFrame()).copy()

        # Ensure date
        if "Billing Date" in sales_df.columns:
            sales_df["Billing Date"] = pd.to_datetime(sales_df["Billing Date"], errors="coerce")

        def fmt_kd(x):
            try:
                return f"KD {float(x):,.0f}"
            except Exception:
                return "KD 0"

        # ✅ Local helper (fix NameError)
        def gm_tag(return_pct, cancel_pct):
            try:
                return_pct = float(return_pct)
                cancel_pct = float(cancel_pct)
            except Exception:
                return "🟢 Normal"

            if return_pct >= 4:
                return "🔴 High Return"
            if cancel_pct >= 4:
                return "🟠 High Cancel"
            if return_pct >= 2 or cancel_pct >= 2:
                return "🟡 Watch"
            return "🟢 Normal"

        # ------------------------------------------------
        # 2) Filters
        # ------------------------------------------------
        st.subheader("🎛 GM Scope")

        min_date = pd.to_datetime(sales_df["Billing Date"].min())
        max_date = pd.to_datetime(sales_df["Billing Date"].max())

        f1, f2, f3 = st.columns([2, 1, 2])
        with f1:
            date_range = st.date_input(
                "Select GM period",
                value=(min_date.date(), max_date.date())
            )
        with f2:
            top_n = st.slider("Top N", 3, 15, 5, 1)
        with f3:
            sm_list = []
            if "Driver Name EN" in sales_df.columns:
                sm_list = sorted([x for x in sales_df["Driver Name EN"].dropna().unique()])
            selected_sm = st.multiselect("Salesmen (optional)", sm_list, default=[])

        if isinstance(date_range, (list, tuple)) and len(date_range) == 2:
            start = pd.to_datetime(date_range[0])
            end = pd.to_datetime(date_range[1])
        else:
            start, end = min_date, max_date

        df = sales_df[(sales_df["Billing Date"] >= start) & (sales_df["Billing Date"] <= end)].copy()

        if selected_sm and "Driver Name EN" in df.columns:
            df = df[df["Driver Name EN"].isin(selected_sm)].copy()

        if df.empty:
            st.info("No data in selected period.")
            st.stop()

        # ------------------------------------------------
        # 3) Billing Types & Core GM numbers
        # ------------------------------------------------
        SALES_BT  = {"ZFR", "YKF2"}              # Sales only (Presales + HHT)
        RETURN_BT = {"YKRE", "ZRE"}              # Returns
        CANCEL_BT = {"YKS1", "YKS2", "ZCAN"}     # Cancels

        if "Billing Type" not in df.columns:
            df["Billing Type"] = ""
        df["Billing Type"] = df["Billing Type"].astype(str).str.upper().str.strip()

        sales_val = float(df[df["Billing Type"].isin(SALES_BT)]["Net Value"].sum())
        returns_raw = float(df[df["Billing Type"].isin(RETURN_BT)]["Net Value"].sum())
        cancel_raw  = float(df[df["Billing Type"].isin(CANCEL_BT)]["Net Value"].sum())

        # Make effect negative for Net Sales (safe even if data has positive values)
        returns_effect = returns_raw if returns_raw < 0 else -abs(returns_raw)
        cancel_effect  = cancel_raw  if cancel_raw  < 0 else -abs(cancel_raw)

        net_sales = sales_val + returns_effect + cancel_effect
        returns_val = abs(returns_raw)
        cancel_val = abs(cancel_raw)

        return_pct = (returns_val / sales_val * 100) if sales_val else 0
        cancel_pct = (cancel_val / sales_val * 100) if sales_val else 0

        # ------------------------------------------------
        # 4) Retail vs E-com mix (Sales only)
        # ------------------------------------------------
        retail_sales = 0.0
        ecom_sales = 0.0

        if (not channels_df.empty) and {"PY Name 1", "Channels"}.issubset(channels_df.columns) and "PY Name 1" in df.columns:
            tmp = df[df["Billing Type"].isin(SALES_BT)].copy()

            tmp["_py_norm"] = tmp["PY Name 1"].astype(str).str.strip().str.lower()
            ch = channels_df.copy()
            ch["_py_norm"] = ch["PY Name 1"].astype(str).str.strip().str.lower()

            tmp = tmp.merge(ch[["_py_norm", "Channels"]], on="_py_norm", how="left")
            tmp["Channels"] = tmp["Channels"].astype(str).str.lower().str.strip()
            tmp.loc[tmp["Channels"].isin(["", "nan", "none"]), "Channels"] = "retail"

            e_mask = tmp["Channels"].str.contains("e-com|ecom|ecommerce|online|talabat", regex=True, na=False)
            ecom_sales = float(tmp[e_mask]["Net Value"].sum())
            retail_sales = float(tmp[~e_mask]["Net Value"].sum())
        else:
            # fallback
            retail_sales = float(df[df["Billing Type"].isin(SALES_BT)]["Net Value"].sum())
            ecom_sales = 0.0

        mix_total = retail_sales + ecom_sales
        retail_mix = (retail_sales / mix_total * 100) if mix_total else 0
        ecom_mix   = (ecom_sales / mix_total * 100) if mix_total else 0

        # ------------------------------------------------
        # 5) GM Header (Status line)
        # ------------------------------------------------
        active_sm = df["Driver Name EN"].dropna().nunique() if "Driver Name EN" in df.columns else 0
        status = "🟢 Stable" if (return_pct < 3 and cancel_pct < 3) else "🟠 Needs Attention" if (return_pct < 5 and cancel_pct < 5) else "🔴 At Risk"

        st.markdown(
            f"**Period:** {start.date()} → {end.date()}  |  "
            f"**Active Salesmen:** {active_sm}  |  "
            f"**Status:** {status}"
        )

        # ------------------------------------------------
        # 6) GM KPIs (No duplicate Key Metrics)
        # ------------------------------------------------
        st.markdown("---")
        st.subheader("📊 GM Executive KPIs")

        k1, k2, k3, k4 = st.columns(4)
        k1.metric("Net Sales", fmt_kd(net_sales))
        k2.metric("Return Rate %", f"{return_pct:.1f}%")
        k3.metric("Cancel Rate %", f"{cancel_pct:.1f}%")
        k4.metric("Retail / E-Com Mix", f"{retail_mix:.0f}% / {ecom_mix:.0f}%")

        # ------------------------------------------------
        # 7) LY same dates + Forecast month end + YTD vs LY + Forecast year end
        # ------------------------------------------------
        def _sales_sum(df_):
            if df_.empty:
                return 0.0
            if "Billing Type" in df_.columns:
                _tmp = df_.copy()
                _tmp["Billing Type"] = _tmp["Billing Type"].astype(str).str.upper().str.strip()
                return float(_tmp[_tmp["Billing Type"].isin(SALES_BT)]["Net Value"].sum())
            return float(df_["Net Value"].sum())

        cur_start = pd.to_datetime(start)
        cur_end   = pd.to_datetime(end)

        # Same date-to-date last year
        ly_start = cur_start - pd.DateOffset(years=1)
        ly_end   = cur_end   - pd.DateOffset(years=1)

        cur_period_sales = _sales_sum(sales_df[(sales_df["Billing Date"] >= cur_start) & (sales_df["Billing Date"] <= cur_end)])
        ly_period_sales  = _sales_sum(sales_df[(sales_df["Billing Date"] >= ly_start) & (sales_df["Billing Date"] <= ly_end)])

        # Month forecast
        month_start = cur_end.replace(day=1)
        days_in_month = calendar.monthrange(cur_end.year, cur_end.month)[1]
        month_end = cur_end.replace(day=days_in_month)

        mtd_df = sales_df[(sales_df["Billing Date"] >= month_start) & (sales_df["Billing Date"] <= cur_end)].copy()
        mtd_sales = _sales_sum(mtd_df)

        mtd_days_with_data = int(mtd_df["Billing Date"].dt.date.nunique()) if not mtd_df.empty else 0
        if mtd_days_with_data > 0:
            forecast_month_end = (mtd_sales / mtd_days_with_data) * days_in_month
        else:
            forecast_month_end = 0.0

        # YTD vs LY
        ytd_start = cur_end.replace(month=1, day=1)
        ly_ytd_start = (cur_end - pd.DateOffset(years=1)).replace(month=1, day=1)
        ly_ytd_end = cur_end - pd.DateOffset(years=1)

        ytd_sales = _sales_sum(sales_df[(sales_df["Billing Date"] >= ytd_start) & (sales_df["Billing Date"] <= cur_end)])
        ly_ytd_sales = _sales_sum(sales_df[(sales_df["Billing Date"] >= ly_ytd_start) & (sales_df["Billing Date"] <= ly_ytd_end)])

        # Year forecast
        is_leap = (cur_end.year % 4 == 0 and cur_end.year % 100 != 0) or (cur_end.year % 400 == 0)
        days_in_year = 366 if is_leap else 365

        ytd_df_range = sales_df[(sales_df["Billing Date"] >= ytd_start) & (sales_df["Billing Date"] <= cur_end)].copy()
        ytd_days_with_data = int(ytd_df_range["Billing Date"].dt.date.nunique()) if not ytd_df_range.empty else 0

        if ytd_days_with_data > 0:
            forecast_year_end = (ytd_sales / ytd_days_with_data) * days_in_year
        else:
            forecast_year_end = 0.0

        # ------------------------------------------------
        # 7) GM Comparison & Forecast (FIXED using YTD sheet)
        # ------------------------------------------------
        st.markdown("---")
        st.subheader("📌 GM Comparison & Forecast")

        # ✅ Use YTD sheet for history if available, else fallback to sales_df
        ytd_df = st.session_state.get("ytd_df", pd.DataFrame()).copy()
        src = ytd_df if (not ytd_df.empty and "Billing Date" in ytd_df.columns and "Net Value" in ytd_df.columns) else sales_df

        # Ensure date + numeric
        src["Billing Date"] = pd.to_datetime(src["Billing Date"], errors="coerce")
        src["Net Value"] = pd.to_numeric(src["Net Value"], errors="coerce").fillna(0)

        # IMPORTANT:
        # - YTD sheet usually already contains "Net Value" final numbers (no billing type split)
        # - So for comparisons/forecasts, use Net Value sum directly.
        def sum_value(df_):
            if df_.empty:
                return 0.0
            return float(df_["Net Value"].sum())

        cur_start = pd.to_datetime(start)
        cur_end   = pd.to_datetime(end)

        # Same date-to-date last year (from YTD sheet)
        ly_start = cur_start - pd.DateOffset(years=1)
        ly_end   = cur_end   - pd.DateOffset(years=1)

        cur_period_val = sum_value(src[(src["Billing Date"] >= cur_start) & (src["Billing Date"] <= cur_end)])
        ly_period_val  = sum_value(src[(src["Billing Date"] >= ly_start) & (src["Billing Date"] <= ly_end)])

        # ---- Month forecast (if month already completed => forecast = actual month) ----
        month_start = cur_end.replace(day=1)
        days_in_month = calendar.monthrange(cur_end.year, cur_end.month)[1]
        month_end = cur_end.replace(day=days_in_month)

        month_actual = sum_value(src[(src["Billing Date"] >= month_start) & (src["Billing Date"] <= month_end)])
        mtd_actual   = sum_value(src[(src["Billing Date"] >= month_start) & (src["Billing Date"] <= cur_end)])

        elapsed_days = (cur_end.date() - month_start.date()).days + 1  # calendar elapsed days

        if cur_end.date() >= month_end.date():
            forecast_month_end = month_actual   # ✅ month finished
        else:
            forecast_month_end = (mtd_actual / elapsed_days) * days_in_month if elapsed_days > 0 else 0.0

        # ---- YTD vs LY (Jan 1 -> current end) ----
        ytd_start = cur_end.replace(month=1, day=1)
        ly_ytd_start = (cur_end - pd.DateOffset(years=1)).replace(month=1, day=1)
        ly_ytd_end   = cur_end - pd.DateOffset(years=1)

        ytd_val    = sum_value(src[(src["Billing Date"] >= ytd_start) & (src["Billing Date"] <= cur_end)])
        ly_ytd_val = sum_value(src[(src["Billing Date"] >= ly_ytd_start) & (src["Billing Date"] <= ly_ytd_end)])

        # ---- Year forecast (if year finished => forecast = actual YTD) ----
        is_leap = (cur_end.year % 4 == 0 and cur_end.year % 100 != 0) or (cur_end.year % 400 == 0)
        days_in_year = 366 if is_leap else 365

        day_of_year = (cur_end.date() - ytd_start.date()).days + 1  # calendar YTD days

        if cur_end.month == 12 and cur_end.day == 31:
            forecast_year_end = ytd_val  # ✅ year finished
        else:
            forecast_year_end = (ytd_val / day_of_year) * days_in_year if day_of_year > 0 else 0.0

        # ---- Cards ----
        g1, g2, g3, g4 = st.columns(4)

        g1.metric(
            "LY Same Dates",
            fmt_kd(ly_period_val),
            delta=f"{((cur_period_val-ly_period_val)/ly_period_val*100):.1f}% vs LY" if ly_period_val > 0 else None
        )

        g2.metric(
            "Forecast Month End",
            fmt_kd(forecast_month_end) if forecast_month_end > 0 else "N/A"
        )

        g3.metric(
            "YTD vs LY",
            fmt_kd(ytd_val),
            delta=f"{((ytd_val-ly_ytd_val)/ly_ytd_val*100):.1f}% vs LY" if ly_ytd_val > 0 else None
        )

        g4.metric(
            "Forecast Year End",
            fmt_kd(forecast_year_end) if forecast_year_end > 0 else "N/A"
        )

        st.caption(
            f"Source used: {'YTD sheet' if src is ytd_df else 'Sales sheet'} | "
            f"LY Same Dates: {ly_start.date()} → {ly_end.date()} | "
            f"YTD: Jan 1 → {cur_end.date()} (vs LY Jan 1 → {ly_ytd_end.date()})"
        )
        
        # ==========================================================
        # ✅ FULL AI INTELLIGENCE MODE (NO REPEAT) – GM Structured Notes
        # Paste this block AFTER your existing:
        #  - GM Snapshot KPIs
        #  - GM Comparison & Forecast cards
        # So it will NOT repeat target/forecast/mix totals again.
        # ==========================================================

        st.markdown("---")
        st.subheader("🧠 Full AI Intelligence Mode (GM Structured Notes)")

        # ---------------- Helpers ----------------
        def _safe_pct(a, b):
            return (a / b * 100) if b else 0.0

        def _sum_sales_only(df_):
            """Sales-only = ZFR + YKF2 (pre-sales + HHT)"""
            if df_ is None or df_.empty:
                return 0.0
            d = df_.copy()
            if "Billing Type" in d.columns:
                d["Billing Type"] = d["Billing Type"].astype(str).str.upper().str.strip()
                return float(d[d["Billing Type"].isin({"ZFR", "YKF2"})]["Net Value"].sum())
            return float(d["Net Value"].sum())

        def _sum_returns(df_):
            if df_ is None or df_.empty or "Billing Type" not in df_.columns:
                return 0.0
            d = df_.copy()
            d["Billing Type"] = d["Billing Type"].astype(str).str.upper().str.strip()
            return float(d[d["Billing Type"].isin({"YKRE", "ZRE"})]["Net Value"].sum())

        def _sum_cancels(df_):
            if df_ is None or df_.empty or "Billing Type" not in df_.columns:
                return 0.0
            d = df_.copy()
            d["Billing Type"] = d["Billing Type"].astype(str).str.upper().str.strip()
            return float(d[d["Billing Type"].isin({"YKS1", "YKS2", "ZCAN"})]["Net Value"].sum())

        def calc_mix(df_src):
            """
            Returns (retail_value, ecom_value) using channels_df mapping.
            Fallback: (total, 0) if channels not available.
            """
            if df_src is None or df_src.empty:
                return (0.0, 0.0)

            if channels_df is None or channels_df.empty or not {"PY Name 1", "Channels"}.issubset(channels_df.columns):
                total = float(df_src["Net Value"].sum())
                return (total, 0.0)

            if "PY Name 1" not in df_src.columns:
                total = float(df_src["Net Value"].sum())
                return (total, 0.0)

            tmp = df_src.copy()
            tmp["_py_norm"] = tmp["PY Name 1"].astype(str).str.strip().str.lower()

            ch = channels_df.copy()
            ch["_py_norm"] = ch["PY Name 1"].astype(str).str.strip().str.lower()

            tmp = tmp.merge(ch[["_py_norm", "Channels"]], on="_py_norm", how="left")
            tmp["Channels"] = tmp["Channels"].astype(str).str.lower().str.strip()
            tmp.loc[tmp["Channels"].isin(["", "nan", "none"]), "Channels"] = "retail"

            e_mask = tmp["Channels"].str.contains("e-com|ecom|ecommerce|online|talabat", regex=True, na=False)
            e = float(tmp[e_mask]["Net Value"].sum())
            r = float(tmp[~e_mask]["Net Value"].sum())
            return (r, e)

        def fmt_kd(x):
            try:
                return f"KD {float(x):,.0f}"
            except Exception:
                return "KD 0"

        # ---------------- Base frames ----------------
        # df_ai (your filtered period) should exist. If not, fallback:
        try:
            df_cur = df_ai.copy()
        except Exception:
            df_cur = df.copy()

        df_cur = df_cur.copy()
        df_cur["Billing Date"] = pd.to_datetime(df_cur["Billing Date"], errors="coerce")
        df_cur["Net Value"] = pd.to_numeric(df_cur["Net Value"], errors="coerce").fillna(0)

        cur_start = pd.to_datetime(start)
        cur_end = pd.to_datetime(end)

        # For LY comparisons: use YTD sheet if it has dates+values
        ytd_df = st.session_state.get("ytd_df", pd.DataFrame()).copy()
        use_ytd = (not ytd_df.empty and {"Billing Date", "Net Value"}.issubset(ytd_df.columns))
        hist_src = ytd_df if use_ytd else sales_df

        hist_src = hist_src.copy()
        hist_src["Billing Date"] = pd.to_datetime(hist_src["Billing Date"], errors="coerce")
        hist_src["Net Value"] = pd.to_numeric(hist_src["Net Value"], errors="coerce").fillna(0)

        ly_start = cur_start - pd.DateOffset(years=1)
        ly_end = cur_end - pd.DateOffset(years=1)

        # ---------------- A) Executive Summary (NO repeat) ----------------
        st.markdown("### 📝 Executive Summary (GM Notes)")

        # YoY same dates (value-based) – from hist source
        cur_same_val = float(hist_src[(hist_src["Billing Date"] >= cur_start) & (hist_src["Billing Date"] <= cur_end)]["Net Value"].sum())
        ly_same_val  = float(hist_src[(hist_src["Billing Date"] >= ly_start) & (hist_src["Billing Date"] <= ly_end)]["Net Value"].sum())
        yoy_pct = (_safe_pct(cur_same_val - ly_same_val, ly_same_val) if ly_same_val > 0 else None)

        # Momentum (7d vs prev7d) – from current period daily totals
        ts_daily = df_cur.groupby(df_cur["Billing Date"].dt.date)["Net Value"].sum().sort_index()
        if len(ts_daily) >= 7:
            last7 = ts_daily.tail(7).mean()
            prev7 = ts_daily.tail(14).head(7).mean() if len(ts_daily) >= 14 else None
        else:
            last7, prev7 = None, None

        mom = None
        if prev7 is not None and prev7 != 0 and last7 is not None:
            mom = (last7 - prev7) / prev7 * 100

        notes = []
        notes.append(f"Period: {cur_start.date()} → {cur_end.date()} | LY reference: {ly_start.date()} → {ly_end.date()} | Source: {'YTD sheet' if use_ytd else 'Sales sheet'}")

        if yoy_pct is not None:
            notes.append(f"Same-date YoY: **{yoy_pct:+.1f}%** (Current {fmt_kd(cur_same_val)} vs LY {fmt_kd(ly_same_val)}).")
        else:
            notes.append("Same-date YoY: **N/A** (LY data not found for same dates).")

        # Only risk flags (no repeating net/mix/forecast cards)
        if "return_pct" in globals():
            if return_pct >= 3:
                notes.append(f"🔴 Returns risk: **{return_pct:.1f}%** is high.")
            elif return_pct >= 2:
                notes.append(f"🟠 Returns watch: **{return_pct:.1f}%**.")
        if "cancel_pct" in globals():
            if cancel_pct >= 5:
                notes.append(f"🔴 Cancels risk: **{cancel_pct:.1f}%** is very high.")
            elif cancel_pct >= 3:
                notes.append(f"🟠 Cancels watch: **{cancel_pct:.1f}%**.")

        if mom is not None:
            tag = "🟢 improving" if mom >= 5 else "🟡 stable" if mom > -5 else "🔴 slowing"
            notes.append(f"Momentum: **{mom:+.1f}%** ({tag}) vs previous 7 days.")

        for n in notes[:6]:
            st.write("• " + n)

# ================= MIX SHIFT vs LY (NET values) + DEPENDENCY TABLE (THEMED) =================
        st.markdown("---")
        st.markdown("### 🔄 Mix Shift vs Last Year (Clear)")

        # ---------- helpers ----------
        def _safe_pct(a, b):
            a = float(a or 0)
            b = float(b or 0)
            return (a / b * 100.0) if b else 0.0

        def _normalize_channel(x: str) -> str:
            s = str(x).strip().lower()
            if s in ("", "nan", "none"):
                return "retail"
            if any(k in s for k in ["e-com", "ecom", "ecommerce", "online", "talabat"]):
                return "e-com"
            return "retail"

        def _build_channel_map(ch_df):
            # returns dict: py_norm -> ch_norm
            if ch_df is None or ch_df.empty:
                return {}
            need = {"PY Name 1", "Channels"}
            if not need.issubset(ch_df.columns):
                return {}

            tmp = ch_df.copy()
            tmp["_py_norm"] = tmp["PY Name 1"].astype(str).str.strip().str.lower()
            tmp["_ch_norm"] = tmp["Channels"].apply(_normalize_channel)
            tmp = tmp.dropna(subset=["_py_norm"])
            tmp = tmp.drop_duplicates(subset=["_py_norm"], keep="last")
            return dict(zip(tmp["_py_norm"], tmp["_ch_norm"]))

        _ch_map_dict = _build_channel_map(channels_df)

        def _add_channel(df_src):
            t = df_src.copy()
            if "PY Name 1" not in t.columns:
                t["_ch_norm"] = "retail"
                return t
            t["_py_norm"] = t["PY Name 1"].astype(str).str.strip().str.lower()
            t["_ch_norm"] = t["_py_norm"].map(_ch_map_dict).fillna("retail")
            return t

        def _filter_period(df_src, start_dt, end_dt):
            if df_src is None or df_src.empty or "Billing Date" not in df_src.columns:
                return df_src.iloc[0:0].copy() if df_src is not None else pd.DataFrame()
            d = df_src.copy()
            d["Billing Date"] = pd.to_datetime(d["Billing Date"], errors="coerce")
            return d[(d["Billing Date"] >= start_dt) & (d["Billing Date"] <= end_dt)].copy()

        def calc_net_mix(df_src):
            """
            NET Sales per channel = (ZFR+YKF2) - abs(YKRE+ZRE) - abs(YKS1+YKS2+ZCAN)
            Returns: (retail_net, ecom_net)
            """
            if df_src is None or df_src.empty:
                return 0.0, 0.0

            t = _add_channel(df_src)

            # ensure numeric
            if "Net Value" not in t.columns:
                return 0.0, 0.0
            t["Net Value"] = pd.to_numeric(t["Net Value"], errors="coerce").fillna(0.0)

            # ensure billing type
            if "Billing Type" not in t.columns:
                t["Billing Type"] = ""
            t["Billing Type"] = t["Billing Type"].astype(str).str.upper().str.strip()

            sales_mask   = t["Billing Type"].isin(["ZFR", "YKF2"])
            returns_mask = t["Billing Type"].isin(["YKRE", "ZRE"])
            cancel_mask  = t["Billing Type"].isin(["YKS1", "YKS2", "ZCAN"])

            sales   = t.loc[sales_mask].groupby("_ch_norm")["Net Value"].sum()
            returns = t.loc[returns_mask].groupby("_ch_norm")["Net Value"].sum().abs()
            cancel  = t.loc[cancel_mask].groupby("_ch_norm")["Net Value"].sum().abs()

            retail_net = float(sales.get("retail", 0.0) - returns.get("retail", 0.0) - cancel.get("retail", 0.0))
            ecom_net   = float(sales.get("e-com", 0.0) - returns.get("e-com", 0.0) - cancel.get("e-com", 0.0))
            return retail_net, ecom_net

        # ---------- current NET mix ----------
        cur_retail_net, cur_ecom_net = calc_net_mix(df_cur)
        cur_total_net = cur_retail_net + cur_ecom_net
        cur_retail_pct = _safe_pct(cur_retail_net, cur_total_net)
        cur_ecom_pct   = _safe_pct(cur_ecom_net, cur_total_net)

        # ---------- LY NET mix ----------
        ly_start = pd.to_datetime(cur_start) - pd.DateOffset(years=1)
        ly_end   = pd.to_datetime(cur_end)   - pd.DateOffset(years=1)

        ly_src = None
        if ytd_df is not None and not ytd_df.empty and {"Billing Date", "Net Value", "PY Name 1"}.issubset(ytd_df.columns):
            ly_src = _filter_period(ytd_df, ly_start, ly_end)

        if ly_src is None or ly_src.empty:
            ly_src = _filter_period(sales_df, ly_start, ly_end)

        ly_retail_net, ly_ecom_net = calc_net_mix(ly_src)
        ly_total_net = ly_retail_net + ly_ecom_net
        ly_retail_pct = _safe_pct(ly_retail_net, ly_total_net)
        ly_ecom_pct   = _safe_pct(ly_ecom_net, ly_total_net)

        shift_retail = cur_retail_pct - ly_retail_pct
        shift_ecom   = cur_ecom_pct - ly_ecom_pct

        # ---------- UI ----------
        c1, c2, c3 = st.columns(3)
        c1.metric("Retail Share (Current)", f"{cur_retail_pct:.0f}%", delta=f"{shift_retail:+.0f} pts vs LY")
        c2.metric("E-Com Share (Current)",  f"{cur_ecom_pct:.0f}%",   delta=f"{shift_ecom:+.0f} pts vs LY")
        c3.metric("Mix Signal", "🔴 Significant Shift" if abs(shift_retail) >= 8 else "🟢 Normal")

        st.caption(
            f"Current (NET): Retail {fmt_kd(cur_retail_net)} | E-Com {fmt_kd(cur_ecom_net)}  "
            f"|| LY (NET): Retail {fmt_kd(ly_retail_net)} | E-Com {fmt_kd(ly_ecom_net)}"
        )

 # ============================================================
        # ⚠️ Dependency Risk (Top Names + Share)  ✅ TRUE NET VALUES
        # Put this block under your Mix Shift section
        # Requires: pandas as pd, fmt_kd() already defined
        # ============================================================

        st.markdown("---")
        st.markdown("### ⚠️ Dependency Risk (Top Names + Share)")

        def _net_value_series(df):
            if df is None or df.empty or "Net Value" not in df.columns:
                return pd.Series(dtype="float")

            t = df.copy()
            t["Net Value"] = pd.to_numeric(t["Net Value"], errors="coerce").fillna(0.0)

            if "Billing Type" not in t.columns:
                t["Billing Type"] = ""
            bt = t["Billing Type"].astype(str).str.upper().str.strip()

            sales_mask   = bt.isin(["ZFR", "YKF2"])
            returns_mask = bt.isin(["YKRE", "ZRE"])
            cancel_mask  = bt.isin(["YKS1", "YKS2", "ZCAN"])

            net = pd.Series(0.0, index=t.index)
            net.loc[sales_mask]   = t.loc[sales_mask, "Net Value"]
            net.loc[returns_mask] = -t.loc[returns_mask, "Net Value"].abs()
            net.loc[cancel_mask]  = -t.loc[cancel_mask, "Net Value"].abs()
            return net

        def dependency_table_global_total(df_src, group_col, total_net, top_n=5, label="Name"):
            empty_tbl = pd.DataFrame({label: [], "NET (KD)": [], "Share %": []})

            if df_src is None or df_src.empty or group_col not in df_src.columns:
                return empty_tbl, 0.0, 0.0

            t = df_src.copy()
            t["_net"] = _net_value_series(t)

            # NET by group (can be positive/negative)
            g = t.groupby(group_col)["_net"].sum().sort_values(ascending=False)

            # Dependency focus: only positive groups (optional)
            g = g[g > 0].head(top_n)

            top_net = float(g.sum())
            share_pct = (top_net / total_net * 100.0) if total_net else 0.0

            out = g.reset_index()
            out.columns = [label, "_net"]
            out["Share %"] = (out["_net"] / total_net * 100.0).round(1) if total_net else 0.0
            out["NET (KD)"] = out["_net"].apply(fmt_kd)
            out = out[[label, "NET (KD)", "Share %"]]

            return out, float(share_pct), top_net

        # ✅ Use CURRENT period df (df_cur)
        dep_df = df_cur.copy() if df_cur is not None else pd.DataFrame()

        # ---- columns (edit if needed) ----
        SKU_COL  = "Material Description"
        CUST_COL = "PY Name 1"

        # ✅ ONE Global NET total for both tables
        dep_df["_net"] = _net_value_series(dep_df)
        global_total_net = float(dep_df["_net"].sum())

        # tables
        sku_tbl, sku_share, sku_top_net = dependency_table_global_total(
            dep_df, SKU_COL, global_total_net, top_n=5, label="Top 5 SKU"
        )

        cus_tbl, cus_share, cus_top_net = dependency_table_global_total(
            dep_df, CUST_COL, global_total_net, top_n=5, label="Top 5 Customer"
        )

        # ---- Side-by-side ----
        c1, c2 = st.columns(2)

        with c1:
            st.markdown("#### 🧾 Top 5 SKU Dependency (NET)")
            st.dataframe(sku_tbl, use_container_width=True, hide_index=True)
            st.caption(f"Top 5 NET: {fmt_kd(sku_top_net)} | Total NET (Current): {fmt_kd(global_total_net)}")

        with c2:
            st.markdown("#### 👤 Top 5 Customer Dependency (NET)")
            st.dataframe(cus_tbl, use_container_width=True, hide_index=True)
            st.caption(f"Top 5 NET: {fmt_kd(cus_top_net)} | Total NET (Current): {fmt_kd(global_total_net)}")

        st.info(f"Top 5 SKU = **{sku_share:.0f}%**  |  Top 5 Customer = **{cus_share:.0f}%**")

        # ---------------- C) Risk Radar (spike vs last 30 days) ----------------
        st.markdown("---")
        st.markdown("### 🚨 Risk Radar (Spikes vs Last 30 Days)")

        hist30_end = cur_end
        hist30_start = cur_end - pd.Timedelta(days=30)

        hist30 = sales_df[(sales_df["Billing Date"] >= hist30_start) & (sales_df["Billing Date"] <= hist30_end)].copy()
        hist30["Billing Date"] = pd.to_datetime(hist30["Billing Date"], errors="coerce")
        hist30["Net Value"] = pd.to_numeric(hist30["Net Value"], errors="coerce").fillna(0)

        hist_sales_only = _sum_sales_only(hist30)
        hist_ret = abs(_sum_returns(hist30))
        hist_can = abs(_sum_cancels(hist30))

        hist_ret_pct = _safe_pct(hist_ret, hist_sales_only)
        hist_can_pct = _safe_pct(hist_can, hist_sales_only)

        # current rates (use your already computed return_pct/cancel_pct if present)
        cur_sales_only = _sum_sales_only(df_cur)
        cur_ret = abs(_sum_returns(df_cur))
        cur_can = abs(_sum_cancels(df_cur))

        cur_ret_pct = _safe_pct(cur_ret, cur_sales_only)
        cur_can_pct = _safe_pct(cur_can, cur_sales_only)

        spike_ret = cur_ret_pct - hist_ret_pct
        spike_can = cur_can_pct - hist_can_pct

        r1, r2, r3, r4 = st.columns(4)
        r1.metric("Return % (Current)", f"{cur_ret_pct:.1f}%")
        r2.metric("Return % (Last 30d)", f"{hist_ret_pct:.1f}%")
        r3.metric("Cancel % (Current)", f"{cur_can_pct:.1f}%")
        r4.metric("Cancel % (Last 30d)", f"{hist_can_pct:.1f}%")

        if spike_ret > 1:
            st.warning(f"🔴 Return spike: +{spike_ret:.1f}% vs last 30d.")
        if spike_can > 1:
            st.warning(f"🟠 Cancel spike: +{spike_can:.1f}% vs last 30d.")
        if spike_ret <= 1 and spike_can <= 1:
            st.success("✅ No major return/cancel spikes vs last 30 days.")

        # Identify top drivers for returns/cancels (names + values)
        col_a, col_b = st.columns(2)

        with col_a:
            st.markdown("**Top Return Drivers (Current Period)**")
            if "Billing Type" in df_cur.columns and "Driver Name EN" in df_cur.columns:
                dtmp = df_cur.copy()
                dtmp["Billing Type"] = dtmp["Billing Type"].astype(str).str.upper().str.strip()
                ret_sm = dtmp[dtmp["Billing Type"].isin({"YKRE","ZRE"})].groupby("Driver Name EN")["Net Value"].sum().abs().sort_values(ascending=False).head(5)
                if len(ret_sm):
                    ret_tbl = ret_sm.reset_index()
                    ret_tbl.columns = ["Salesman", "Returns"]
                    render_table(ret_tbl, hide_index=True, formats={"Returns": "{:,.0f}"})
                else:
                    st.caption("No return records found.")
            else:
                st.caption("Required columns missing.")

        with col_b:
            st.markdown("**Top Cancel Drivers (Current Period)**")
            if "Billing Type" in df_cur.columns and "Driver Name EN" in df_cur.columns:
                dtmp = df_cur.copy()
                dtmp["Billing Type"] = dtmp["Billing Type"].astype(str).str.upper().str.strip()
                can_sm = dtmp[dtmp["Billing Type"].isin({"YKS1","YKS2","ZCAN"})].groupby("Driver Name EN")["Net Value"].sum().abs().sort_values(ascending=False).head(5)
                if len(can_sm):
                    can_tbl = can_sm.reset_index()
                    can_tbl.columns = ["Salesman", "Cancels"]
                    render_table(can_tbl, hide_index=True, formats={"Cancels": "{:,.0f}"})
                else:
                    st.caption("No cancel records found.")
            else:
                st.caption("Required columns missing.")

 

        # # ---------------- E) Top 5 SKUs by Category ----------------
        # st.markdown("---")
        # st.markdown("### 🏷️ Top 5 SKUs by Category (Value)")

        # possible_cat_cols = ["Category", "Material Group", "Product Group", "Brand", "Division", "Group"]
        # cat_col = next((c for c in possible_cat_cols if c in df_cur.columns), None)

        # if cat_col and "Material Description" in df_cur.columns:
        #     cat_vals = df_cur[cat_col].dropna().astype(str).unique().tolist()
        #     cat_vals = sorted(cat_vals)

        #     sel_cat = st.selectbox("Select Category", cat_vals)
        #     df_cat = df_cur[df_cur[cat_col].astype(str) == str(sel_cat)]

        #     sku_cat = df_cat.groupby("Material Description")["Net Value"].sum().sort_values(ascending=False).head(5).reset_index()
        #     sku_cat.columns = ["SKU", "Sales"]
        #     render_table(sku_cat, hide_index=True, formats={"Sales": "{:,.0f}"})

        #     st.caption(f"Category column used: {cat_col}")
        # else:
        #     st.info("Category-wise view not available (need Category/Group column + Material Description).")

        # ---------------- F) GM Action Plan (Auto, non-repeat) ----------------
        st.markdown("---")
        st.markdown("### ✅ GM Action Plan (Auto)")

        actions = []

        # Based on spikes and momentum
        if spike_can > 1:
            actions.append("Warehouse/Planning: Cancel spike vs last 30 days → check stock accuracy, picking, delivery schedule, and allocation.")
        elif cur_can_pct >= 3:
            actions.append("Operations: Cancel % is high → review cancel reasons and fix root causes.")

        if spike_ret > 1:
            actions.append("QA/Warehouse: Return spike vs last 30 days → check quality, handling, expiry/temperature issues, and top return SKUs/customers.")
        elif cur_ret_pct >= 3:
            actions.append("QA: Return % is high → run return reason audit and corrective actions.")

        if mom is not None:
            if mom < -5:
                actions.append("Sales Leaders: Momentum slowing → enforce weekly push plan, focus top customers and hero SKUs.")
            elif mom > 5:
                actions.append("Management: Momentum improving → scale winning actions (promotions/visibility/assortment) and protect availability.")

        # Mix shift meaning (only if we calculated)
        try:
            if abs(shift_retail) >= 8:
                actions.append("Key Accounts: Channel mix shifted strongly vs LY → review execution, promotions, and supply planning by channel.")
        except Exception:
            pass

        # Dependency actions
        try:
            if top3_share >= 55:
                actions.append("Category Manager: High SKU dependency → diversify mix, push secondary SKUs, reduce single-SKU risk.")
        except Exception:
            pass

        if not actions:
            actions.append("✅ No critical alerts. Maintain execution and review weekly KPIs.")

        for a in actions[:7]:
            st.write("• " + a)
            

        # ------------------------------------------------
        # 8) Top Salesmen Spotlight (Net + Risk)
        # ------------------------------------------------
        st.markdown("---")
        st.subheader("🏆 GM Spotlight – Top Salesmen (Net + Risk)")

        if "Driver Name EN" in df.columns:
            g = df.groupby(["Driver Name EN", "Billing Type"])["Net Value"].sum().unstack(fill_value=0)

            sales_sm = g.reindex(columns=list(SALES_BT), fill_value=0).sum(axis=1)
            ret_sm   = g.reindex(columns=list(RETURN_BT), fill_value=0).sum(axis=1).abs()
            can_sm   = g.reindex(columns=list(CANCEL_BT), fill_value=0).sum(axis=1).abs()

            net_sm = sales_sm - ret_sm - can_sm

            sm_tbl = pd.DataFrame({
                "Salesman": net_sm.index,
                "Net Sales": net_sm.values,
                "Return %": np.where(sales_sm.values != 0, (ret_sm.values / sales_sm.values * 100).round(1), 0),
                "Cancel %": np.where(sales_sm.values != 0, (can_sm.values / sales_sm.values * 100).round(1), 0),
            })
            sm_tbl["Tag"] = sm_tbl.apply(lambda r: gm_tag(r["Return %"], r["Cancel %"]), axis=1)

            sm_tbl = sm_tbl.sort_values("Net Sales", ascending=False).head(top_n)

            render_table(
                sm_tbl,
                hide_index=True,
                formats={
                    "Net Sales": "{:,.0f}",
                    "Return %": "{:.1f}%",
                    "Cancel %": "{:.1f}%"
                }
            )
        else:
            st.info("Salesman column not available.")

        # ------------------------------------------------
        # 9) Customer Risk Table (Returns focus)
        # ------------------------------------------------
        st.markdown("---")
        st.subheader("🔻 Customer Risk – Returns Focus (Top 10)")

        if "PY Name 1" in df.columns:
            g2 = df.groupby(["PY Name 1", "Billing Type"])["Net Value"].sum().unstack(fill_value=0)
            sales_c = g2.reindex(columns=list(SALES_BT), fill_value=0).sum(axis=1)
            ret_c   = g2.reindex(columns=list(RETURN_BT), fill_value=0).sum(axis=1).abs()
            can_c   = g2.reindex(columns=list(CANCEL_BT), fill_value=0).sum(axis=1).abs()
            net_c   = sales_c - ret_c - can_c

            cust_tbl = pd.DataFrame({
                "Customer": net_c.index,
                "Net Sales": net_c.values,
                "Returns": ret_c.values,
                "Return %": np.where(sales_c.values != 0, (ret_c.values / sales_c.values * 100).round(1), 0),
            }).sort_values(["Return %", "Returns"], ascending=False).head(10)

            render_table(
                cust_tbl,
                hide_index=True,
                formats={
                    "Net Sales": "{:,.0f}",
                    "Returns": "{:,.0f}",
                    "Return %": "{:.1f}%"
                }
            )
        else:
            st.info("Customer column not available.")

        # ------------------------------------------------
        # 10) GM Action Notes
        # ------------------------------------------------
        st.markdown("---")
        st.subheader("✅ GM Action Notes")

        notes = []
        if return_pct >= 3:
            notes.append(f"Returns are high ({return_pct:.1f}%). Check top return customers and handling / expiry.")
        if cancel_pct >= 3:
            notes.append(f"Cancels are high ({cancel_pct:.1f}%). Review warehouse cancel reasons and stock accuracy.")
        if mix_total > 0 and ecom_mix >= 50:
            notes.append("E-com share is high. Ensure retail execution (visibility + availability) is not dropping.")

        if notes:
            for n in notes:
                st.write("• " + n)
        else:
            st.success("✅ No major GM risks detected for this period.")
            
            
            
# ---------------- Customer Insights Page (Full Pro Version, Fixed) ----------------
elif choice == texts[lang]["customer_insights"]:
    st.title(texts[lang]["customer_insights_title"])

    # Pre-checks
    if "data_loaded" not in st.session_state:
        st.warning(texts[lang]["no_data_warning"])
        st.stop()

    df_rfm = st.session_state["sales_df"].copy()
    if df_rfm.empty:
        st.warning(texts[lang]["rfm_no_data"])
        st.stop()

    # Apply salesman filter if applicable
    if user_role == "salesman" and salesman_name:
        df_rfm = df_rfm[df_rfm["Driver Name EN"] == salesman_name]

    # Detect columns robustly
    def find_col(df, candidates):
        for n in candidates:
            if n in df.columns:
                return n
        for c in df.columns:
            lc = c.lower()
            for n in candidates:
                if n.lower() in lc:
                    return c
        return None

    cust_col = find_col(df_rfm, ["SP Name1", "SP Name 1", "SP_Name1", "Customer", "PY Name1", "PY Name 1"])
    date_col = find_col(df_rfm, ["Billing Date", "billing date", "Date"])
    amount_col = find_col(df_rfm, ["Net Value", "NetAmount", "Net Amount", "Amount", "Sales Amount"])
    material_col = find_col(df_rfm, ["Material Description", "Item", "Product", "Material"])

    if None in [cust_col, date_col, amount_col]:
        st.warning(texts[lang]["rfm_no_data"])
        st.stop()

    # Normalize date col
    df_rfm[date_col] = pd.to_datetime(df_rfm[date_col], errors="coerce")
    today = pd.Timestamp.today().normalize()

    # --- Fixed robust RFM aggregation ---
    rfm_group = df_rfm.groupby(cust_col)

    rfm_agg = pd.DataFrame({
        "Customer": rfm_group.apply(lambda g: g.name),
        "Recency": rfm_group[date_col].max().apply(lambda d: (today - d).days),
        "Frequency": rfm_group[date_col].count(),
        "Monetary": rfm_group[amount_col].sum()
    }).reset_index(drop=True)

    rfm_agg = rfm_agg[rfm_agg["Monetary"] > 0].set_index("Customer")

    if rfm_agg.empty:
        st.warning(texts[lang]["rfm_no_data"])
        st.stop()

    # Safe qcut as before
    def safe_qcut(series, q=4, reverse=False):
        s = series.copy().fillna(series.max() + 1)
        unique_vals = pd.unique(s)
        n_unique = len(unique_vals)
        if n_unique == 1:
            return pd.Series([1]*len(s), index=s.index)
        if n_unique < q:
            ranks = s.rank(method='dense', ascending=not reverse)
            return ranks.astype(int)
        labels = list(range(q, 0, -1)) if reverse else list(range(1, q+1))
        try:
            return pd.qcut(s, q=q, labels=labels, duplicates='drop')
        except Exception:
            ranks = s.rank(method='dense', ascending=not reverse)
            return ranks.astype(int)

    rfm_agg["R_Score"] = safe_qcut(rfm_agg["Recency"], q=4, reverse=True).astype(int)
    rfm_agg["F_Score"] = safe_qcut(rfm_agg["Frequency"], q=4).astype(int)
    rfm_agg["M_Score"] = safe_qcut(rfm_agg["Monetary"], q=4).astype(int)
    rfm_agg["RFM_Score"] = rfm_agg["R_Score"].astype(str) + rfm_agg["F_Score"].astype(str) + rfm_agg["M_Score"].astype(str)

    # Segmentation
    def rfm_segment(row):
        if row["RFM_Score"] in ["444", "443", "434", "433"]:
            return "Champions"
        if row["R_Score"] >= 3 and row["F_Score"] >= 3:
            return "Loyal Customers"
        if row["R_Score"] >= 3 and row["M_Score"] >= 3:
            return "Potential Loyalists"
        if row["R_Score"] >= 3:
            return "New Customers"
        if row["R_Score"] <= 2 and row["F_Score"] >= 2 and row["M_Score"] >= 2:
            return "At Risk"
        if row["R_Score"] <= 1 and row["F_Score"] >= 2 and row["M_Score"] >= 2:
            return "Hibernating"
        return "Others"

    rfm_agg["Segment"] = rfm_agg.apply(rfm_segment, axis=1)

    # --- Layout Tabs: RFM, Cohort, Weekly CRM ---
    tab_weekly, tab_360 = st.tabs([
    # texts[lang]["rfm_analysis_sub"], 
    # texts[lang]["rfm_cohort_sub"], 
    "CRM & Weekly Operations",
    "Customer 360°"  # ← NEW TAB
])

    # # ---------------- RFM Tab ----------------
    # with tab_rfm:
    #     st.subheader(texts[lang]["rfm_table_sub"])
    #     display_rfm = rfm_agg.copy()
    #     display_rfm[["Recency","Frequency","Monetary"]] = display_rfm[["Recency","Frequency","Monetary"]].astype(int)
    #     st.dataframe(display_rfm.sort_values("Monetary", ascending=False), use_container_width=True, hide_index=True)

    #     # Download with safe sheet name
    #     ts = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    #     safe_sheet = "RFM_Analysis"[:31]
    #     if st.download_button(
    #         texts[lang]["rfm_download"],
    #         data=to_excel_bytes(display_rfm.reset_index(), sheet_name=safe_sheet),
    #         file_name=f"rfm_analysis_{ts}.xlsx"
    #     ):
    #         st.session_state["audit_log"].append({"user": username, "action":"download","details":f"rfm_analysis_{ts}.xlsx","timestamp":datetime.now().isoformat()})

    #     # Segment Pie + metrics
    #     st.subheader("RFM Segment Distribution")
    #     seg_counts = display_rfm["Segment"].value_counts().reset_index()
    #     seg_counts.columns = ["Segment", "Count"]
    #     seg_counts["Percentage"] = (seg_counts["Count"]/seg_counts["Count"].sum()*100).round(1)
    #     seg_avg = display_rfm.groupby("Segment")["Monetary"].mean().round(2).reset_index().rename(columns={"Monetary":"Avg Monetary"})
    #     seg_counts = seg_counts.merge(seg_avg, on="Segment", how="left")
    #     fig_seg = px.pie(seg_counts, names="Segment", values="Count", hole=0.35, hover_data=["Percentage","Avg Monetary"], title="RFM Segment Distribution")
    #     fig_seg.update_traces(textinfo='percent+label')
    #     st.plotly_chart(fig_seg, use_container_width=True)

    #     st.subheader("Key Metrics per Segment")
    #     seg_metrics = display_rfm.groupby("Segment").agg(
    #         mean_recency=("Recency","mean"),
    #         mean_frequency=("Frequency","mean"),
    #         mean_monetary=("Monetary","mean"),
    #         count=("R_Score","count")
    #     ).round(2).rename(columns={"count":"Count"})
    #     st.dataframe(seg_metrics, use_container_width=True, hide_index=True)

    #     st.subheader("Prescriptive Actions per Segment")
    #     recs = {
    #         "Champions":"Reward with exclusive offers & loyalty programs.",
    #         "Loyal Customers":"Upsell & referrals.",
    #         "Potential Loyalists":"Nurture with targeted campaigns.",
    #         "New Customers":"Onboard & incentivize repeat purchase.",
    #         "At Risk":"Win-back campaigns and surveys.",
    #         "Hibernating":"Reactivate with promotions.",
    #         "Others":"Investigate further."
    #     }
    #     for s in seg_metrics.index:
    #         st.write(f"- **{s}**: {recs.get(s,'General engagement strategies recommended.')}")

    #     st.subheader(texts[lang]["rfm_chart_sub"])
    #     fig_rfm = px.scatter(display_rfm.reset_index(), x="Recency", y="Monetary", size="Frequency", color="Segment",
    #                          hover_name=display_rfm.reset_index()["Customer"], title="RFM Scatter (Recency vs Monetary; size=Frequency)")
    #     st.plotly_chart(fig_rfm, use_container_width=True)

    # # ---------------- Cohort Tab (Fixed) ----------------
    # with tab_cohort:
    #     st.subheader(texts[lang]["rfm_cohort_sub"])
    #     st.info(texts[lang]["rfm_cohort_info"])

    #     df_cohort = df_rfm[[cust_col, date_col, amount_col]].dropna().copy().rename(columns={cust_col:"Customer", date_col:"Billing Date"})

    #     # Convert periods to strings
    #     df_cohort["Cohort_Month_Str"] = df_cohort.groupby("Customer")["Billing Date"].transform("min").dt.to_period("M").astype(str)
    #     df_cohort["Period_Month_Str"] = df_cohort["Billing Date"].dt.to_period("M").astype(str)
    #     df_cohort["Cohort_Index"] = (pd.to_datetime(df_cohort["Period_Month_Str"]) - pd.to_datetime(df_cohort["Cohort_Month_Str"])).dt.days // 30

    #     cohort_summary = df_cohort.groupby(["Cohort_Month_Str","Cohort_Index"]).agg(
    #         Customer=("Customer","nunique"),
    #         Monetary=(amount_col,"mean")
    #     ).reset_index()

    #     if not cohort_summary.empty:
    #         cohort_pivot = cohort_summary.pivot(index="Cohort_Month_Str", columns="Cohort_Index", values="Monetary").fillna(0)
    #         fig_cohort = px.imshow(cohort_pivot, labels=dict(x="Months after acquisition", y="Cohort", color="Avg Monetary"),
    #                                title="Cohort Monetary Heatmap", text_auto=True, aspect="auto")
    #         st.plotly_chart(fig_cohort, use_container_width=True)

    #         st.subheader(texts[lang]["rfm_cohort_table_sub"])
    #         cohort_table = cohort_summary.pivot(index="Cohort_Month_Str", columns="Cohort_Index", values="Customer").fillna(0).astype(int)
    #         st.dataframe(cohort_table, use_container_width=True, hide_index=True)
    #     else:
    #         st.warning(texts[lang]["rfm_cohort_no_data"])

    # ---------------- CRM & Weekly Operations Tab ----------------
    with tab_weekly:
        st.subheader("CRM Operations — Weekly Tracker, Products & Manager Dashboard")

        # Date selector
        col_left, col_right = st.columns([3,1])
        with col_left:
            auto_date = st.toggle("Use current date automatically", value=True)
            selected_date = datetime.now().date() if auto_date else st.date_input("Select visit date manually", datetime.now().date())
            st.session_state["visit_date"] = selected_date

        with col_right:
            show_manager = st.checkbox("Show Manager KPIs", value=True)
            refresh = st.button("🔄 Refresh")  # triggers rerun

        # Prepare commonly used variables
        ytd_df = st.session_state.get("ytd_df", pd.DataFrame()).copy()
        sales_df = st.session_state.get("sales_df", pd.DataFrame()).copy()

        # Ensure date cols
        if date_col in ytd_df.columns:
            ytd_df[date_col] = pd.to_datetime(ytd_df[date_col], errors="coerce")
        if date_col in sales_df.columns:
            sales_df[date_col] = pd.to_datetime(sales_df[date_col], errors="coerce")

        # ---------------- Weekly Visit Tracker (robust) ----------------
        st.markdown("### Weekly Visit Tracker")
        last_3_months = pd.Timestamp(selected_date) - pd.DateOffset(months=3)
        recent_ytd = ytd_df[ytd_df.get(date_col, pd.Series()) >= last_3_months] if not ytd_df.empty else pd.DataFrame()
        customer_list = pd.Series(recent_ytd.get(cust_col, pd.Series()).dropna().unique()).astype(str) if not recent_ytd.empty else pd.Series(sales_df.get(cust_col, pd.Series()).dropna().unique()).astype(str)

        if customer_list.empty:
            st.info("No customers found in YTD or Sales for the last 3 months.")
        else:
            # last 7-day window
            days_dt = [pd.Timestamp(selected_date) - pd.Timedelta(days=i) for i in range(6, -1, -1)]
            days_str = [d.strftime("%Y-%m-%d") for d in days_dt]
            sales_window_start = days_dt[0]
            sales_window_end = pd.Timestamp(selected_date) + pd.Timedelta(days=1)

            sales_last7 = sales_df[(sales_df[date_col] >= sales_window_start) & (sales_df[date_col] < sales_window_end)].copy()
            sales_last7[cust_col] = sales_last7[cust_col].astype(str)
            sales_last7[amount_col] = pd.to_numeric(sales_last7[amount_col], errors="coerce").fillna(0.0)
            sales_last7["__date_str"] = sales_last7[date_col].dt.strftime("%Y-%m-%d")

            pivot7 = (sales_last7.groupby([cust_col, "__date_str"])[amount_col].sum().reset_index()
                      .pivot(index=cust_col, columns="__date_str", values=amount_col).reindex(columns=days_str, fill_value=0.0).reset_index())
            pivot7 = pivot7.rename(columns={cust_col: "Customer"})
            base = pd.DataFrame({"Customer": customer_list})
            visit_df = base.merge(pivot7, on="Customer", how="left").fillna(0.0)
            visit_df.insert(1, "Visit Date", selected_date)

            existing_days = [c for c in days_str if c in visit_df.columns]
            visit_df["Weekly Total"] = visit_df[existing_days].sum(axis=1) if existing_days else 0

            # 4-week totals
            end_date = pd.Timestamp(selected_date)
            start_date = end_date - pd.Timedelta(weeks=4)
            recent_sales = sales_df[(sales_df[date_col] >= start_date) & (sales_df[date_col] <= end_date)].copy()
            recent_sales[cust_col] = recent_sales[cust_col].astype(str)
            recent_sales[amount_col] = pd.to_numeric(recent_sales[amount_col], errors="coerce").fillna(0.0)
            recent_sales["Week_Number"] = ((recent_sales[date_col] - start_date).dt.days // 7) + 1
            recent_sales.loc[recent_sales["Week_Number"] > 4, "Week_Number"] = 4
            week_totals = (recent_sales.groupby([cust_col, "Week_Number"])[amount_col].sum().unstack(fill_value=0).reset_index().rename(columns={cust_col:"Customer"}))
            week_cols = [c for c in week_totals.columns if c != "Customer"]
            if not week_cols:
                for i in range(1,5):
                    week_totals[i] = 0
                week_cols = [1,2,3,4]
            ordered_week_cols = sorted(week_cols, key=lambda x:int(x))
            week_totals = week_totals[["Customer"] + ordered_week_cols]
            week_headers = []
            for i in range(len(ordered_week_cols)):
                wstart = (start_date + pd.Timedelta(days=i*7)).strftime("%b %d")
                wend = (start_date + pd.Timedelta(days=(i+1)*7-1)).strftime("%b %d")
                week_headers.append(f"Week {i+1} ({wstart}-{wend})")
            new_cols = ["Customer"] + week_headers
            if len(new_cols) != len(week_totals.columns):
                new_cols = ["Customer"] + [f"Week {i+1}" for i in range(len(week_totals.columns)-1)]
            week_totals.columns = new_cols
            visit_df = visit_df.merge(week_totals, on="Customer", how="left").fillna(0.0)

            # Total Sales (3 months)
            total_sales_3m = sales_df[sales_df[date_col] >= last_3_months].groupby(cust_col)[amount_col].sum().reset_index().rename(columns={cust_col:"Customer", amount_col:"Total Sales"})
            visit_df = visit_df.merge(total_sales_3m, on="Customer", how="left").fillna(0.0)

            # Alerts & recommended action
            def compute_alert(row):
                q80 = visit_df["Total Sales"].quantile(0.8) if "Total Sales" in visit_df.columns else 0
                q50 = visit_df["Total Sales"].quantile(0.5) if "Total Sales" in visit_df.columns else 0
                if row.get("Weekly Total",0) == 0:
                    if row.get("Total Sales",0) >= q80:
                        return "🔴 High", "Visit immediately"
                    if row.get("Total Sales",0) >= q50:
                        return "🟠 Medium", "Call / Email"
                    return "🟢 Low", "Monitor"
                return "✅ Visited", "No action"
            visit_df[["Alert Level","Recommended Action"]] = visit_df.apply(lambda r: pd.Series(compute_alert(r)), axis=1)

            # numeric formatting
            numeric_cols = [c for c in existing_days + ["Weekly Total","Total Sales"] + week_headers if c in visit_df.columns]
            if numeric_cols:
                visit_df[numeric_cols] = visit_df[numeric_cols].fillna(0).round(0).astype(int)

            # show table
            with st.expander("Show Weekly Visit Table", expanded=True):
                st.dataframe(visit_df.sort_values(["Total Sales","Weekly Total"], ascending=[False,False]), use_container_width=True, hide_index=True)

            # manager KPIs
            if show_manager:
                st.markdown("### Manager Dashboard")
                col1, col2, col3, col4 = st.columns(4)
                last7_revenue = sales_last7[amount_col].sum() if not sales_last7.empty else 0
                col1.metric("Revenue (Last 7 days)", f"KD {last7_revenue:,.0f}")
                col2.metric("Customers Visited (7d)", int((visit_df["Weekly Total"]>0).sum()))
                col3.metric("High-Value Missed", int((visit_df["Alert Level"]=="🔴 High").sum()))
                col4.metric("Avg Weekly Sales per Customer", f"KD {visit_df['Weekly Total'].mean():,.0f}")

                top_missed = visit_df[visit_df["Alert Level"]=="🔴 High"].sort_values("Total Sales", ascending=False).head(10)
                if not top_missed.empty:
                    st.subheader("Top High-Value Missed Customers")
                    st.bar_chart(top_missed.set_index("Customer")["Total Sales"])

            # Save Visit as record (mini CRM)
            st.markdown("### Visit Planner / Notes")
            planner_customer = st.selectbox("Select Customer to plan visit", options=sorted(visit_df["Customer"].astype(str).unique()))
            note = st.text_area("Note / Follow-up", key=f"note_{planner_customer}")
            next_visit = st.date_input("Next Visit Date", value=(pd.Timestamp(selected_date)+pd.Timedelta(days=7)).date(), key=f"nextvisit_{planner_customer}")
            if st.button("💾 Save Visit Plan"):
                if "visit_plans" not in st.session_state:
                    st.session_state["visit_plans"] = []
                st.session_state["visit_plans"].append({
                    "Customer": planner_customer,
                    "Next Visit": next_visit.isoformat() if hasattr(next_visit, "isoformat") else str(next_visit),
                    "Note": note,
                    "Created By": username,
                    "Created At": datetime.now().isoformat()
                })
                st.success(f"Saved visit plan for {planner_customer}")

            if "visit_plans" in st.session_state and st.session_state["visit_plans"]:
                st.subheader("Saved Visit Plans")
                st.dataframe(pd.DataFrame(st.session_state["visit_plans"]), use_container_width=True, hide_index=True)
                # Download visit plans
                safe_file_name = f"visit_plans_{selected_date}.csv"
                st.download_button("⬇️ Download Visit Plans", data=pd.DataFrame(st.session_state["visit_plans"]).to_csv(index=False).encode('utf-8'), file_name=safe_file_name, mime="text/csv")

        # ---------------- 15-Day Product Analysis ----------------
        st.markdown("### Customer Product Activity (Last 15 Days)")
        product_start_date = pd.Timestamp(selected_date) - pd.Timedelta(days=15)
        prod_sales = sales_df[(sales_df[date_col] >= product_start_date) & (sales_df[date_col] <= pd.Timestamp(selected_date))].copy()

        if prod_sales.empty or material_col is None:
            st.info("No product-sales data available for the last 15 days.")
        else:
            prod_sales[cust_col] = prod_sales[cust_col].astype(str)
            prod_sales[amount_col] = pd.to_numeric(prod_sales[amount_col], errors="coerce").fillna(0.0)
            all_products = sorted(sales_df[material_col].dropna().unique())

            prod_summary = prod_sales.groupby([cust_col, material_col])[amount_col].sum().reset_index().rename(columns={cust_col:"Customer", material_col:"Product", amount_col:"Sales Amount"})
            customers_prod = sorted(prod_summary["Customer"].unique())

            if customers_prod:
                sel_cust = st.selectbox("Select a Customer to inspect products", options=customers_prod)
                sold_by_cust = prod_summary[prod_summary["Customer"]==sel_cust].sort_values("Sales Amount", ascending=False)
                sold_set = set(sold_by_cust["Product"].dropna())
                not_sold = [p for p in all_products if p not in sold_set]
                df_not_sold = pd.DataFrame({"Product": not_sold, "Status":"❌ Not Purchased"})

                with st.expander(f"{sel_cust} - Purchased (Last 15 Days)", expanded=True):
                    if not sold_by_cust.empty:
                        st.dataframe(sold_by_cust, use_container_width=True, hide_index=True)
                    else:
                        st.info("No purchases by this customer in last 15 days.")


                with st.expander(f"{sel_cust} - Not Purchased (Last 15 Days)", expanded=False):
                    if not df_not_sold.empty:
                        st.dataframe(df_not_sold, use_container_width=True, hide_index=True)
                    else:
                        st.info("All products purchased!")


                # Safe sheet name
                safe_sheet_name = (sel_cust+"_Purchased")[:31]
                if st.download_button(f"⬇️ Download {sel_cust} Purchased (15d)",
                                      data=to_excel_bytes(sold_by_cust, sheet_name=safe_sheet_name),
                                      file_name=f"{sel_cust}_purchased_15days_{selected_date}.xlsx"):
                    st.success("Download ready!")

    # ──────────────────────── TAB 4: Customer 360° (FIXED) ────────────────────────
    with tab_360:
        st.markdown("### **Customer 360° – One-Click Profile**")
        
        # === 1. COLUMN DETECTION (Safe & Reusable) ===
        def find_col(df, candidates):
            for c in candidates:
                if c in df.columns:
                    return c
            return None

        cust_col = find_col(df_rfm, ["SP Name1", "SP Name 1", "Customer", "PY Name 1", "Customer Name"])
        date_col = find_col(df_rfm, ["Billing Date", "Date", "Invoice Date"])
        amount_col = find_col(df_rfm, ["Net Value", "Amount", "Sales"])
        driver_col = find_col(df_rfm, ["Driver Name EN", "Salesman", "Driver", "Rep"])
        material_col = find_col(df_rfm, ["Material", "Material Description", "Item", "SKU"])

        if not all([cust_col, date_col, amount_col]):
            st.error("Missing required columns. Check your Excel file.")
            st.stop()

        # === 2. CUSTOMER SELECTOR ===
        all_customers = sorted(df_rfm[cust_col].dropna().unique())
        selected_cust = st.selectbox("**Select Customer**", all_customers, key="cust360_select")

        # === 3. FILTER DATA ===
        cust_sales = df_rfm[df_rfm[cust_col] == selected_cust].copy()
        cust_sales[date_col] = pd.to_datetime(cust_sales[date_col], errors='coerce')
        cust_sales = cust_sales.dropna(subset=[date_col])  # Remove invalid dates
        today = pd.Timestamp.today()

        if cust_sales.empty:
            st.warning(f"No sales data for **{selected_cust}**")
        else:
            # ========================================
            # KPI CARDS – CUSTOMER 360°
            # ========================================
            col1, col2, col3, col4 = st.columns(4)
            
            # --- 1. ENSURE amount_col is numeric ---
            cust_sales[amount_col] = pd.to_numeric(cust_sales[amount_col], errors='coerce').fillna(0)

            # --- 2. NORMALIZE Billing Type ---
            if "Billing Type" in cust_sales.columns:
                cust_sales["Billing Type"] = cust_sales["Billing Type"].astype(str).str.strip().str.upper()

            # --- 3. TOTAL SALES & ORDERS ---
            total_sales = cust_sales[amount_col].sum()
            order_count = len(cust_sales)
            last_visit = cust_sales[date_col].max()
            days_since = (today - last_visit).days if pd.notna(last_visit) else 999

            # --- 4. RETURNS ONLY (YKRE, ZRE) – CANCELLATIONS EXCLUDED ---
            return_codes = ["YKRE", "ZRE"]  # Only Returns
            if "Billing Type" in cust_sales.columns:
                returns_mask = cust_sales["Billing Type"].isin(return_codes)
            else:
                returns_mask = cust_sales[date_col].notna() & False  # no returns if column missing
            returns_df = cust_sales[returns_mask]

            # --- 5. RETURN VALUE = ABSOLUTE SUM (handles negative values) ---
            returns_value = returns_df[amount_col].abs().sum()
            return_rate = (returns_value / total_sales * 100) if total_sales > 0 else 0

            # === KPI CARD 1: Total Sales ===
            col1.metric(
                label="**Total Sales**",
                value=f"KD {total_sales:,.0f}",
                delta=None
            )

            # === KPI CARD 2: Orders ===
            col2.metric(
                label="**Orders**",
                value=f"{order_count:,}",
                delta=None
            )

            # === KPI CARD 3: Last Visit ===
            if pd.notna(last_visit):
                col3.metric(
                    label="**Last Visit**",
                    value=last_visit.strftime("%b %d, %Y"),
                    delta=f"{days_since} days ago" if days_since <= 365 else "Over 1 year",
                    delta_color="inverse" if days_since > 30 else "normal"
                )
            else:
                col3.metric(
                    label="**Last Visit**",
                    value="Never",
                    delta="No data"
                )

            # === KPI CARD 4: Return Rate + Value (COMBINED) ===
            if returns_value > 0:
                col4.metric(
                    label="**Return Rate**",
                    value=f"{return_rate:.2f}%",
                    delta=f"KD {returns_value:,.0f} returned",
                    delta_color="inverse"  # Red = high return
                )
            else:
                col4.metric(
                    label="**Return Rate**",
                    value="0.00%",
                    delta="No returns",
                    delta_color="normal"  # Green = good
                )
            
            # === MINI TABS ===
            mini_tab1, mini_tab2, mini_tab3, mini_tab4, mini_tab5 = st.tabs([
                "Sales Trend", "RFM", "Visits", "Issues", "Actions"
            ])

            # ───────── Mini Tab 1: Sales Trend ─────────
            with mini_tab1:
                daily = cust_sales.groupby(cust_sales[date_col].dt.date)[amount_col].sum().reset_index()
                daily.columns = ["Date", "Sales"]
                fig = px.line(daily, x="Date", y="Sales", title="Sales Trend", markers=True)
                fig.update_layout(height=300)
                st.plotly_chart(fig, use_container_width=True)

            # ───────── Mini Tab 2: RFM ─────────
            with mini_tab2:
                if selected_cust in rfm_agg.index:
                    r = rfm_agg.loc[selected_cust]
                    c1, c2, c3 = st.columns(3)
                    c1.metric("**Recency**", f"{int(r['Recency'])} days")
                    c2.metric("**Frequency**", f"{int(r['Frequency'])}")
                    c3.metric("**Monetary**", f"KD {r['Monetary']:,.0f}")
                    st.success(f"**Segment:** {r['Segment']}")
                else:
                    st.info("RFM not calculated yet")

            # ───────── Mini Tab 3: Visits ─────────
            with mini_tab3:
                if driver_col and driver_col in cust_sales.columns:
                    visits = cust_sales[[date_col, driver_col]].drop_duplicates()
                    visits = visits.sort_values(date_col, ascending=False).head(20)
                    visits["Date"] = visits[date_col].dt.strftime("%Y-%m-%d")
                    visits["Salesman"] = visits[driver_col]
                    st.dataframe(visits[["Date", "Salesman"]], use_container_width=True, hide_index=True)
                else:
                    st.info("No salesman data available")

                        # ───────── Mini Tab 4: Issues (Returns + Material Details) ─────────
                        # ───────── Mini Tab 4: Issues (Returns + Material Details) ─────────
            with mini_tab4:
                issues = cust_sales[returns_mask]

                if not issues.empty:
                    # Main returns summary by invoice/date
                    st.error(f"**Returns: KD {returns_value:,.0f} ({return_rate:.2f}%)**")
                    st.dataframe(
                        issues[[date_col, "Billing Type", amount_col]].rename(
                            columns={date_col: "Billing Date", amount_col: "Return Value"}
                        ),
                        use_container_width=True
                    , hide_index=True)

                    # ---- Return Material Details table ----
                    st.markdown("#### Return Material Details")

                    # Prefer Material Description-like columns first
                    desc_candidates = [
                        "Material Description",
                        "Material Desc",
                        "Material Description EN",
                        "Material Description AR",
                        "MAT Description",
                    ]
                    desc_col = None
                    for c in desc_candidates:
                        if c in issues.columns:
                            desc_col = c
                            break

                    # Fallback: if no description column, use material_col (code)
                    if not desc_col and material_col and material_col in issues.columns:
                        desc_col = material_col

                    if desc_col:
                        issues_mat = issues.copy()
                        issues_mat[amount_col] = (
                            pd.to_numeric(issues_mat[amount_col], errors="coerce")
                            .fillna(0.0)
                            .abs()
                        )

                        # Group by the chosen description column
                        mat_summary = (
                            issues_mat
                            .groupby(desc_col)[amount_col]
                            .sum()
                            .reset_index()
                            .rename(columns={desc_col: "Material Description", amount_col: "Return Value"})
                            .sort_values("Return Value", ascending=False)
                        )

                        # Add Total row at end
                        total_val = mat_summary["Return Value"].sum()
                        total_row = {
                            "Material Description": "Total",
                            "Return Value": total_val
                        }
                        mat_summary = pd.concat(
                            [mat_summary, pd.DataFrame([total_row])],
                            ignore_index=True
                        )

                        st.dataframe(mat_summary, use_container_width=True, hide_index=True)
                    else:
                        st.info("No material description column found for returns.")
                else:
                    st.success("**No Returns – Perfect!**")

            # ───────── Mini Tab 5: Actions ─────────
            with mini_tab5:
                st.markdown("#### **Smart Actions**")
                actions = []
                if days_since > 30:
                    actions.append("**URGENT:** Schedule visit TODAY")
                if return_rate > 10:
                    actions.append("Call about quality issues")
                if total_sales > 5000:
                    actions.append("Offer premium products")
                if order_count > 15:
                    actions.append("Send loyalty reward")

                for a in actions:
                    st.markdown(f"• {a}")

                note_key = f"note_{selected_cust}"
                note = st.text_area("**Add Note**", value=st.session_state.get(note_key, ""), height=80)
                if st.button("**Save Note**", type="primary"):
                    st.session_state[note_key] = note
                    st.success("Note saved!")

                # Download Profile
                profile = pd.DataFrame({
                    "Metric": ["Customer", "Total Sales", "Orders", "Last Visit", "Days Since", "Return Rate %", "Note"],
                    "Value": [
                        selected_cust,
                        total_sales,
                        order_count,
                        last_visit.strftime("%Y-%m-%d") if pd.notna(last_visit) else "N/A",
                        days_since,
                        return_rate,
                        note
                    ]
                })
                st.download_button(
                    "**Download Profile (Excel)**",
                    data=to_excel_bytes(profile),
                    file_name=f"{selected_cust.replace(' ', '_')}_360.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

# --- Material Forecast Page ---
elif choice == texts[lang]["material_forecast"]:
    st.title(texts[lang]["material_forecast_title"])

    if "data_loaded" not in st.session_state:
        st.warning(texts[lang]["no_data_warning"])
        st.stop()

    # Use sales_df directly
    df_sales = st.session_state["sales_df"].copy()

    required_cols = ["Billing Date", "Material", "Material Description"]
    missing_cols = [col for col in required_cols if col not in df_sales.columns]
    if df_sales.empty or missing_cols:
        st.warning(f"⚠️ Sales data is missing required columns: {missing_cols}")
        st.stop()

    # Optional: if you have salesman/user view
    try:
        if "Driver Name EN" in df_sales.columns and "user_role" in st.session_state:
            # Keep existing logic (do not force filter here)
            pass
    except Exception:
        pass

    # Ensure date column is datetime
    df_sales["Billing Date"] = pd.to_datetime(df_sales["Billing Date"], errors="coerce")
    df_sales = df_sales.dropna(subset=["Billing Date"]).copy()

    # Ensure numeric columns
    if "Quantity" in df_sales.columns:
        df_sales["Quantity"] = pd.to_numeric(df_sales["Quantity"], errors="coerce").fillna(0)
    else:
        df_sales["Quantity"] = 0

    if "Net Value" in df_sales.columns:
        df_sales["Net Value"] = pd.to_numeric(df_sales["Net Value"], errors="coerce").fillna(0)

    # Extract Month & Year
    df_sales["Year"] = df_sales["Billing Date"].dt.year.astype(int)
    df_sales["Month"] = df_sales["Billing Date"].dt.month.astype(int)

    # ---------------- Settings ----------------
    with st.expander("⚙️ Forecast Settings", expanded=True):
        metric_choice = st.radio(
            "Forecast Based On",
            options=["Quantity", "Value (Net Value)"],
            horizontal=True,
            index=0
        )

        if metric_choice == "Value (Net Value)" and "Net Value" not in df_sales.columns:
            st.warning("⚠️ 'Net Value' column not found. Switching to Quantity.")
            metric_choice = "Quantity"

        value_col = "Quantity" if metric_choice == "Quantity" else "Net Value"

        # Materials
        all_mats = sorted(df_sales["Material Description"].dropna().astype(str).unique().tolist())
        if not all_mats:
            st.info("No materials found in the data.")
            st.stop()

        # Performance helper: Top-N default when too many
        use_topn = st.checkbox("Use Top-N Materials (recommended for large lists)", value=(len(all_mats) > 60))
        topn = st.slider("Top N Materials", 5, min(200, max(5, len(all_mats))), min(30, len(all_mats))) if use_topn else None

        if use_topn and topn:
            mat_rank = (
                df_sales.groupby("Material Description")[value_col].sum()
                .sort_values(ascending=False)
                .head(topn)
                .index.astype(str)
                .tolist()
            )
            default_mats = mat_rank
        else:
            # User requested: full materials when not selected
            default_mats = all_mats

        selected_mats = st.multiselect(
            "Select Materials (leave as default for all)",
            options=all_mats,
            default=default_mats
        )

        # If user clears selection, fall back to ALL (so nothing becomes empty)
        if not selected_mats:
            selected_mats = all_mats

        exclude_returns = st.checkbox("Exclude Returns (YKRE / ZRE)", value=False)
        exclude_cancels = st.checkbox("Exclude Cancellations (YKS1 / YKS2 / ZCAN)", value=False)

    # Apply optional exclusions
    df_work = df_sales.copy()
    if exclude_returns and "Billing Type" in df_work.columns:
        df_work = df_work[~df_work["Billing Type"].astype(str).str.upper().isin(["YKRE", "ZRE"])].copy()
    if exclude_cancels and "Billing Type" in df_work.columns:
        df_work = df_work[~df_work["Billing Type"].astype(str).str.upper().isin(["YKS1", "YKS2", "ZCAN"])].copy()

    df_work = df_work[df_work["Material Description"].astype(str).isin([str(x) for x in selected_mats])].copy()

    # Tabs for Monthly & Yearly Forecast
    tab_month, tab_year = st.tabs(["Monthly Forecast", "Yearly Forecast"])

    # ---------------- Monthly Forecast ----------------
    with tab_month:
        st.subheader("Monthly Material Forecast")

        years = sorted(df_work["Year"].dropna().unique().tolist())
        if not years:
            st.info("No valid years found after filters.")
            st.stop()

        selected_year = st.selectbox("Select Year:", years, index=len(years)-1)
        df_monthly = df_work[df_work["Year"] == selected_year].copy()

        monthly = (
            df_monthly.groupby(["Month", "Material Description"])[value_col]
            .sum()
            .reset_index()
        )

        # Fill missing months for each material
        all_months = pd.DataFrame({"Month": list(range(1, 13))})
        all_materials = pd.DataFrame({"Material Description": sorted(df_monthly["Material Description"].dropna().astype(str).unique().tolist())})
        full_index = all_materials.merge(all_months, how="cross")

        monthly = full_index.merge(
            monthly, on=["Month", "Material Description"], how="left"
        ).fillna({value_col: 0})

        # Plot
        fig = px.line(
            monthly,
            x="Month",
            y=value_col,
            color="Material Description",
            markers=True,
            title=f"Monthly Trend ({selected_year}) – {metric_choice}"
        )
        st.plotly_chart(fig, use_container_width=True)

        # Pivot
        pivot_table = monthly.pivot(
            index="Material Description", columns="Month", values=value_col
        ).fillna(0)

        st.dataframe(pivot_table, hide_index=True)

        # Download Excel
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        excel_bytes = to_excel_bytes(monthly, sheet_name="Monthly_Forecast")
        if st.download_button(
            texts[lang].get("download_excel", "⬇️ Download Excel"),
            data=excel_bytes,
            file_name=f"monthly_forecast_{selected_year}_{timestamp}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            st.session_state["audit_log"].append({
                "user": username,
                "action": "download",
                "details": f"monthly_forecast_{selected_year}_{timestamp}.xlsx",
                "timestamp": datetime.now().isoformat()
            })

    # ---------------- Yearly Forecast ----------------
    with tab_year:
        st.subheader("Yearly Material Forecast")

        years = sorted(df_work["Year"].dropna().unique().tolist())
        if not years:
            st.info("No valid years found after filters.")
            st.stop()

        # Let user pick years to compare
        default_years = years[-3:] if len(years) >= 3 else years
        selected_years = st.multiselect("Select Year(s):", options=years, default=default_years)
        if not selected_years:
            selected_years = years

        df_year = df_work[df_work["Year"].isin(selected_years)].copy()

        yearly = (
            df_year.groupby(["Year", "Material Description"])[value_col]
            .sum()
            .reset_index()
        )

        fig = px.bar(
            yearly,
            x="Year",
            y=value_col,
            color="Material Description",
            barmode="group",
            text=value_col,
            title=f"Yearly Trend – {metric_choice}"
        )
        st.plotly_chart(fig, use_container_width=True)

        pivot_table_year = yearly.pivot(
            index="Material Description", columns="Year", values=value_col
        ).fillna(0)

        st.dataframe(pivot_table_year, hide_index=True)

        # Download Excel
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        excel_bytes_year = to_excel_bytes(yearly, sheet_name="Yearly_Forecast")
        if st.download_button(
            texts[lang].get("download_excel", "⬇️ Download Excel"),
            data=excel_bytes_year,
            file_name=f"yearly_forecast_{timestamp}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            st.session_state["audit_log"].append({
                "user": username,
                "action": "download",
                "details": f"yearly_forecast_{timestamp}.xlsx",
                "timestamp": datetime.now().isoformat()
            })
                      

# ================= PROFIT & MARGIN PAGE =================
elif choice == "💰 Profit & Margin":
    st.title("💰 Profit & Margin Analysis")

    # ─── Helper: fuzzy column finder ──────────────────────────────────────
    def find_column(df, possible_names):
        if df.empty:
            return None
        possible = [str(n).lower().replace(" ", "").replace("_", "") for n in possible_names]
        for col in df.columns:
            clean = str(col).lower().replace(" ", "").replace("_", "")
            if clean in possible or any(p in clean for p in possible):
                return col
        return None

    # ─── Price list status ─────────────────────────────────────────────────
    price_df = st.session_state.get("price_df", pd.DataFrame())
    if price_df.empty:
        st.info("Price list sheet ('price list') not found → cost & discount will show as missing.")
    else:
        st.success(f"Price list loaded ({len(price_df):,} rows)")
        with st.expander("Price List Preview (first 8 rows)", expanded=False):
            st.dataframe(price_df.head(8))

    # ─── Data source selection ─────────────────────────────────────────────
    data_source = st.radio(
        "Analyze which data?",
        ["Current Month only (sales data sheet)", "Full Year / Historical (YTD sheet)"],
        index=0
    )

    if data_source == "Current Month only (sales data sheet)":
        if "sales_df" not in st.session_state or st.session_state["sales_df"].empty:
            st.warning("No current month sales data loaded.")
            st.stop()
        base_df = st.session_state["sales_df"].copy()
    else:
        if "ytd_df" not in st.session_state or st.session_state["ytd_df"].empty:
            st.warning("No YTD data → falling back to current month sales.")
            if "sales_df" not in st.session_state or st.session_state["sales_df"].empty:
                st.stop()
            base_df = st.session_state["sales_df"].copy()
        else:
            base_df = st.session_state["ytd_df"].copy()

    # ─── Column discovery ──────────────────────────────────────────────────
    MATERIAL_COL = find_column(base_df, ["Material Description", "Mat Description", "Description", "Item Description"])
    QTY_COL      = find_column(base_df, ["Quantity", "Qty", "Sales Qty"])
    UOM_COL      = find_column(base_df, ["UOM", "Unit of Measure", "Unit"])
    NET_COL      = find_column(base_df, ["Net Value", "Net Amount", "Net Sales", "Amount"])

    PRICE_MAT_COL = find_column(price_df, ["Material Description", "Mat Description", "Description"])
    COST_COL      = find_column(price_df, ["Cost Price", "Cost", "Unit Cost"])
    PACK_COL      = find_column(price_df, ["Pack Size", "Pack", "Pack Qty"])
    CATEGORY_COL  = find_column(price_df, [
        "Category", "Item Category", "Sales Category", "Material Category",
        "Type", "Group", "Product Category", "Exchange Category",
        "Transaction Type", "Doc Type", "Bill Type", "Sale Type", "Doc Category"
    ])

    missing = []
    if MATERIAL_COL is None: missing.append("Material Description (main data)")
    if QTY_COL      is None: missing.append("Quantity")
    if NET_COL      is None: missing.append("Net Value / Net Amount")
    if not price_df.empty:
        if PRICE_MAT_COL is None: missing.append("Material Description (price list)")
        if COST_COL      is None: missing.append("Cost Price")

    if missing:
        st.error("Cannot calculate Profit & Margin — missing critical columns:\n• " + "\n• ".join(missing))
        st.stop()

    # ─── Date range filter ─────────────────────────────────────────────────
    col1, col2 = st.columns(2)
    with col1:
        min_date = base_df["Billing Date"].min().date() if "Billing Date" in base_df.columns else datetime.date.today()
        start_date = st.date_input("Start Date", value=min_date)
    with col2:
        max_date = base_df["Billing Date"].max().date() if "Billing Date" in base_df.columns else datetime.date.today()
        end_date = st.date_input("End Date", value=max_date)

    start_dt = pd.to_datetime(start_date)
    end_dt   = pd.to_datetime(end_date) + pd.Timedelta(days=1)

    df_pm = base_df[
        (base_df["Billing Date"] >= start_dt) &
        (base_df["Billing Date"] < end_dt)
    ].copy()

    if df_pm.empty:
        st.warning("No records found in selected date range.")
        st.stop()

    # ─── Prepare normalized matching ───────────────────────────────────────
    df_val = df_pm.copy()
    df_val["_mat_norm"] = df_val[MATERIAL_COL].astype(str).str.strip().str.upper()

    df_val["Cost Price"]  = pd.NA
    df_val["Pack Size"]   = pd.NA
    df_val["Category"]    = pd.NA
# ─── Customer normalization (for R&R / contracts) ─────────────────────
    if "PY Name 1" in df_val.columns:
        df_val["_py_name_norm"] = (
            df_val["PY Name 1"]
            .astype(str)
            .str.strip()
            .str.lower()
        )
    else:
        df_val["_py_name_norm"] = ""

    # ─── Price list mapping ────────────────────────────────────────────────
    if not price_df.empty and PRICE_MAT_COL:
        price_df["_mat_norm"] = price_df[PRICE_MAT_COL].astype(str).str.strip().str.upper()
        price_map = price_df.set_index("_mat_norm")

        df_val["Cost Price"] = df_val["_mat_norm"].map(price_map.get(COST_COL))
        if PACK_COL:
            df_val["Pack Size"] = df_val["_mat_norm"].map(price_map.get(PACK_COL))
        if CATEGORY_COL:
            df_val["Category"] = df_val["_mat_norm"].map(price_map.get(CATEGORY_COL))

    # ─── Cost calculation ──────────────────────────────────────────────────
    def calculate_line_cost(row):
        if pd.isna(row["Cost Price"]):
            return None
        qty   = pd.to_numeric(row.get(QTY_COL), errors='coerce') or 0
        cost  = pd.to_numeric(row["Cost Price"], errors='coerce')
        uom   = str(row.get(UOM_COL, "")).strip().upper()
        pack  = pd.to_numeric(row.get("Pack Size"), errors='coerce')

        if uom == "KAR":
            if pd.isna(pack) or pack <= 0:
                return None
            return qty * pack * cost
        return qty * cost

    df_val["Total Cost"] = df_val.apply(calculate_line_cost, axis=1)

    # ─── Discount calculation ──────────────────────────────────────────────
    df_val[NET_COL] = pd.to_numeric(df_val[NET_COL], errors='coerce').fillna(0)

    # ✅ FIX #1: Return/Cancel rows cost reversal (prevents wrong Discount Value)
    # If Net Value is negative but Total Cost is positive, flip cost to negative (cost reversal)
    df_val.loc[
        (df_val[NET_COL] < 0) &
        (df_val["Total Cost"].notna()) &
        (df_val["Total Cost"] > 0),
        "Total Cost"
    ] *= -1

    # Now Discount Value becomes correct for both Sales and Returns
    df_val["Discount Value"] = df_val[NET_COL] - df_val["Total Cost"].fillna(0)

    # ✅ FIX #2: Discount % only for SALES (Net >= 0) and based on Cost
    cost_safe = df_val["Total Cost"].replace(0, np.nan)
    df_val["Discount %"] = np.where(
        df_val[NET_COL] >= 0,
        (abs(df_val["Discount Value"]) / cost_safe * 100),
        0
    )
    df_val["Discount %"] = pd.to_numeric(df_val["Discount %"], errors="coerce").fillna(0).round(2)

    df_val["⚠ Cost Missing"] = df_val["Cost Price"].isna()

    if UOM_COL in df_val.columns:
        df_val["⚠ Pack Missing (KAR)"] = (
            df_val[UOM_COL].astype(str).str.upper() == "KAR"
        ) & df_val["Pack Size"].isna()
    else:
        df_val["⚠ Pack Missing (KAR)"] = False

    # ============================================================
    # ✅ NEW: R&R (Rebate % + Display Rental value) mapping & logic
    # ============================================================
    rr_df = st.session_state.get("rr_df", pd.DataFrame())

    RR_CUST_COL   = find_column(rr_df, ["PY Name 1", "Customer", "Customer Name"])
    REBATE_COL    = find_column(rr_df, ["Rebate %", "Rebate", "Rebate Percent", "Rebate%"])
    RENTAL_COL    = find_column(rr_df, ["Display Rental value", "Display Rental", "Rental", "Annual Rental", "Rental Value"])

    # defaults (so page never breaks)
    df_val["Rebate %"] = 0.0
    df_val["Display Rental value"] = 0.0

    if not rr_df.empty and RR_CUST_COL and (REBATE_COL or RENTAL_COL) and ("PY Name 1" in df_val.columns):
        rr_tmp = rr_df.copy()
        rr_tmp["_py_name_norm"] = rr_tmp[RR_CUST_COL].astype(str).str.strip().str.upper()

        df_val["_py_name_norm"] = df_val["PY Name 1"].astype(str).str.strip().str.upper()

        if REBATE_COL:
            rr_tmp["Rebate %"] = pd.to_numeric(rr_tmp[REBATE_COL], errors="coerce").fillna(0)
        else:
            rr_tmp["Rebate %"] = 0.0

        if RENTAL_COL:
            rr_tmp["Display Rental value"] = pd.to_numeric(rr_tmp[RENTAL_COL], errors="coerce").fillna(0)
        else:
            rr_tmp["Display Rental value"] = 0.0

        rr_map = rr_tmp.set_index("_py_name_norm")[["Rebate %", "Display Rental value"]]

        df_val["Rebate %"] = df_val["_py_name_norm"].map(rr_map["Rebate %"]).fillna(0.0)
        df_val["Display Rental value"] = df_val["_py_name_norm"].map(rr_map["Display Rental value"]).fillna(0.0)
    else:
        # keep defaults; optional info (no stop)
        pass

    # Rebate Value applies only on positive sales (avoid confusion on returns)
    # ============================================================
    # ✅ CONTRACT-COMPLIANT REBATE CALCULATION (NET VALUE BASED)
    # ============================================================

    # Net Value must remain untouched (contract rule)
    df_val[NET_COL] = pd.to_numeric(df_val[NET_COL], errors="coerce").fillna(0)

    # Rebate % numeric safety
    df_val["Rebate %"] = pd.to_numeric(
        df_val["Rebate %"],
        errors="coerce"
    ).fillna(0)
    # --- R&R mapping done above ---

    

    # ============================================================
    # ✅ REBATE – STRICTLY BASED ON NET VALUE (CONTRACT RULE)
    # ============================================================

    # Safety
    df_val[NET_COL] = pd.to_numeric(df_val[NET_COL], errors="coerce").fillna(0)
    df_val["Rebate %"] = pd.to_numeric(df_val["Rebate %"], errors="coerce").fillna(0)

    # 1️⃣ Customer-level NET VALUE (sales + returns)
    cust_net = (
        df_val
        .groupby("_py_name_norm")[NET_COL]
        .sum()
    )

    # 2️⃣ Customer-level REBATE (only on net)
    cust_rebate = (
        cust_net *
        (df_val.groupby("_py_name_norm")["Rebate %"].first() / 100)
    ).fillna(0)

    # 3️⃣ Map back to dataframe (same value per customer)
    df_val["Rebate Value"] = df_val["_py_name_norm"].map(cust_rebate).fillna(0)


    

    # Audit helper flag (optional but useful)
    df_val["Rebate Applied"] = df_val[NET_COL] > 0


    # Allocate annual rental to the selected period (by days in range), then distribute by sales share per customer
    period_days = max((pd.to_datetime(end_date) - pd.to_datetime(start_date)).days + 1, 1)
    year_factor = period_days / 365.0

    # annual rental for period
    df_val["_period_rental"] = df_val["Display Rental value"] * year_factor

    # sales share per customer (only positive net)
    pos_sales = df_val[df_val[NET_COL] > 0].groupby("_py_name_norm")[NET_COL].sum()
    df_val["_cust_pos_sales"] = df_val["_py_name_norm"].map(pos_sales)
    df_val["_sales_share"] = np.where(
        (df_val[NET_COL] > 0) & (df_val["_cust_pos_sales"].notna()) & (df_val["_cust_pos_sales"] != 0),
        df_val[NET_COL] / df_val["_cust_pos_sales"],
        0.0
    )

    df_val["Allocated Rental"] = (df_val["_period_rental"] * df_val["_sales_share"]).fillna(0.0)

    # Effective profit after Discount + Rebate + Rental
    df_val["Effective Profit"] = (
        df_val[NET_COL]
        - df_val["Total Cost"].fillna(0)
        - df_val["Rebate Value"].fillna(0)
        - df_val["Allocated Rental"].fillna(0)
    )

    df_val["Effective Margin %"] = np.where(
        df_val[NET_COL] != 0,
        (df_val["Effective Profit"] / df_val[NET_COL]) * 100,
        0.0
    )
    df_val["Effective Margin %"] = pd.to_numeric(df_val["Effective Margin %"], errors="coerce").fillna(0).round(2)

    # clean helper cols used for allocation
    for _c in ["_period_rental", "_cust_pos_sales", "_sales_share"]:
        if _c in df_val.columns:
            pass  # keep internal (not shown), safe

    # ─── Category logic ────────────────────────────────────────────────────
    sales_category_value  = 0.0
    return_category_value = 0.0

    has_category = "Category" in df_val.columns and df_val["Category"].notna().any()

    if has_category:
        cat_clean = df_val["Category"].astype(str).str.lower().str.strip()
        pos_mask = cat_clean.str.contains(
            r"sale|revenue|positive|income|normal|inv|bill|pos|good|regular",
            na=False, regex=True
        )
        ret_mask = cat_clean.str.contains(
            r"return|ret|refund|negative|adjustment|cr|debit|reverse|cancel|cn|zre|ykre|credit|deduction",
            na=False, regex=True
        )
        sales_category_value  = df_val[pos_mask][NET_COL].sum()
        return_category_value = df_val[ret_mask][NET_COL].sum()
    else:
        sales_category_value  = df_val[df_val[NET_COL] >= 0][NET_COL].sum()
        return_category_value = df_val[df_val[NET_COL] <  0][NET_COL].sum()

    # ─── Executive KPIs ────────────────────────────────────────────────────
    st.markdown("## 📌 Executive Summary")

    total_net      = df_val[NET_COL].sum()
    total_cost     = df_val["Total Cost"].sum(min_count=1)
    total_discount = df_val["Discount Value"].sum(min_count=1)

    # ✅ (Keep same KPI style) but make Discount % based on Total Cost safely
    overall_disc_pct = (total_discount / total_cost * 100) if total_cost != 0 else 0

    cols = st.columns(7)
    cols[0].metric("Total Net",           f"{total_net:,.2f}")
    cols[1].metric("Total Cost",          f"{total_cost:,.2f}")
    cols[2].metric("Discount Value",      f"{total_discount:,.2f}")
    cols[3].metric("Discount %",          f"{overall_disc_pct:.2f}%")
    cols[4].metric("Cost Missing Rows",   f"{df_val['⚠ Cost Missing'].sum():,}")
    cols[5].metric("Positive / Sales Val", f"{sales_category_value:,.2f}")
    cols[6].metric("Returns / Negative",   f"{return_category_value:,.2f}")

    # ✅ NEW: R&R Executive Add-on (keeps your old KPIs unchanged)
    st.markdown("## 🧾 Contract (R&R) Impact")

    total_rebate = df_val["Rebate Value"].sum(min_count=1)
    total_rental = df_val["Allocated Rental"].sum(min_count=1)
    total_eff_profit = df_val["Effective Profit"].sum(min_count=1)

    cA, cB, cC, cD = st.columns(4)
    cA.metric("Rebate Value", f"{total_rebate:,.2f}")
    cB.metric("Allocated Rental (Period)", f"{total_rental:,.2f}")
    cC.metric("Effective Profit (after R&R)", f"{total_eff_profit:,.2f}")
    # effective margin (sales only)
    sales_only_net = df_val[df_val[NET_COL] > 0][NET_COL].sum()
    sales_only_eff_profit = df_val[df_val[NET_COL] > 0]["Effective Profit"].sum()
    eff_margin_pct = (sales_only_eff_profit / sales_only_net * 100) if sales_only_net else 0
    cD.metric("Effective Margin % (Sales)", f"{eff_margin_pct:.2f}%")

    # ============================================================
    # 📊 EXECUTIVE SUMMARY – CATEGORY PERFORMANCE (FINAL)
    # ============================================================

    st.markdown("## 📊 Executive Summary – Category Performance")

    if "Category" not in df_val.columns or df_val["Category"].isna().all():
        st.warning("⚠️ Category not mapped — check price list material matching")
    else:
        category_summary = (
            df_val
            .groupby("Category", dropna=False)[NET_COL]
            .sum()
            .reset_index()
            .rename(columns={NET_COL: "Net Value"})
            .sort_values("Net Value", ascending=False)
        )

        st.markdown("### 🧾 Category-wise Net Value")

        kpi_per_row = 4
        data = category_summary.to_dict("records")

        for i in range(0, len(data), kpi_per_row):
            cols = st.columns(len(data[i:i + kpi_per_row]))
            for col, row in zip(cols, data[i:i + kpi_per_row]):
                with col:
                    st.metric(
                        label=str(row["Category"]),
                        value=f"KD {row['Net Value']:,.0f}"
                    )

        st.markdown("### 📌 Category Contribution")

        category_contribution = (
            df_val
            .groupby("Category", dropna=False)[NET_COL]
            .sum()
            .reset_index()
            .rename(columns={NET_COL: "Net Value"})
        )

        total_net_value = category_contribution["Net Value"].sum()

        category_contribution["Contribution %"] = (
            category_contribution["Net Value"] / total_net_value * 100
        ).round(1)

        insight_lines = [
            f"• **{row['Category']}** → {row['Contribution %']}%"
            for _, row in category_contribution
                .sort_values("Contribution %", ascending=False)
                .iterrows()
        ]

        st.markdown("  \n".join(insight_lines))

    # ─── Data Quality ──────────────────────────────────────────────────────
    st.markdown("### 🧪 Data Quality")
    dq1, dq2, dq3 = st.columns(3)
    dq1.metric("Cost Mapped %", f"{100 - df_val['⚠ Cost Missing'].mean()*100:.1f}%")
    dq2.metric("KAR Pack Mapped %", f"{100 - df_val.get('⚠ Pack Missing (KAR)', pd.Series(0)).mean()*100:.1f}%")
    dq3.metric("Negative Discount Rows", int((df_val["Discount Value"] < 0).sum()))

    # ─── Material level summary ────────────────────────────────────────────
    st.markdown("### 🔍 Material Discount Hotspots")

    agg_dict = {}
    if NET_COL in df_val.columns:
        agg_dict["Net_Value"] = (NET_COL, "sum")
    if "Discount %" in df_val.columns:
        agg_dict["Avg_Discount_Pct"] = ("Discount %", "mean")
    if "Discount Value" in df_val.columns:
        agg_dict["Discount_Value"] = ("Discount Value", "sum")
    if "Total Cost" in df_val.columns:
        agg_dict["Total_Cost"] = ("Total Cost", "sum")
    # ✅ NEW (optional) include effective margin
    if "Effective Profit" in df_val.columns:
        agg_dict["Effective_Profit"] = ("Effective Profit", "sum")

    if agg_dict:
        material_summary = (
            df_val.groupby(MATERIAL_COL, dropna=True)
                  .agg(**agg_dict)
                  .reset_index()
        )

        if "Discount_Value" in material_summary.columns:
            cA, cB = st.columns(2)
            with cA:
                st.caption("🔴 Top biggest discount / loss makers")
                st.dataframe(
                    material_summary.sort_values("Discount_Value").head(12),
                    use_container_width=True
                )
            with cB:
                st.caption("🟢 Top most profitable (lowest discount)")
                st.dataframe(
                    material_summary.sort_values("Discount_Value", ascending=False).head(12),
                    use_container_width=True
                )
        else:
            st.dataframe(material_summary, use_container_width=True)
    else:
        st.info("No numeric columns available for material summary.")

    # ─── Quick Filters ─────────────────────────────────────────────────────
    st.markdown("### 🎯 Quick Filters")
    f1, f2, f3 = st.columns(3)
    show_negative     = f1.checkbox("Only negative discount", key="pm_neg")
    show_cost_missing = f2.checkbox("Only cost missing",      key="pm_cost")
    show_kar_issues   = f3.checkbox("Only KAR pack issues",   key="pm_kar")

    display_df = df_val.copy()
    if show_negative:
        display_df = display_df[display_df["Discount Value"] < 0]
    if show_cost_missing:
        display_df = display_df[display_df["⚠ Cost Missing"]]
    if show_kar_issues:
        display_df = display_df[display_df.get("⚠ Pack Missing (KAR)", False)]

    # ─── Final table ───────────────────────────────────────────────────────
    st.subheader(f"Result: {len(display_df):,} rows  •  {start_date} → {end_date}")

    final_cols = [
        "Billing Date", "Driver Name EN", "PY Name 1",
        MATERIAL_COL, QTY_COL, UOM_COL,
        "Cost Price", "Pack Size", "Total Cost",
        NET_COL, "Discount Value", "Discount %",
        # ✅ NEW: R&R columns
        "Rebate %", "Rebate Value", "Display Rental value", "Allocated Rental",
        "Effective Profit", "Effective Margin %",
        "⚠ Cost Missing", "⚠ Pack Missing (KAR)"
    ]
    avail_cols = [c for c in final_cols if c in display_df.columns]

    st.dataframe(
        display_df[avail_cols].sort_values("Billing Date", ascending=False),
        use_container_width=True,
        hide_index=True
    )

    # ✅ NEW: Customer contract compliance summary (R&R + Discount)
    st.markdown("## 📜 Customer Contract Compliance (Discount + Rebate + Rental)")

    if "PY Name 1" in df_val.columns:
        cust_sum = (
            df_val.groupby("PY Name 1")
                  .agg(
                      Net_Sales=(NET_COL, "sum"),
                      Total_Cost=("Total Cost", "sum"),
                      Discount_Value=("Discount Value", "sum"),
                      Rebate_Value=("Rebate Value", "max"),
                      Rental_Allocated=("Allocated Rental", "sum"),
                      Effective_Profit=("Effective Profit", "sum"),
                  )
                  .reset_index()
        )

        # Leakage % on positive sales only (safer)
        cust_pos_sales = df_val[df_val[NET_COL] > 0].groupby("PY Name 1")[NET_COL].sum()
        cust_sum["_pos_sales"] = cust_sum["PY Name 1"].map(cust_pos_sales).fillna(0)

        cust_sum["Total Leakage"] = (
            cust_sum["Discount_Value"].abs()
            + cust_sum["Rebate_Value"].abs()
            + cust_sum["Rental_Allocated"].abs()
        )

        cust_sum["Leakage % (on Sales)"] = np.where(
            cust_sum["_pos_sales"] != 0,
            (cust_sum["Total Leakage"] / cust_sum["_pos_sales"]) * 100,
            0
        ).round(2)

        cust_sum["Effective Margin %"] = np.where(
            cust_sum["Net_Sales"] != 0,
            (cust_sum["Effective_Profit"] / cust_sum["Net_Sales"]) * 100,
            0
        ).round(2)

        show_cols = [
            "PY Name 1", "Net_Sales", "Total_Cost",
            "Discount_Value", "Rebate_Value", "Rental_Allocated",
            "Total Leakage", "Leakage % (on Sales)",
            "Effective_Profit", "Effective Margin %"
        ]

        st.dataframe(
            cust_sum[show_cols].sort_values("Leakage % (on Sales)", ascending=False),
            use_container_width=True,
            hide_index=True
        )
    else:
        st.info("Customer column 'PY Name 1' not found for contract summary.")

    # ============================================================
    # 📊 Category Contribution % (ONE-LINE INSIGHT)
    # ============================================================

    st.markdown("### 📊 Category Contribution")

    df_val["Category"] = (
        df_val["Category"]
        .fillna("Unmapped")
        .astype(str)
        .str.strip()
    )

    category_contribution = (
        df_val
        .groupby("Category")[NET_COL]
        .sum()
        .reset_index()
        .rename(columns={NET_COL: "Net Value"})
    )

    total_net_value = category_contribution["Net Value"].sum()

    for _, row in category_contribution.sort_values("Net Value", ascending=False).iterrows():
        pct = (row["Net Value"] / total_net_value * 100) if total_net_value else 0
        st.markdown(f"• **{row['Category']}** → {pct:.1f}%")

    # ============================================================
    # 💰 Profit View – Margin by Category
    # ============================================================

    st.markdown("## 💰 Profit View – Margin by Category")

    margin_by_category = (
        df_val
        .groupby("Category")
        .agg(
            Net_Value=(NET_COL, "sum"),
            Total_Cost=("Total Cost", "sum"),
        )
        .reset_index()
    )

    margin_by_category["Gross Profit"] = (
        margin_by_category["Net_Value"] - margin_by_category["Total_Cost"]
    )

    margin_by_category["Margin %"] = (
        margin_by_category["Gross Profit"] /
        margin_by_category["Net_Value"].replace(0, np.nan) * 100
    ).round(1)

    st.dataframe(
        margin_by_category.sort_values("Net_Value", ascending=False),
        use_container_width=True,
        hide_index=True
    )

    # ============================================================
    # 🏬 Customer × Category Sales Matrix
    # ============================================================

    st.markdown("## 🏬 Customer-wise Sales by Category")

    CUSTOMER_COL = "PY Name 1"

    customer_category = (
        df_val
        .pivot_table(
            index=CUSTOMER_COL,
            columns="Category",
            values=NET_COL,
            aggfunc="sum",
            fill_value=0
        )
    )

    customer_category["Total"] = customer_category.sum(axis=1)
    customer_category = customer_category.sort_values("Total", ascending=False)

    total_row = customer_category.sum().to_frame().T
    total_row.index = ["TOTAL"]

    customer_category_final = pd.concat([customer_category, total_row])

    st.dataframe(
        customer_category_final,
        use_container_width=True
    )

    # ─── Downloads ─────────────────────────────────────────────────────────
    st.markdown("### ⬇️ Export")
    c1, c2 = st.columns(2)
    with c1:
        csv = display_df.to_csv(index=False).encode("utf-8-sig")
        st.download_button("⬇️ CSV", csv, f"profit_margin_{start_date}_to_{end_date}.csv", "text/csv")
    with c2:
        st.download_button("⬇️ Excel", data=to_excel_bytes(display_df),
                           file_name=f"profit_margin_{start_date}_to_{end_date}.xlsx")
        
        


from datetime import date

def build_daily_email_summary(
    total_target,
    total_sales,
    salesman_df,
    customer_sales
):
    today_str = date.today().strftime("%d %b %Y")
    balance = total_target - total_sales

    subject = f"Daily Sales Summary – {today_str}"

    body = f"""
📊 OVERALL PERFORMANCE
Total Target : {total_target:,.0f} KD
Achieved     : {total_sales:,.0f} KD
Balance      : {balance:,.0f} KD

👤 SALESMAN PERFORMANCE
--------------------------------------------------
Salesman        Target      Achieved     Balance
--------------------------------------------------
"""

    for _, r in salesman_df.iterrows():
        bal = r["Target"] - r["Achieved"]
        body += (
            f"{r['Driver Name EN']:<15}"
            f"{r['Target']:>12,.0f}"
            f"{r['Achieved']:>14,.0f}"
            f"{bal:>14,.0f}\n"
        )

    body += """
--------------------------------------------------

🏪 CUSTOMER SALES SUMMARY
--------------------------------------------------
Customer              Sales (KD)
--------------------------------------------------
"""

    for cust, val in customer_sales.items():
        body += f"{cust:<22}{val:>12,.0f}\n"

    body += """
--------------------------------------------------

Regards,
Sales Dashboard
"""

    return subject, body


if choice == "🧭 Management Command Center":

    st.title("🧭 Management Command Center")

    # ================= SAFETY CHECK =================
    if sales_df is None or sales_df.empty:
        st.warning("Please load sales data first")
        st.stop()

    df = sales_df.copy()
    df["Billing Date"] = pd.to_datetime(df["Billing Date"], errors="coerce")

    # ================= DATE & WORKING DAYS =================
    today = pd.to_datetime("today").normalize()
    month_start = today.replace(day=1)
    month_end = month_start + pd.offsets.MonthEnd(1)

    all_days = pd.date_range(month_start, month_end, freq="D")
    working_days = all_days[all_days.weekday != 4]  # Exclude Friday only

    total_working_days = len(working_days)
    days_completed = max(1, len(working_days[working_days <= today]))

    # ================= TARGET DATA =================
    if "target_df" in globals() and "KA Target" in target_df.columns:
        ka_target_map = target_df.set_index("Driver Name EN")["KA Target"]
    else:
        ka_target_map = pd.Series(dtype=float)

    # ================= OVERALL SALES =================
    total_sales = float(df["Net Value"].sum())
    total_ka_target = float(ka_target_map.sum()) if not ka_target_map.empty else 0.0

    # ================= DAILY PACE =================
    ka_target_per_day = round(
        total_ka_target / total_working_days, 0
    ) if total_working_days > 0 else 0

    ka_actual_per_day = round(
        total_sales / days_completed, 0
    )

    def pace_status(actual_day, target_day):
        if target_day <= 0:
            return "🟢 GREEN"
        ratio = actual_day / target_day
        if ratio >= 1.0:
            return "🟢 GREEN"
        elif ratio >= 0.95:
            return "🟠 AMBER"
        else:
            return "🔴 RED"

    overall_ka_status = pace_status(ka_actual_per_day, ka_target_per_day)

    # ================= 1️⃣ EXECUTIVE DASHBOARD =================
    st.subheader("1️⃣ Executive RAG Dashboard (Daily Pace)")

    # ─── Get Date Range From Current Data ─────────────────────
    DATE_COL = "Billing Date"   # change only if column name different

    df_dates = df.copy()
    df_dates[DATE_COL] = pd.to_datetime(df_dates[DATE_COL], errors="coerce")

    from_dt = df_dates[DATE_COL].min()
    to_dt   = df_dates[DATE_COL].max()

    from_txt = from_dt.strftime("%d %b %Y") if pd.notna(from_dt) else "-"
    to_txt   = to_dt.strftime("%d %b %Y") if pd.notna(to_dt) else "-"

    # ─── KPI Columns ───────────────────────────────────────────
    c1, c2, c3, c4, c5 = st.columns(5)

    # ---- Total Sales KPI
    c1.metric("Total Sales", f"KD {total_sales:,.0f}")

    # ---- Date range below KPI (Left = From, Right = To)
    c1.markdown(
        f"""
        <div style="display:flex;justify-content:space-between;
                    font-size:12px;margin-top:-8px;color:#334155;">
            <span><b>From:</b> {from_txt}</span>
            <span><b>To:</b> {to_txt}</span>
        </div>
        """,
        unsafe_allow_html=True
    )

    # ---- Other KPIs
    c2.metric("Total KA Target", f"KD {total_ka_target:,.0f}")
    c3.metric("KA Target / Day", f"KD {ka_target_per_day:,.0f}")
    c4.metric("KA Actual / Day", f"KD {ka_actual_per_day:,.0f}")
    c5.metric("Overall KA Status", overall_ka_status)
    # ================= 2️⃣ SALESMAN TABLE =================
    st.subheader("2️⃣ Salesman Performance")

    salesman_df = (
        df.groupby("Driver Name EN")["Net Value"]
        .sum()
        .reset_index(name="Achieved")
    )

    salesman_df["Target"] = salesman_df["Driver Name EN"].map(ka_target_map).fillna(0)
    salesman_df["Balance"] = salesman_df["Target"] - salesman_df["Achieved"]

    st.dataframe(
        salesman_df[["Driver Name EN", "Target", "Achieved", "Balance"]],
        use_container_width=True
    )

    # ================= 3️⃣ MANAGEMENT INSIGHTS =================
    st.subheader("3️⃣ Action-based Management Insights")

    st.write(
        "🟢 Overall KA pace ON TRACK"
        if overall_ka_status == "🟢 GREEN"
        else "🟠 Overall KA pace NEEDS PUSH"
        if overall_ka_status == "🟠 AMBER"
        else "🚨 Overall KA pace CRITICAL"
    )

    # ================= 4️⃣ EMAIL SUMMARY =================
    st.subheader("📧 Daily Email Summary")

    # ---- Customer sales (Top 10) ----
    customer_sales = (
        df.groupby("PY Name 1")["Net Value"]
        .sum()
        .sort_values(ascending=False)
        .head(10)
    )

    subject, body = build_daily_email_summary(
        total_ka_target,
        total_sales,
        salesman_df,
        customer_sales
    )

    st.text_area(
        "📄 Email Preview",
        value=f"Subject: {subject}\n\n{body}",
        height=420
    )

    mailto_link = (
        f"mailto:?subject={urllib.parse.quote(subject)}"
        f"&body={urllib.parse.quote(body)}"
    )

    st.markdown(
        f"""
        <a href="{mailto_link}">
            <button style="
                background-color:#2563eb;
                color:white;
                padding:10px 18px;
                border:none;
                border-radius:6px;
                font-size:16px;
                cursor:pointer;">
                📧 Send Daily Summary Email
            </button>
        </a>
        """,
        unsafe_allow_html=True
    )


# Admin-only Audit Logs View
if user_role == "admin":
    st.sidebar.markdown("---")
    st.sidebar.subheader("Admin Tools")
    if st.sidebar.button("View Audit Logs"):
        st.title("📋 Audit Logs")
        log_df = pd.DataFrame(st.session_state["audit_log"])
        st.dataframe(log_df, hide_index=True)
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        if st.download_button(
            "⬇️ Download Audit Logs (Excel)",
            data=to_excel_bytes(log_df, sheet_name="Audit_Logs", index=False),
            file_name=f"audit_logs_{timestamp}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            st.session_state["audit_log"].append({
                "user": username,
                "action": "download",
                "details": f"audit_logs_{timestamp}.xlsx",
                "timestamp": datetime.now()
            })

if "audit_log" not in st.session_state:
    st.session_state["audit_log"] = []
    
    
